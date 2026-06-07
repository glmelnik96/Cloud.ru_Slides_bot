"""Python-script nodes — wrap vendored skill scripts (no LLM calls).

Each node:
- Reads the upstream artefact it needs (raises if missing).
- Emits a progress event.
- Writes its output under a stable artefacts[] key.

Output shapes are typed through the corresponding Pydantic model in
``schemas/slides.py`` wherever feasible.

Pending wiring (S3 round-trip for input/output, donor-slot-map loading
inside distribute_node, LibreOffice render, native chart engine) is
tracked by inline FIXME(next-chunk) comments.
"""
from __future__ import annotations

import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

import structlog

from schemas.session import SessionState, Stage
from schemas.slides import (
    BrandReport,
    ParsedDeck,
    Plan,
    PlanSlide,
    VerifierVerdict,
)
from worker import progress, skill_bridge

logger = structlog.get_logger(__name__)


def _artefacts(state: SessionState) -> dict[str, Any]:
    return dict(state.artefacts)


def _emit(state: SessionState, stage: Stage, pct: int, detail: str) -> None:
    progress.stage(state.session_id, stage, pct=pct, detail=detail)


# Anything that isn't a word char (Unicode-aware, so Cyrillic survives), dot,
# or hyphen is collapsed to a single underscore in the output filename.
_SAFE_NAME_RE = re.compile(r"[^\w.-]+", re.UNICODE)


def _output_filename(session_id: str, source_filename: str | None) -> str:
    """Name the built deck after its run ID, suffixed with the sanitised source
    stem when available: ``{session_id}_{source}.pptx`` (falls back to
    ``{session_id}.pptx``). Keeps a deck traceable to its run in logs/Telegram.
    """
    if source_filename:
        stem = _SAFE_NAME_RE.sub("_", Path(source_filename).stem).strip("_")
        if stem:
            return f"{session_id}_{stem}.pptx"
    return f"{session_id}.pptx"


def _resolve_input_path(state: SessionState) -> Path | None:
    """Locate the input draft on disk.

    Resolution order (M3 interim — S3 layer lands in M5):
        1. ``state.artefacts['input_path']`` — explicit override used by
           integration tests and the local-disk worker path.
        2. ``state.input_s3_key`` — treated as a local path while
           ``storage/s3.py`` is not yet implemented. The schema field
           keeps its name so the contract with the bot doesn't churn.
    Returns ``None`` if neither is set or the path does not exist.
    """
    override = state.artefacts.get("input_path") if state.artefacts else None
    candidates = [override, state.input_s3_key]
    for c in candidates:
        if not c:
            continue
        p = Path(str(c))
        if p.is_file():
            return p
    return None


def _render_first_slide_png(pptx_path: Path) -> str | None:
    """Render slide 1 of ``pptx_path`` to PNG via LibreOffice headless.

    Runs ``render_slides.py`` in a subprocess so its ``sys.exit(1)`` on
    failure (used by the vendored script) doesn't tear down the worker.
    Returns the absolute PNG path as a string (consumable by
    ``llm.client.VisionImage = str | bytes | Path``) or ``None`` if
    LibreOffice / pdftoppm aren't available — Brief Reader degrades to
    its 1×1 placeholder grounding.
    """
    script = Path(skill_bridge.SKILL_SCRIPTS) / "render_slides.py"
    if not script.is_file():  # pragma: no cover — vendored file guaranteed
        return None
    out_dir = Path(tempfile.mkdtemp(prefix="slidesbot_render_"))
    try:
        result = subprocess.run(
            [sys.executable, str(script), str(pptx_path), str(out_dir)],
            capture_output=True, text=True, timeout=180,
        )
        if result.returncode != 0:
            logger.warning("node.parse.render_failed",
                           stderr=result.stderr[-500:] if result.stderr else "")
            shutil.rmtree(out_dir, ignore_errors=True)
            return None
        pngs = sorted(out_dir.glob("slide*.png"))
        if not pngs:
            shutil.rmtree(out_dir, ignore_errors=True)
            return None
        # Keep only slide 1 — Kimi grounding needs one anchor; per-deck
        # cost balloons fast if we ship every slide upstream.
        first = pngs[0]
        return str(first)
    except (subprocess.TimeoutExpired, FileNotFoundError) as e:
        logger.warning("node.parse.render_unavailable", error=str(e))
        shutil.rmtree(out_dir, ignore_errors=True)
        return None


def _render_all_slides_png(pptx_path: Path) -> dict[int, str]:
    """Render every slide to PNG via render_slides.py. Returns {slide_num: path}.

    Empty dict if soffice/pdftoppm unavailable (caller degrades). Output files
    are named slide-01.png, slide-02.png, … (1-based) by the vendored script.
    Caller owns cleanup of the returned paths' parent temp dir.
    """
    script = Path(skill_bridge.SKILL_SCRIPTS) / "render_slides.py"
    if not script.is_file():  # pragma: no cover — vendored file guaranteed
        return {}
    out_dir = Path(tempfile.mkdtemp(prefix="slidesbot_renderall_"))
    try:
        result = subprocess.run(
            [sys.executable, str(script), str(pptx_path), str(out_dir)],
            capture_output=True, text=True, timeout=300,
        )
        if result.returncode != 0:
            logger.warning("node.parse.render_all_failed",
                           stderr=result.stderr[-500:] if result.stderr else "")
            return {}
        mapping = {}
        for png in sorted(out_dir.glob("slide-*.png")):
            stem = png.stem.split("-")[-1]
            if stem.isdigit():
                mapping[int(stem)] = str(png)
        return mapping
    except (subprocess.TimeoutExpired, FileNotFoundError) as e:
        logger.warning("node.parse.render_all_unavailable", error=str(e))
        return {}


def extract_images_extract(pptx_path, out_dir, manifest=None):
    """Thin wrapper around the vendored extractor (monkeypatchable in tests)."""
    skill_bridge.install()
    import extract_images
    return extract_images.extract(str(pptx_path), str(out_dir), manifest)


def _media_prep_for_slide(pptx_path, slide_num, visual_kind, extract_dir, render_pngs):
    """Return an image_path (str) for a raster/opaque slide, or None.

    raster → largest extracted picture on that slide.
    opaque → pre-rendered full-slide PNG from ``render_pngs`` (slide_num→path).
    Falls through to None when nothing is available (caller logs WARN).
    """
    if visual_kind == "raster":
        try:
            manifest = extract_images_extract(pptx_path, extract_dir)
            imgs = [im for im in manifest.get("images", []) if im.get("slide_num") == slide_num]
        except Exception as e:
            logger.warning("node.parse.media_prep_extract_failed", slide=slide_num, error=str(e))
            imgs = []
        if imgs:
            best = max(imgs, key=lambda im: (im.get("width_px") or 0) * (im.get("height_px") or 0))
            return str(Path(extract_dir) / best["file"])
        # Extraction yielded no raster (commonly a group-nested picture the
        # non-recursive extractor misses) → fall back to the full-slide render
        # (design §5: "extract_images yields no raster → switch to B").
        return render_pngs.get(slide_num)
    if visual_kind == "opaque":
        return render_pngs.get(slide_num)
    return None


# ─── parse_node ──────────────────────────────────────────────────────────────

def parse_node(state: SessionState) -> dict[str, Any]:
    """Parse the uploaded draft into ``ParsedDeck`` (artefacts['parsed_deck']).

    Verstai/audit modes ship a .pptx. Brief mode (.docx) and a future
    markdown path are gated: their classified-slide shape is unrelated
    to ``ParsedDeck`` and lands with M5 (three modes).

    Side effect: renders slide 1 to PNG (subprocess, LibreOffice
    headless) and stores its path under ``artefacts['original_pngs']``
    for Brief Reader vision grounding. Best-effort — skipped silently
    if soffice/pdftoppm aren't on the host.

    FIXME(next-chunk):
        * Pull input bytes from S3 once ``storage/s3.py`` lands; for now
          ``state.input_s3_key`` is treated as a local path.
        * Wire .docx / .md branches when M5 brief mode arrives — they
          should populate ``artefacts['classified_deck']`` instead.
        * Honour the 50 MB Telegram cap once the bot persists files.
    """
    _emit(state, Stage.PARSING, pct=5, detail="разбор файла")
    arts = _artefacts(state)

    path = _resolve_input_path(state)
    if path is None:
        raise RuntimeError(
            "parse_node: no input file — set state.input_s3_key (local path "
            "until S3 lands) or artefacts['input_path']"
        )

    ext = path.suffix.lower()
    if ext != ".pptx":
        raise NotImplementedError(
            f"parse_node: only .pptx supported in M3 (got {ext}); "
            ".docx/.md land with M5 brief mode"
        )

    skill_bridge.install()
    import parse_pptx  # vendored — sys.path mounted by skill_bridge

    raw = parse_pptx.parse(str(path))
    # parse_pptx writes 'slide_size' with width_emu/height_emu keys; ParsedDeck
    # has model_config extra='allow' so the extra keys roundtrip.
    deck = ParsedDeck.model_validate(raw)
    arts["parsed_deck"] = deck.model_dump()

    slides = deck.model_dump().get("slides", [])
    # raster/opaque ground the Brief Reader vision pass. structured slides also
    # get a rendered full-slide PNG stashed as a B-fallback (used by the
    # classifier injector only when native reconstruction yields too few
    # cards) so no visual slide can ever drop — design §5 fallback chain.
    ground_nums = [s["num"] for s in slides
                   if s.get("visual_kind") in ("raster", "opaque")]
    render_nums = ground_nums + [s["num"] for s in slides
                                 if s.get("visual_kind") == "structured"]

    render_pngs = {}
    if render_nums:
        render_pngs = _render_all_slides_png(path)

    # Grounding: slide 1 always; plus every raster/opaque slide we rendered.
    # Reuse slide 1 from the full render when we already paid for it — avoids
    # a second full-deck LibreOffice pass on visual decks.
    grounding = []
    first = render_pngs.get(1) or _render_first_slide_png(path)
    if first is not None:
        grounding.append(first)
    for n in ground_nums:
        p = render_pngs.get(n)
        if p and p not in grounding:
            grounding.append(p)
    if grounding:
        arts["original_pngs"] = grounding
    else:
        arts.pop("original_pngs", None)

    # Resolve image_path per visual slide and stitch it back into parsed_deck.
    # FIXME(media_prep): the raster path re-extracts the whole deck once per
    # raster slide, and extract_dir/render temp trees live in the system temp
    # root (not the session dir) so they aren't swept on session cleanup.
    # Tolerable today (decks carry ~1 raster slide) but worth a single
    # up-front extraction + session-scoped temp dir if raster decks grow.
    extract_dir = Path(tempfile.mkdtemp(prefix="slidesbot_extract_"))
    pd = arts["parsed_deck"]
    for s in pd.get("slides", []):
        vk = s.get("visual_kind")
        if vk == "structured":
            # Stash the rendered full slide as a silent B-fallback; the
            # classifier injector prefers native flow and only uses this when
            # the group nodes can't fill a card grid.
            fb = render_pngs.get(s["num"])
            if fb:
                s["image_path"] = fb
            continue
        if vk not in ("raster", "opaque"):
            continue
        img_path = _media_prep_for_slide(
            pptx_path=path, slide_num=s["num"], visual_kind=vk,
            extract_dir=extract_dir, render_pngs=render_pngs,
        )
        if img_path:
            s["image_path"] = img_path
            # raster that fell back to the rendered PNG is really an image_b.
            rendered = img_path in render_pngs.values()
            route = "image_a" if (vk == "raster" and not rendered) else "image_b"
            logger.info("node.parse.visual_route", slide=s["num"],
                        kind=vk, route=route, image_path=img_path)
        else:
            logger.warning("node.parse.no_image", slide=s["num"], kind=vk)
    arts["parsed_deck"] = pd

    logger.info(
        "node.parse.done",
        session_id=state.session_id,
        path=str(path),
        slide_count=deck.slide_count,
        grounded=bool(grounding),
    )
    return {"artefacts": arts, "stage": Stage.PARSING.value, "progress_pct": 10}


# ─── assemble_plan_node ──────────────────────────────────────────────────────

_NATIVE_BLOCK_KEYS = ("kpi", "chart", "table", "flow", "image")


def _sanitize_native_block(slide_type: str, key: str, block: Any) -> Any:
    """Patch known shape mismatches between classifier output and vendored
    renderers in skill_assets/scripts/.

    These come from prompt-tuning gaps in Agent 02 — the right long-term
    fix is to tighten the classifier system prompt and re-capture cassettes.
    Until then, sanitize at the orchestration boundary so build_v9 has a
    fighting chance:

    * flow_diagram_native with ``grid=false`` AND no explicit x/y/w/h on any
      block: force ``grid=true`` so ``flow_renderer.compose_grid`` derives
      coords from row/col defaults. Blocks then stack but at least render.
    """
    if not isinstance(block, dict):
        return block
    if key == "flow" and slide_type == "flow_diagram_native":
        # Preset archetypes own their layout — flow_renderer dispatches to the
        # preset function and returns before the grid/blocks path. Forcing
        # grid=true here would be a no-op at best, so leave preset blocks alone.
        if (block.get("preset") or "").strip():
            return block
        blocks = block.get("blocks") or []
        if not block.get("grid"):
            has_coords = all(
                isinstance(b, dict) and all(b.get(c) is not None for c in ("x", "y", "w", "h"))
                for b in blocks
            ) if blocks else False
            if not has_coords:
                block = dict(block)
                block["grid"] = True
                # cols default — flow_renderer derives from blocks if absent,
                # but supplying a sane fallback avoids divide-by-zero edges.
                if not block.get("cols"):
                    block["cols"] = max(1, min(len(blocks), 4))
    return block


def _native_block_is_usable(slide_type: str, cls: dict[str, Any]) -> bool:
    """True iff the classifier supplied enough data for build_v9 to render
    the native ``slide_type`` without crashing.

    Build_v9 raises ValueError on empty headers/data/series — that takes down
    the whole pipeline. Detecting it here lets ``assemble_plan_node`` skip
    just the offending slide and keep the rest of the deck.

    Empirical: Agent 02 (Slide Classifier) sometimes picks ``table_native``
    for a slide where the source draft has no tabular data, leaving the
    ``table`` block as ``{"header": "…", "subtitle": "", "style": "…"}``
    with no ``headers`` / ``data``. Same pattern observed for chart_pptx and
    kpi natives.
    """
    if slide_type == "table_native":
        tbl = cls.get("table") or {}
        if not isinstance(tbl, dict):
            return False
        if not (tbl.get("headers") or []):
            return False
        if not (tbl.get("data") or []):
            return False
        return True
    if slide_type == "chart_pptx_native":
        chart = cls.get("chart") or {}
        if not isinstance(chart, dict):
            return False
        if not (chart.get("series") or []):
            return False
        return True
    if slide_type == "kpi_native":
        kpi = cls.get("kpi") or {}
        if not isinstance(kpi, dict):
            return False
        nums = kpi.get("numbers") or []
        if not nums:
            return False
        # Belt-and-braces numeric guard: kpi_renderer.render_kpi() supports only
        # 1-3 numbers and KPI values must be numeric (a bare word like "Прогноз"
        # renders as a giant non-number). classify_node._coerce_overflow_kpis is
        # the primary validation/demotion site; this catches any path that
        # bypasses it (autofix re-injection, future agents). We drop non-numeric
        # values and reject the slide if none remain or >3 survive — assemble
        # then skips it rather than crashing the pipeline.
        #
        # NOTE: the digit check below intentionally mirrors agents._kpi_value_has_digit
        # (same `any(ch.isdigit() …)` rule). It is duplicated rather than imported:
        # pipeline.py is a sibling node module that does not import from agents.py,
        # and there is no shared util home — a one-line check is not worth a new
        # module or a cross-node dependency just to dedupe.
        #
        # NOTE: the `len(valid) > 3` case here deliberately REJECTS (returns
        # False → assemble skips the slide), unlike _coerce_overflow_kpis which
        # REBUILDS >3 into a card_grid to preserve every pair. This is the
        # last-resort guard for paths that bypass the classifier coercion; by
        # the time we reach here the card-grid rebuild window has passed, so the
        # asymmetry is intentional — do NOT "fix" it to rebuild here.
        valid = [n for n in nums
                 if any(ch.isdigit() for ch in str((n or {}).get("value", "")))]
        if not valid or len(valid) > 3:
            return False
        if valid != nums:
            kpi["numbers"] = valid
            cls["kpi"] = kpi
        return True
    if slide_type == "flow_diagram_native":
        flow = cls.get("flow") or {}
        if not isinstance(flow, dict):
            return False
        # Preset archetypes (card_grid / numbered_rows / numbered_columns /
        # hero_statement) carry their content in preset-specific keys, NOT in
        # the generic `blocks` list. Treat the slide as buildable when the
        # preset's required data is present (flow_renderer reads these keys).
        preset = (flow.get("preset") or "").strip()
        if preset:
            preset_data_keys = {
                "card_grid": "cards",
                "numbered_rows": "rows",
                "numbered_columns": "columns",
                "hero_statement": "statement",
            }
            data_key = preset_data_keys.get(preset)
            if data_key is None:
                # Unknown preset → fall through to blocks check below.
                pass
            elif flow.get(data_key):
                return True
            else:
                return False
        if not (flow.get("blocks") or []):
            return False
        return True
    if slide_type == "image_native":
        img = cls.get("image") or {}
        if not isinstance(img, dict):
            return False
        if not (img.get("image_path") or img.get("path")):
            return False
        return True
    # Unknown slide_type → let build_v9 try; if it crashes we'll widen this.
    return True


def _by_num(items: list[dict[str, Any]], key: str = "num") -> dict[int, dict[str, Any]]:
    """Index a list of slide-shaped dicts by their slide number key.

    Different agent outputs use different number keys: Classifier/Layouts
    use ``num``, Distributor/Copy Editor use ``slide_num``. Pass the
    right key per source.
    """
    out: dict[int, dict[str, Any]] = {}
    for item in items or []:
        n = item.get(key) if isinstance(item, dict) else None
        if isinstance(n, int):
            out[n] = item
    return out


# ─── #5 body-line recovery ───────────────────────────────────────────────────
# The Content Distributor (Agent 03) is instructed "слотов < контента →
# объедини … или отбрось наименее важные". When the chosen donor has fewer
# body slots than the brief has content lines, the LLM may silently drop a
# line (live 81673 slide 5: 3 brief lines → 2 distributed → 1 lost). The
# prompt now says "never drop"; this deterministic safeguard is the actual
# guarantee — it re-appends any brief body line the distributed body fails
# to represent, so nothing is lost regardless of LLM compliance.

# Source briefs encode soft line breaks as the OOXML vertical tab (\v /
# U+000B) as well as \n; split on both so we count logical lines correctly.
_BODY_LINE_SPLIT_RE = re.compile(r"[\n\v]+")
# Drop distributor decoration (** key-phrase markup) and punctuation when
# comparing — the distributor reformats copy ("…домен:" → "…домен."), so we
# tokenise on word characters only.
#
# NOTE: we use a LOCAL tokeniser instead of reusing
# ``skill_assets/scripts/text_sanitize.sanitize_text``. That helper lives in
# the vendored skill scripts dir which is only on ``sys.path`` inside the
# worker/skill-bridge subprocess (see build_v9.py:422 — a deferred import),
# not when this host node module imports at top level. Pulling it in here
# would add a fragile cross-package dependency. Our need is narrower anyway —
# strip markdown emphasis + punctuation and tokenise — so a self-contained
# regex is clearer and import-safe. ``\w+`` already discards ``*`` / ``**``
# and punctuation, so no separate markdown pass is required.
_WORD_RE = re.compile(r"\w+", re.UNICODE)

# Tuning constants for the body-recovery matcher (#5). Promoted from magic
# numbers so the false-positive/false-negative trade-off is reviewable.
_MIN_SIG_WORD_LEN = 3          # ignore short connectives (и, в, на, the…)
# A brief line is "covered" by a distributed line when their significant-word
# overlap, measured against the BRIEF line, reaches this ratio. Per-line (not
# pooled) so scattered shared words across other lines can't mask a real drop.
_COVERAGE_THRESHOLD = 0.6
# When re-appending a supposedly-dropped line, suppress it if it is a near
# duplicate / subset of an existing distributed line. Measured as the share of
# the DISTRIBUTED line's words that the brief candidate also contains — i.e.
# "is some kept line essentially a subset of this candidate?" — neutralising
# the rephrase/merge double-content risk from a different angle than coverage.
_DUPLICATE_THRESHOLD = 0.6


# ─── cover title/subtitle swap-guard ─────────────────────────────────────────
# Russian month names (genitive form, as written on slides: "9 Июня 2026").
_DATE_MONTHS_RE = re.compile(
    r"\b(январ|феврал|март|апрел|мая|май|июн|июл|август|сентябр|октябр|"
    r"ноябр|декабр)\w*",
    re.IGNORECASE | re.UNICODE,
)
# A bare 4-digit year in the plausible deck range.
_DATE_YEAR_RE = re.compile(r"\b20(2[0-9]|3[0-5])\b")
# dd.mm.yyyy / dd/mm/yyyy / dd-mm-yyyy numeric date.
_DATE_NUMERIC_RE = re.compile(r"\b\d{1,2}[./-]\d{1,2}[./-]\d{2,4}\b")
# Event markers that, combined with a date signal, mark an event/venue line.
_EVENT_MARKER_RE = re.compile(
    r"(\bдень\b|\bday\b|конференц|форум|митап|саммит|вебинар|"
    r"tech\s*day|go\s*cloud)",
    re.IGNORECASE | re.UNICODE,
)
# Product/title words: their presence vetoes the date/event classification —
# a line naming the product is a title, not a date line, even if it carries a
# trailing year.
_PRODUCT_WORD_RE = re.compile(
    r"(платформ|возможност|решени|сервис|продукт|облак|cloud\.ru|advanced)",
    re.IGNORECASE | re.UNICODE,
)
# Conservative length cap — a real product title can be long; a date/event
# line is short. Beyond this we never classify as date/event.
_DATE_EVENT_MAX_LEN = 60


def _looks_like_date_or_event(text: str) -> bool:
    """Conservative test: does ``text`` look like a date/event/venue line that
    must NOT win the cover slide's title slot?

    Designed to MISS rather than over-fire (a wrong swap is worse than a
    missed one). A line qualifies only when it is short AND carries no obvious
    product/title word AND shows a strong date/event signal:
      • a Russian month name or a plausible 4-digit year (2020–2035), or
      • a numeric dd.mm.yyyy-style date, or
      • an event marker ("День"/"Day"/"конференция"/"Tech Day"/…) — and a
        bare event marker only counts together with a "·" name·date join or a
        date token, so a generic word like "День" alone can't trip it.
    """
    if not text:
        return False
    s = text.strip()
    if not s or len(s) > _DATE_EVENT_MAX_LEN:
        return False
    if _PRODUCT_WORD_RE.search(s):
        return False
    has_month = bool(_DATE_MONTHS_RE.search(s))
    has_year = bool(_DATE_YEAR_RE.search(s))
    has_numeric = bool(_DATE_NUMERIC_RE.search(s))
    has_event = bool(_EVENT_MARKER_RE.search(s))
    has_date = has_month or has_year or has_numeric
    if has_date:
        return True
    # Event marker without a date token only counts when it's a "name · date"
    # style join (the "·"/"|" separator), e.g. "Cloud Tech Day · …".
    if has_event and ("·" in s or "•" in s or "|" in s):
        return True
    return False


def _body_lines(text: str) -> list[str]:
    """Split a body string into non-empty logical lines (\n or \v separated)."""
    if not text:
        return []
    return [ln.strip() for ln in _BODY_LINE_SPLIT_RE.split(text) if ln.strip()]


def _sig_words(text: str) -> set[str]:
    """Lowercased significant-word set of a line.

    Keeps only tokens ≥ ``_MIN_SIG_WORD_LEN`` chars so short connective words
    don't inflate overlap. ``\\w+`` strips markdown (``**``/``*``) and
    punctuation, normalising distributor decoration before comparison.
    """
    return {
        w for w in (m.group(0).lower() for m in _WORD_RE.finditer(text))
        if len(w) >= _MIN_SIG_WORD_LEN
    }


def _max_line_overlap(brief_words: set[str],
                      distributed_line_words: list[set[str]]) -> float:
    """Best per-line overlap ratio of ``brief_words`` against ANY single
    distributed line (NOT a global pooled set). Ratio is over the brief line's
    own significant words, so a brief line counts as represented only when one
    distributed line carries most of its content — scattered shared words
    across different lines can't add up to a false match."""
    if not brief_words:
        return 1.0  # no significant content → treat as represented (noise-safe)
    best = 0.0
    for dwords in distributed_line_words:
        if not dwords:
            continue
        ratio = len(brief_words & dwords) / len(brief_words)
        if ratio > best:
            best = ratio
    return best


def _recover_dropped_body_lines(
    brief_raw_body: list[str],
    slots: dict[str, Any],
    body_slot_names: list[str],
) -> list[str]:
    """Append GENUINELY-dropped brief body lines into the last body slot.
    Mutates ``slots`` in place. Returns the recovered lines (empty in the
    common no-op case).

    Matcher (per-line + distinctive-token, with two anti-duplication gates):
      • Tokenise each distributed line separately; coverage of a brief line is
        its MAX overlap against any SINGLE distributed line, never a pooled
        set. This kills false positives where a dropped line's common words
        happen to be scattered across the kept lines (IMPORTANT 1).
      • A brief line is "covered" — and so NOT re-appended — when that max
        overlap ≥ ``_COVERAGE_THRESHOLD`` (treats reformat/light-edit as kept).
      • WHOLESALE-REPHRASE gate (CRITICAL 2): we only recover when at least one
        OTHER brief line is strongly covered (an anchor proving the distributor
        worked line-by-line — kept some lines and dropped this one). When NO
        brief line is strongly covered yet the line count shrank, the
        distributor MERGED/COMPRESSED everything into a reworded line; its
        surface words barely overlap any original, so re-appending originals
        would duplicate the merged meaning. In that case we trust the
        distributor and recover nothing.
      • DUPLICATE/SUBSET gate: even with an anchor, suppress a candidate drop
        if some distributed line is largely a subset of it (overlap ≥
        ``_DUPLICATE_THRESHOLD``) — it would re-introduce kept content.
      • NON-BODY COVERAGE gate (Task A): suppress a candidate when it is
        BIDIRECTIONALLY contained with any line of any NON-body slot (every slot
        in ``slots`` whose name is NOT in ``body_slot_names`` — the title slot +
        section-header placeholders on multi-section "multicolumn" donors).
        Two directions, each gated at ``_COVERAGE_THRESHOLD``: (a) the candidate
        is a SUBSET of a non-body line (denominator = candidate's words; the
        heading is the longer text), and (b) a non-body line is a SUBSET of the
        candidate (denominator = the non-body line's words; the candidate is a
        SUPERSET that extends the title/heading with extra words). Direction (b)
        catches title-VARIANT superset leaks the candidate-denominator check
        misses — deck3 s12, where the body bullet "ТРЕБОВАНИЯ ДЛЯ УСПЕШНОГО
        ЗАПУСКА ПРОДУКТОВ, ИСПОЛЬЗУЮЩИХ ТЕХНОЛОГИИ ИИ-АГЕНТОВ" (~8 sig words —
        "ИИ" drops below _MIN_SIG_WORD_LEN) is a superset of the 4-word title:
        |title ∩ cand|/|cand| = 0.5 < 0.6 (a) misses it, but
        |cand ∩ title|/|title| = 1.0 ≥ 0.6 (b) suppresses it.
        The slide already renders that content in a non-body placeholder, so
        re-appending it as a trailing body bullet merely overflows the last
        column (deck3 slides 9/12: section headings living in subN slots + a
        title-variant were dumped into body[-1]).
        This SUBSUMES the old exact-equality title gate (the title is one
        non-body slot) while also catching title VARIANTS and headings.
        Accepted cost of the broadening from equality→overlap: a genuinely
        distinct body line that happens to share ≥ ``_COVERAGE_THRESHOLD`` of
        its significant words with a heading/title will be dropped (deemed
        worth it — overflow off-slide is the worse failure). The BODY-only
        fast-gate and anchor remain BODY-only so genuine dropped-body-line
        recovery is unaffected.

    Guard rails: needs ≥1 body slot; fast-gates out when the distributed body
    already has ≥ as many non-empty lines as the brief (distributor split/kept
    the budget → trust it).
    """
    if not body_slot_names:
        return []
    brief_lines: list[str] = []
    for chunk in brief_raw_body or []:
        if isinstance(chunk, str):
            brief_lines.extend(_body_lines(chunk))
    if not brief_lines:
        return []

    distributed_text = "\n".join(
        str(slots.get(name) or "") for name in body_slot_names
    )
    distributed_lines = _body_lines(distributed_text)
    # Fast gate: distributor kept (or expanded) the line budget → trust it.
    if len(distributed_lines) >= len(brief_lines):
        return []

    distributed_line_words = [_sig_words(dl) for dl in distributed_lines]
    # Non-body coverage gate (Task A): tokenise every line of every NON-body
    # slot (slots whose name is not a body slot — title + section-header
    # placeholders). A candidate covered by one of these already renders in a
    # non-body placeholder, so re-appending it as a body bullet just overflows
    # the last column. This subsumes the old exact-equality title gate.
    non_body_line_words: list[set[str]] = []
    for name, value in slots.items():
        if name in body_slot_names:
            continue
        for line in _body_lines(str(value or "")):
            words = _sig_words(line)
            if words:
                non_body_line_words.append(words)

    # Pre-compute per-brief-line coverage so we can apply the wholesale-rephrase
    # anchor gate before deciding what (if anything) to recover.
    coverage = [
        _max_line_overlap(_sig_words(bl), distributed_line_words)
        for bl in brief_lines
    ]
    has_anchor = any(ov >= _COVERAGE_THRESHOLD for ov in coverage)
    if not has_anchor:
        # No 1:1 anchor → distributor compressed/rephrased wholesale; recovering
        # originals would duplicate the merged meaning (CRITICAL 2). Trust it.
        return []

    missing: list[str] = []
    for bl, overlap in zip(brief_lines, coverage):
        if overlap >= _COVERAGE_THRESHOLD:
            continue  # already represented (kept / lightly reformatted)
        bwords = _sig_words(bl)
        # Non-body coverage gate (Task A): drop a candidate already represented
        # in a non-body slot (title / section-header). Subsumes the old
        # exact-equality title gate and also catches title variants + the
        # slide's own section headings that live in column-header placeholders.
        # Bidirectional containment: (a) candidate-denominator — candidate is a
        # SUBSET of a non-body line (heading is longer); (b) non-body-denominator
        # — a non-body line is a SUBSET of the candidate (candidate is a SUPERSET
        # that extends the title/heading with extra words). Direction (b) catches
        # title-VARIANT superset leaks (deck3 s12: body bullet "<title> ПРОДУКТОВ,
        # ИСПОЛЬЗУЮЩИХ ТЕХНОЛОГИИ ИИ-АГЕНТОВ" — |title ∩ cand|/|cand| = 0.5 misses
        # it, but |cand ∩ title|/|title| = 1.0 suppresses it).
        if _max_line_overlap(bwords, non_body_line_words) >= _COVERAGE_THRESHOLD:
            continue
        if any(
            nb_words and len(bwords & nb_words) / len(nb_words) >= _COVERAGE_THRESHOLD
            for nb_words in non_body_line_words
        ):
            continue
        # Duplicate/subset gate: skip if any distributed line is largely a
        # subset of this candidate — it would re-introduce kept content.
        if any(
            dwords and len(bwords & dwords) / len(dwords) >= _DUPLICATE_THRESHOLD
            for dwords in distributed_line_words
        ):
            continue
        missing.append(bl)

    if not missing:
        return []

    target = body_slot_names[-1]
    existing = str(slots.get(target) or "").strip()
    addition = "\n".join(missing)
    slots[target] = f"{existing}\n{addition}" if existing else addition
    return missing


def assemble_plan_node(state: SessionState) -> dict[str, Any]:
    """Fold classification + layouts + copy-edited content + icons +
    infographics into a single ``Plan`` consumable by ``build_v9.py``.

    For each slide:
        * Classification.slide_type set → native PlanSlide carrying the
          matching kpi/chart/table/flow/image block + ``dark``.
        * Otherwise → donor PlanSlide(clone_from_slide=layout_idx,
          slots, slot_styles_override). Slot keys are translated from
          ``ph_idx`` to canonical slot names via
          ``donor_map.slot_name_by_ph_idx`` so build_v9 finds them.
        * Infographic shapes (Agent 06) for non-``none`` types are
          attached under the ``infographic`` extra field. ``PlanSlide``
          uses ``extra='allow'``, so build_v9 sees the key verbatim.
        * Icon assignments are attached the same way under ``icons``.

    Slides Classifier rejected as ``slide_type=None`` AND with no donor
    (layout_idx=0 + no native block) are skipped with a warning — they
    are unbuildable.
    """
    from graph import donor_map  # noqa: WPS433 — keep cycle local

    _emit(state, Stage.DESIGNING, pct=78, detail="сборка плана")
    arts = _artefacts(state)

    classification_slides = (arts.get("classification") or {}).get("slides", [])
    layouts_slides = (arts.get("layouts") or {}).get("slides", [])
    # copy_edited is the cleaned form of `content` — prefer it; fall back.
    content_src = arts.get("copy_edited") or arts.get("content") or {}
    content_slides = content_src.get("slides", [])
    infographics_slides = (arts.get("infographics") or {}).get("slides", [])
    icons_slides = (arts.get("icons") or {}).get("slides", [])

    # D3 fix (2026-06-05): brief-derived title fallback for donor slides
    # whose distributor output omits the title placeholder. Live run2.slide1
    # rendered the cover with empty title because Agent 03 emitted no
    # placeholder_assignment for ph_idx=1 of donor 4 — the slot stayed empty
    # and build_v9 cleared the donor's "Заголовок" placeholder text. Source
    # of truth for the topic is brief.topic (deck title) and BriefSlide
    # .raw_title (per-slide title).
    brief_data = arts.get("brief") or {}
    brief_topic = (brief_data.get("topic") or "").strip()
    brief_slides_by_num: dict[int, str] = {}
    brief_body_by_num: dict[int, list[str]] = {}
    for bs in (brief_data.get("slides") or []):
        if isinstance(bs, dict):
            n = bs.get("num")
            rt = (bs.get("raw_title") or "").strip()
            if isinstance(n, int) and rt:
                brief_slides_by_num[n] = rt
            rb = bs.get("raw_body")
            if isinstance(n, int) and isinstance(rb, list):
                brief_body_by_num[n] = rb

    cls_by_num = _by_num(classification_slides, key="num")
    lay_by_num = _by_num(layouts_slides, key="num")
    content_by_num = _by_num(content_slides, key="slide_num")
    info_by_num = _by_num(infographics_slides, key="slide_num")
    icons_by_num = _by_num(icons_slides, key="slide_num")

    plan_slides: list[PlanSlide] = []
    skipped: list[int] = []

    # Drive by classification (canonical slide order). Each classification
    # slide either lands as a donor PlanSlide or a native one.
    for cls in classification_slides:
        if not isinstance(cls, dict):
            continue
        num = cls.get("num")
        if not isinstance(num, int):
            continue

        slide_type = cls.get("slide_type")
        layout = lay_by_num.get(num) or {}
        donor = layout.get("layout_idx") or layout.get("donor") or 0

        try:
            if slide_type:
                # Native route — carry the matching data block straight from
                # classification (Agent 02 produces typed blocks per slide_type).
                # Skip if the classifier under-filled the native block — better
                # to lose one slide than crash the whole deck inside build_v9.
                if not _native_block_is_usable(slide_type, cls):
                    logger.warning("node.assemble.native_block_empty",
                                   session_id=state.session_id,
                                   num=num, slide_type=slide_type)
                    skipped.append(num)
                    continue
                kwargs: dict[str, Any] = {
                    "slide_type": slide_type,
                    "dark": bool(cls.get("dark", False)),
                }
                for k in _NATIVE_BLOCK_KEYS:
                    block = cls.get(k)
                    if block is not None:
                        kwargs[k] = _sanitize_native_block(slide_type, k, block)
                ps = PlanSlide(**kwargs)
            else:
                if not donor:
                    # Classifier left slide_type empty AND Designer routed
                    # to native (donor=0). Nothing to build — skip.
                    skipped.append(num)
                    continue
                # Donor route — translate ph_idx → canonical slot name.
                slot_name_map = donor_map.slot_name_by_ph_idx(int(donor))
                if not slot_name_map:
                    # Donor isn't in donor-slot-map.yaml — build_v9 won't
                    # match any slots and template defaults will leak through.
                    # design_node should have caught this; log loudly so we
                    # can extend the YAML or tighten the designer.
                    logger.warning(
                        "node.assemble.donor_unmapped",
                        session_id=state.session_id,
                        num=num,
                        donor=int(donor),
                    )
                cont = content_by_num.get(num) or {}
                phs = cont.get("placeholder_assignments") or []
                slots: dict[str, Any] = {}
                for pa in phs:
                    if not isinstance(pa, dict):
                        continue
                    ph_idx = pa.get("ph_idx")
                    name = slot_name_map.get(int(ph_idx)) if isinstance(ph_idx, int) else None
                    # If donor schema doesn't know this ph_idx, key by the raw
                    # index — build_v9 will warn but still surface the content.
                    key = name or (f"ph_{ph_idx}" if ph_idx is not None else None)
                    if key is None:
                        continue
                    slots[key] = pa.get("content", "")

                # D3 fix (2026-06-05): if the donor has a "title" slot but
                # the Distributor never assigned it, fall back to per-slide
                # raw_title from brief, then deck topic for slide 1. Without
                # this, build_v9 clears the donor's "Заголовок" placeholder
                # and the slide renders with an empty title bar.
                donor_slot_names = set(slot_name_map.values())
                if "title" in donor_slot_names and not (slots.get("title") or "").strip():
                    fallback = brief_slides_by_num.get(num) or (
                        brief_topic if num == 1 else ""
                    )
                    if fallback:
                        slots["title"] = fallback
                        logger.info(
                            "node.assemble.title_fallback",
                            session_id=state.session_id,
                            num=num, donor=int(donor),
                            source="brief.raw_title" if num in brief_slides_by_num else "brief.topic",
                        )

                # Cover swap-guard (2026-06-07): on the COVER/title slide,
                # parse_pptx assigns the FIRST short text run as the title, so
                # a source cover that lists the event/date line BEFORE the real
                # product title ("Cloud Tech Day · 9 Июня 2026" then "Cloud.ru
                # Advanced: …") demotes the product name to the subtitle. Runs
                # AFTER the title slot is populated (so it sees the Distributor's
                # assignment) and AFTER the empty-title D3 fallback, so the two
                # don't conflict. Conservative: only swaps when the title looks
                # like a date/event line, the subtitle does NOT, and the
                # subtitle is a plausible (longer) title.
                if (
                    num == 1
                    and "title" in donor_slot_names
                    and "subtitle" in donor_slot_names
                ):
                    title_txt = (slots.get("title") or "").strip()
                    sub_txt = (slots.get("subtitle") or "").strip()
                    if (
                        title_txt
                        and sub_txt
                        and _looks_like_date_or_event(title_txt)
                        and not _looks_like_date_or_event(sub_txt)
                        and len(sub_txt) >= len(title_txt)
                    ):
                        slots["title"], slots["subtitle"] = sub_txt, title_txt
                        logger.info(
                            "node.assemble.title_subtitle_swap",
                            session_id=state.session_id,
                            num=num, donor=int(donor),
                            old_title=title_txt, new_title=sub_txt,
                        )

                # #5 fix (2026-06-07): recover brief body lines the Distributor
                # dropped (live 81673 slide 5: 3 brief lines → 2 distributed,
                # 1 lost). Append any unrepresented line into the last body slot
                # so no source content is silently discarded.
                #
                # Brief lookup must mirror the established split-renumber
                # convention (agents.py:206 et al). The classifier SPLITS long
                # brief slides and CONTINUES the numeration, so deck ``num`` no
                # longer aligns with brief ``num`` after any split:
                #   • _split_part set → this slide is a FRAGMENT of a split brief
                #     slide whose body is divided across deck slides; per-part
                #     recovery is unsafe (we can't attribute brief lines to a
                #     fragment) → skip recovery entirely.
                #   • otherwise key off _source_slide (falls back to num) so we
                #     compare against the CORRECT brief slide, never a foreign
                #     one (avoids cross-slide contamination after a split).
                body_slot_names = [
                    name for name in slot_name_map.values()
                    if donor_map._slot_name_to_ooxml(name) == "BODY"
                ]
                if cls.get("_split_part") or donor_map.is_timeline_donor(int(donor)):
                    # Timeline donors (Task A Fix 2): fixed-capacity stepN_body
                    # slots — appending overflow into "the last body slot" is
                    # semantically wrong and overflows the slide (deck1 s7,
                    # donor 60). Skip recovery for them, same as split fragments.
                    brief_body_for_recovery: list[str] = []
                else:
                    src = cls.get("_source_slide") or num
                    brief_body_for_recovery = brief_body_by_num.get(src) or []
                recovered = _recover_dropped_body_lines(
                    brief_body_for_recovery,
                    slots,
                    body_slot_names,
                )
                if recovered:
                    logger.warning(
                        "node.assemble.body_recovered",
                        session_id=state.session_id,
                        num=num, donor=int(donor),
                        recovered_lines=len(recovered),
                    )

                ps = PlanSlide(
                    clone_from_slide=int(donor),
                    slots=slots,
                    slot_styles_override=layout.get("slot_styles_override") or {},
                )
        except Exception as e:  # noqa: BLE001 — Pydantic ValidationError + lookups
            logger.warning("node.assemble.slide_skip",
                           session_id=state.session_id, num=num, error=str(e))
            skipped.append(num)
            continue

        # Donor-table data wiring (2026-06-06): table donors (53/54 —
        # ``fixed_png_content`` carrying ``remove_if_user_provides_table``)
        # expect build_v9 to draw a real table from a ``table_data`` extra.
        # The donor route only carries text slots, so the classifier's
        # ``table`` block was silently dropped: build saw no table_data and
        # kept the PNG-stub placeholder ("Столбец 1…/Строка 1…/+" — live dl1
        # slide 4 "DNS Resolvers"). Convert the classifier TableConfig
        # (headers + data) into the list-of-rows build_v9 wants so the stub is
        # stripped and a branded table renders in its place.
        if ps.clone_from_slide is not None:
            donor_def = donor_map._load().get(int(ps.clone_from_slide)) or {}
            is_table_donor = bool(donor_def.get("remove_if_user_provides_table")) or (
                donor_def.get("donor_type") == "fixed_png_content"
                and str(donor_def.get("category", "")).startswith("table")
            )
            tbl = cls.get("table") if isinstance(cls.get("table"), dict) else None
            if is_table_donor and tbl:
                headers = [str(h) for h in (tbl.get("headers") or [])]
                rows = [
                    [str(c) for c in r]
                    for r in (tbl.get("data") or [])
                    if isinstance(r, list)
                ]
                if headers and rows:
                    ps_dump = ps.model_dump()
                    ps_dump["table_data"] = [headers] + rows
                    ps = PlanSlide.model_validate(ps_dump)
                    logger.info("node.assemble.donor_table_data",
                                session_id=state.session_id, num=num,
                                donor=int(ps.clone_from_slide),
                                rows=len(rows) + 1, cols=len(headers))
            elif is_table_donor:
                logger.warning("node.assemble.donor_table_no_data",
                               session_id=state.session_id, num=num,
                               donor=int(ps.clone_from_slide))

        # Attach Agent 06 infographic shapes for slides where they apply.
        info = info_by_num.get(num) or {}
        info_type = info.get("infographic_type")
        if info_type and info_type != "none":
            shapes = info.get("shapes") or []
            # Task 5 (2026-06-07): cap process/timeline cards to the layout
            # capacity at the FEED point so the renderer never receives more
            # cards than the horizontal step row fits (~8). Overflow card
            # text is merged into the last shown card — no source word is
            # clipped. Lazy import: infographic_renderer pulls python-pptx,
            # which is only mounted on the skill-scripts path by skill_bridge.
            try:
                skill_bridge.install()
                from infographic_renderer import cap_process_items  # noqa: WPS433
                capped = cap_process_items(info_type, shapes)
                if len(capped) != len(shapes):
                    logger.warning(
                        "node.assemble.infographic_capped",
                        session_id=state.session_id, num=num,
                        type=info_type,
                        before=len(shapes), after=len(capped),
                    )
                shapes = capped
            except Exception as e:  # noqa: BLE001 — never fail assemble on cap
                logger.warning("node.assemble.infographic_cap_failed",
                               session_id=state.session_id, num=num, error=str(e))
            ps_dump = ps.model_dump()
            ps_dump["infographic"] = {
                "type": info_type,
                "shapes": shapes,
            }
            # Re-validate so the extras roundtrip cleanly through Plan.
            ps = PlanSlide.model_validate(ps_dump)

        # Attach icon assignments (Agent 05) under a single 'icons' key.
        # Drop entries whose icon_path is null — the SVG library hasn't been
        # populated yet (only brand_arrow.svg ships with M2), so unresolved
        # picks would land in the plan as "ghost icons" that build_v9 can't
        # render. Visual Verifier then reports them as plan↔PNG mismatches.
        icon_entry = icons_by_num.get(num) or {}
        icon_assigns = [
            a for a in (icon_entry.get("icon_assignments") or [])
            if a.get("icon_path")
        ]
        if icon_assigns:
            ps_dump = ps.model_dump()
            ps_dump["icons"] = icon_assigns
            ps = PlanSlide.model_validate(ps_dump)

        plan_slides.append(ps)

    plan = Plan(slides=plan_slides)
    arts["plan"] = plan.model_dump()
    # Debug artefacts (2026-06-06): persist brief + classification next to
    # plan.json so donor-table / coercion regressions can be diagnosed offline
    # without re-running the LLM pipeline.
    try:
        import json as _dbg_json
        _wd = _session_workdir(state.session_id)
        _wd.mkdir(parents=True, exist_ok=True)
        (_wd / "brief.json").write_text(
            _dbg_json.dumps(arts.get("brief") or {}, ensure_ascii=False),
            encoding="utf-8")
        (_wd / "classification.json").write_text(
            _dbg_json.dumps(arts.get("classification") or {}, ensure_ascii=False),
            encoding="utf-8")
    except Exception as _dbg_e:  # noqa: BLE001
        logger.warning("node.assemble.debug_dump_failed", error=str(_dbg_e))
    logger.info(
        "node.assemble.done",
        session_id=state.session_id,
        slides=len(plan_slides),
        skipped=skipped,
    )
    return {"artefacts": arts, "stage": Stage.DESIGNING.value, "progress_pct": 80}


# ─── build_node — skeleton ───────────────────────────────────────────────────

def _session_workdir(session_id: str) -> Path:
    """Per-session scratch dir.

    Created lazily; not cleaned up between nodes so subsequent nodes
    (build → brand → render) can share artefacts on disk. The worker's
    session-end cleanup hook (M3 close-out) will own teardown.

    Honors ``SLIDESBOT_WORKDIR`` so bot and worker containers can share
    a volume — without this the worker writes ``result.pptx`` to its own
    ``/tmp`` and the bot's send_document on terminal DONE can't see it.
    """
    root = os.environ.get("SLIDESBOT_WORKDIR") or str(Path(tempfile.gettempdir()) / "slidesbot")
    d = Path(root) / session_id
    d.mkdir(parents=True, exist_ok=True)
    return d


def build_node(state: SessionState) -> dict[str, Any]:
    """Call ``build_v9.build()`` with the assembled Plan, write output to
    the per-session workdir, and record its path under
    ``artefacts['built_pptx_path']``.

    S3 upload lands with M5 (storage/s3.py). For M3 the worker keeps the
    file on local disk and downstream nodes (brand_guard, render_png)
    read from the same path. ``state.result_s3_key`` is set to the path
    string as an interim so the bot side can pick it up via the same
    field once S3 is wired.

    FIXME(next-chunk):
        * Replace local-path stash with S3 upload — populate
          ``state.result_s3_key`` with a real key.
    """
    import json as _json  # local — not used elsewhere in this module

    _emit(state, Stage.RENDERING, pct=82, detail="сборка pptx")
    arts = _artefacts(state)
    plan = arts.get("plan")
    if not plan:
        raise RuntimeError("build_node: artefacts['plan'] missing — assemble_plan_node didn't run")

    skill_bridge.install()
    import build_v9  # vendored

    workdir = _session_workdir(state.session_id)
    plan_path = workdir / "plan.json"
    out_path = workdir / _output_filename(state.session_id, state.source_filename)
    # Strip explicit nulls before handing off to vendored build_v9. Several
    # of its renderers do ``cfg.get("key", {})`` which only fires the default
    # when the key is *absent* — an explicit ``null`` (which Pydantic emits
    # for Optional fields with default=None) crashes them. Pydantic-level
    # dump from assemble_plan_node carries those nulls; we drop them here at
    # the boundary so the in-memory state stays canonical.
    #
    # by_alias=True ensures schema fields with Python-keyword aliases
    # (FlowArrow.src→"from", FlowArrow.dst→"to") emit the JSON form that
    # flow_renderer.py reads. Without this, arrows silently disappear.
    plan_for_build = Plan.model_validate(plan).model_dump(
        exclude_none=True, by_alias=True,
    )
    with plan_path.open("w", encoding="utf-8") as f:
        _json.dump(plan_for_build, f, ensure_ascii=False, indent=None)

    build_v9.build(
        str(plan_path),
        str(skill_bridge.TEMPLATE_PATH),
        str(out_path),
        str(skill_bridge.DONOR_SLOT_MAP),
    )
    if not out_path.is_file():
        raise RuntimeError(f"build_node: build_v9 did not produce {out_path}")

    arts["built_pptx_path"] = str(out_path)
    logger.info(
        "node.build.done",
        session_id=state.session_id,
        path=str(out_path),
        size_bytes=out_path.stat().st_size,
    )
    return {
        "artefacts": arts,
        "stage": Stage.RENDERING.value,
        "progress_pct": 84,
        # Interim: surface the local path through the same field S3 will own.
        "result_s3_key": str(out_path),
    }


# ─── brand_guard_node — skeleton ─────────────────────────────────────────────

def _brand_issue_to_dict(item: dict[str, Any], severity: str) -> dict[str, Any]:
    """Translate a brand_guardian violation/warning entry into BrandViolation
    shape (severity / rule / msg / fix). The script tags issues with
    ``type``/``msg``/``fix``/``shape_idx``; we keep the extras (shape_idx)
    through extra='allow'.
    """
    out = dict(item)  # preserve extras (shape_idx, etc.)
    out["severity"] = severity
    out["rule"] = item.get("type", "unknown")
    out["msg"] = item.get("msg", "")
    out["fix"] = item.get("fix", "")
    return out


def brand_guard_node(state: SessionState) -> dict[str, Any]:
    """Run ``brand_guardian.validate_slide`` over each slide of the built
    pptx and aggregate into a ``BrandReport``.

    Per-slide verdict precedence: any violation → FAIL, else any warning
    → WARN, else OK. Deck verdict follows the same precedence over the
    slide list. Score avg is the rounded mean of slide scores.

    The vendored ``brand_guardian.main()`` does CLI I/O and ``sys.exit``;
    we go straight to ``validate_slide`` so we can stay in-process.
    """
    _emit(state, Stage.VALIDATING, pct=85, detail="бренд-проверка")
    arts = _artefacts(state)
    pptx_path = arts.get("built_pptx_path")
    if not pptx_path or not Path(pptx_path).is_file():
        raise RuntimeError(
            f"brand_guard_node: built pptx not found at {pptx_path!r} — "
            "build_node didn't run or produced no output"
        )

    skill_bridge.install()
    import brand_guardian  # vendored
    from pptx import Presentation

    deck = Presentation(pptx_path)
    slide_reports = []
    for i, slide in enumerate(deck.slides, start=1):
        raw = brand_guardian.validate_slide(slide, i)
        violations = [
            _brand_issue_to_dict(v, "FAIL") for v in raw.get("violations", [])
        ]
        warnings = [
            _brand_issue_to_dict(w, "WARN") for w in raw.get("warnings", [])
        ]
        if violations:
            verdict = "FAIL"
        elif warnings:
            verdict = "WARN"
        else:
            verdict = "OK"
        slide_reports.append({
            "slide_num": i,
            "verdict": verdict,
            # BrandViolation list: keep both severities under `violations`.
            # extra='allow' lets us roundtrip warnings under their own key too.
            "violations": violations + warnings,
            "score": int(raw.get("score", 100)),
            "layout_name": raw.get("layout_name", ""),
        })

    # Deck-level rollup
    has_fail = any(s["verdict"] == "FAIL" for s in slide_reports)
    has_warn = any(s["verdict"] == "WARN" for s in slide_reports)
    deck_verdict = "FAIL" if has_fail else ("WARN" if has_warn else "OK")
    score_avg = (
        round(sum(s["score"] for s in slide_reports) / len(slide_reports))
        if slide_reports else 100
    )

    report = BrandReport.model_validate({
        "verdict": deck_verdict,
        "score_avg": score_avg,
        "slides": slide_reports,
    })
    arts["brand_report"] = report.model_dump()
    logger.info(
        "node.brand.done",
        session_id=state.session_id,
        verdict=deck_verdict,
        score_avg=score_avg,
        slides=len(slide_reports),
    )
    return {
        "artefacts": arts,
        "stage": Stage.VALIDATING.value,
        "progress_pct": 87,
        "brand_score": score_avg,
    }


# ─── render_png_node — skeleton ──────────────────────────────────────────────

_MAX_VERIFIER_SLIDES = 20
"""Hard cap on PNGs handed to Visual Verifier. Empirically ~10s per
big-deck Kimi call already; beyond 20 slides verifier accuracy degrades
faster than reasoning budget can compensate (see memory/cloudru_fm_api.md
token budget rule of thumb)."""


def render_png_node(state: SessionState) -> dict[str, Any]:
    """Render the built pptx to per-slide PNGs (LibreOffice headless via
    ``render_slides.py``) and stash the paths under
    ``artefacts['rendered_pngs']`` for the Visual Verifier.

    Capped at ``_MAX_VERIFIER_SLIDES`` to bound Kimi vision cost. Stores
    string paths (``VisionImage = str | bytes | Path`` so the LLM client
    will read them lazily).
    """
    _emit(state, Stage.RENDERING, pct=88, detail="рендер PNG")
    arts = _artefacts(state)
    pptx_path = arts.get("built_pptx_path")
    if not pptx_path or not Path(pptx_path).is_file():
        raise RuntimeError(
            f"render_png_node: built pptx not found at {pptx_path!r}"
        )

    workdir = _session_workdir(state.session_id) / "pngs"
    workdir.mkdir(parents=True, exist_ok=True)

    script = Path(skill_bridge.SKILL_SCRIPTS) / "render_slides.py"
    result = subprocess.run(
        [sys.executable, str(script), str(pptx_path), str(workdir)],
        capture_output=True, text=True, timeout=300,
    )
    if result.returncode != 0:
        # Soft-fail: Visual Verifier will still run on whatever we managed
        # to render (possibly zero — it then falls back per its prompt).
        logger.warning(
            "node.render_png.soffice_failed",
            session_id=state.session_id,
            stderr=result.stderr[-500:] if result.stderr else "",
        )

    pngs = sorted(workdir.glob("slide*.png"))[:_MAX_VERIFIER_SLIDES]
    arts["rendered_pngs"] = [str(p) for p in pngs]
    logger.info(
        "node.render_png.done",
        session_id=state.session_id,
        rendered=len(pngs),
        capped=len(pngs) == _MAX_VERIFIER_SLIDES,
    )
    return {"artefacts": arts, "stage": Stage.RENDERING.value, "progress_pct": 89}


# ─── process_verify_node — skeleton ──────────────────────────────────────────

def process_verify_node(state: SessionState) -> dict[str, Any]:
    """Agent 09 — synthesise validate_plan errors + brand_guardian +
    LLM Visual Verifier into the single READY / NEEDS_REWORK gate.

    Decision rule:
        * validate_plan error → NEEDS_REWORK (plan is broken upstream).
        * brand verdict == FAIL → NEEDS_REWORK.
        * visual llm_verdict == NEEDS_REWORK → NEEDS_REWORK.
        * else → READY (warnings are surfaced but don't block).

    Per-slide ``checklist_results`` carry per-source issue lists for the
    UI summary; deck ``blockers`` aggregate hard failures, ``warnings``
    everything else.

    M4 will hand a NEEDS_REWORK verdict to an autofix branch; for M3 the
    finalize node just surfaces it to the user.
    """
    _emit(state, Stage.VALIDATING, pct=94, detail="свод проверок")
    arts = _artefacts(state)
    plan = arts.get("plan") or {}
    brand = arts.get("brand_report") or {}
    visual = arts.get("visual_verdict") or {}

    # 1. Re-run validate_plan for the assembled plan. We pass the loaded
    # donors dict directly so the script doesn't re-read the YAML.
    skill_bridge.install()
    import validate_plan as vp  # vendored
    from graph import donor_map  # noqa: WPS433
    donors = donor_map._load()  # noqa: SLF001 — internal cache reuse

    blockers: list[str] = []
    warnings: list[str] = []
    checklist: dict[str, Any] = {}

    # validate_plan emits "auto-added canonical color=..." every time a slot
    # lacks an explicit colour and build_v9 backfills it with the brand
    # default — i.e. on every well-formed slide. These messages are
    # informational, not actionable, and inflate the warning count (11/21
    # warnings on the 2026-06-04 live run were of this shape, dragging the
    # verifier score down). Keep them in ``checklist`` for the UI summary
    # but drop them from the deck-level warnings roll-up.
    def _is_noise(msg: str) -> bool:
        return "auto-added canonical" in msg or "применено canonical правило" in msg

    plan_slides = plan.get("slides") or []
    for idx, slide in enumerate(plan_slides, start=1):
        _, errs, warns = vp.validate_slide(idx, slide, donors)
        checklist[str(idx)] = {
            "checks_passed": int(not errs),
            "issues": [*errs, *warns],
        }
        blockers.extend(f"slide {idx}: {e}" for e in errs)
        warnings.extend(
            f"slide {idx}: {w}" for w in warns if not _is_noise(w)
        )

    # 2. Roll in Brand Guardian fails as blockers.
    brand_verdict = brand.get("verdict", "WARN")
    if brand_verdict == "FAIL":
        for sb in brand.get("slides", []) or []:
            if sb.get("verdict") == "FAIL":
                num = sb.get("slide_num")
                for v in sb.get("violations") or []:
                    if v.get("severity") == "FAIL":
                        blockers.append(f"brand slide {num}: {v.get('rule')} — {v.get('msg')}")

    # 3. Visual Verifier rejects.
    visual_verdict = visual.get("llm_verdict", "READY")
    for sv in visual.get("slides", []) or []:
        sv_v = sv.get("slide_verdict")
        if sv_v in ("REJECT", "NEEDS_REWORK"):
            num = sv.get("num")
            for iss in sv.get("issues") or []:
                sev = iss.get("severity")
                tag = f"visual slide {num}: {iss.get('rule')} — {iss.get('msg')}"
                (blockers if sev == "FAIL" else warnings).append(tag)

    deck_verdict = "NEEDS_REWORK" if (
        blockers or brand_verdict == "FAIL" or visual_verdict == "NEEDS_REWORK"
    ) else "READY"

    # Score: prefer visual score (rubric-based), fall back to brand score.
    score_avg = int(visual.get("score_avg") or brand.get("score_avg") or 0)

    verdict = VerifierVerdict(
        verdict=deck_verdict,
        score_avg=score_avg,
        checklist_results=checklist,
        blockers=blockers,
        warnings=warnings,
        next_actions=list(visual.get("next_actions") or []),
    )
    arts["verifier_verdict"] = verdict.model_dump()
    logger.info(
        "node.process_verify.done",
        session_id=state.session_id,
        verdict=deck_verdict,
        blockers=len(blockers),
        warnings=len(warnings),
        score=score_avg,
    )
    return {
        "artefacts": arts,
        "stage": Stage.VALIDATING.value,
        "progress_pct": 96,
    }


# ─── finalize_node — terminal ────────────────────────────────────────────────

def finalize_node(state: SessionState) -> dict[str, Any]:
    """Publish the terminal progress event with the built .pptx path so the
    bot side can send the result to the user.

    M3 interim: ``result_path`` is a local filesystem path (worker and bot
    share the same machine). M5 will swap this for an S3 key without
    changing the field name.
    """
    arts = _artefacts(state)
    verdict = arts.get("verifier_verdict", {}).get("verdict", "NEEDS_REWORK")
    built_path = arts.get("built_pptx_path") or state.result_s3_key
    notes = list(state.notes)

    if verdict == "READY":
        notes.append("Готово")
    elif verdict == "NEEDS_REWORK":
        # Build still produced a .pptx — bot can deliver it with a caveat.
        notes.append("Готов черновик, но верификатор просит доработку.")

    progress.done(
        state.session_id,
        detail="готово" if verdict == "READY" else "draft",
        result_path=built_path,
    )
    logger.info("node.finalize.done", session_id=state.session_id,
                verdict=verdict, has_built=bool(built_path))
    return {
        "stage": Stage.DONE.value,
        "progress_pct": 100,
        "notes": notes,
        "artefacts": arts,
    }

"""Deck E — editorial growth chart + phrase-emphasis body: exercises D1 + D2."""
from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation


def build(out: Path) -> Path:
    prs = Presentation()
    title_l = prs.slide_layouts[0]
    bullet = prs.slide_layouts[1]

    s = prs.slides.add_slide(title_l)
    s.shapes.title.text = "Рост платформы за 4 года"
    s.placeholders[1].text = "Ключевые показатели и динамика"

    # Single-metric growth → editorial bar chart (D1).
    s = prs.slides.add_slide(bullet)
    s.shapes.title.text = "Число клиентов растёт кратно"
    tf = s.placeholders[1].text_frame
    tf.text = "2023: 120 клиентов"
    for line in ["2024: 340 клиентов", "2025: 720 клиентов", "2026: 1280 клиентов"]:
        p = tf.add_paragraph(); p.text = line

    # Two-column body with a key phrase to emphasize (D2).
    s = prs.slides.add_slide(bullet)
    s.shapes.title.text = "Почему выбирают нас"
    tf = s.placeholders[1].text_frame
    tf.text = ("Платформа обеспечивает полную изоляцию данных каждого клиента "
               "и сертифицированную защиту по 152-ФЗ")
    for line in [
        "Команда поддержки на связи круглосуточно без выходных",
        "Миграция с других облаков выполняется бесплатно за счёт провайдера",
    ]:
        p = tf.add_paragraph(); p.text = line

    s = prs.slides.add_slide(bullet)
    s.shapes.title.text = "Платформа в цифрах"
    tf = s.placeholders[1].text_frame
    tf.text = "99.95% — доступность сервисов"
    for line in ["1280 — корпоративных клиентов", "6 — регионов"]:
        p = tf.add_paragraph(); p.text = line

    s = prs.slides.add_slide(bullet)
    s.shapes.title.text = "Начните сегодня"
    s.placeholders[1].text_frame.text = "cloud.ru — грант на тестирование"

    prs.save(str(out))
    return out


if __name__ == "__main__":
    out = Path(sys.argv[1] if len(sys.argv) > 1 else "deck_e.pptx")
    p = build(out)
    print(f"wrote {p} ({len(Presentation(str(p)).slides._sldIdLst)} slides)")

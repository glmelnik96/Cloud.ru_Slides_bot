"""Deck D — screenshot-heavy: exercises D5 browser-chrome frame."""
from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation


def build(out: Path) -> Path:
    prs = Presentation()
    title_l = prs.slide_layouts[0]
    bullet = prs.slide_layouts[1]

    s = prs.slides.add_slide(title_l)
    s.shapes.title.text = "Обзор интерфейса платформы"
    s.placeholders[1].text = "Скриншоты ключевых экранов консоли"

    s = prs.slides.add_slide(bullet)
    s.shapes.title.text = "Главная консоль управления"
    s.placeholders[1].text_frame.text = (
        "Скриншот: единая панель мониторинга ресурсов, виджеты нагрузки, "
        "статусы сервисов и быстрые действия в одном окне")

    s = prs.slides.add_slide(bullet)
    s.shapes.title.text = "Экран биллинга"
    s.placeholders[1].text_frame.text = (
        "Скриншот интерфейса биллинга: детализация по сервисам, прогноз "
        "расходов и история платежей")

    s = prs.slides.add_slide(bullet)
    s.shapes.title.text = "Возможности платформы"
    tf = s.placeholders[1].text_frame
    tf.text = "Автомасштабирование. Ресурсы растут под нагрузку автоматически"
    for line in [
        "Мониторинг. Метрики и алерты из коробки",
        "Резервные копии. Снапшоты по расписанию",
        "Доступ. Гранулярные роли и политики IAM",
    ]:
        p = tf.add_paragraph(); p.text = line

    s = prs.slides.add_slide(bullet)
    s.shapes.title.text = "Начните бесплатно"
    s.placeholders[1].text_frame.text = "Регистрация на cloud.ru"

    prs.save(str(out))
    return out


if __name__ == "__main__":
    out = Path(sys.argv[1] if len(sys.argv) > 1 else "deck_d.pptx")
    p = build(out)
    print(f"wrote {p} ({len(Presentation(str(p)).slides._sldIdLst)} slides)")

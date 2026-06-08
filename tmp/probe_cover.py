from pptx import Presentation
from pptx.util import Emu
from renderers.designer import primitives as P
from worker import skill_bridge
skill_bridge.install()
prs = Presentation()
prs.slide_width = Emu(1280 * P.EMU_PER_PX); prs.slide_height = Emu(720 * P.EMU_PER_PX)
s = prs.slides.add_slide(prs.slide_layouts[6])
P.background(s, "green")
P.portal(s, (980, 560, 110), n=4)
P.display_title(s, "Облачная платформа Cloud.ru", (60, 150, 860, 340), color=P.GRAPHITE)
prs.save("tmp/probe_cover.pptx"); print("saved")

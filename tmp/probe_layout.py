import sys
from pptx import Presentation
from pptx.util import Emu
from renderers.designer import primitives as P
from renderers.designer import layouts as L
from worker import skill_bridge
skill_bridge.install()


def render(fn_name, content, dark=False):
    prs = Presentation()
    prs.slide_width = Emu(1280 * P.EMU_PER_PX)
    prs.slide_height = Emu(720 * P.EMU_PER_PX)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    getattr(L, fn_name)(s, content, dark=dark)
    out = f"tmp/probe_{fn_name}.pptx"
    prs.save(out)
    print(out)


pts = [{"head": f"Направление {i}",
        "text": "Краткое описание возможности и её ценности для клиента."}
       for i in range(1, 9)]

if __name__ == "__main__":
    which = sys.argv[1] if len(sys.argv) > 1 else "all"
    if which in ("all", "cover_green"):
        render("cover_green", {"title": "Облачная платформа Cloud.ru",
                               "subtitle": "Инфраструктура нового поколения"})
    if which in ("all", "cover_dark"):
        render("cover_dark", {"title": "Cloud.ru Evolution Stack",
                              "subtitle": "Спикер / дата / подзаголовок"})
    if which in ("all", "section_divider"):
        render("section_divider", {"title": "Архитектура платформы",
                                   "kicker": "Раздел 02"})
    if which in ("all", "points_3"):
        render("points_3", {"title": "Три направления развития", "points": pts[:3]})
    if which in ("all", "points_4"):
        render("points_4", {"title": "Четыре опоры платформы", "points": pts[:4]})
    if which in ("all", "points_6"):
        render("points_6", {"title": "Шесть ключевых возможностей", "points": pts[:6]})
    if which in ("all", "points_8"):
        render("points_8", {"title": "Восемь управляемых сервисов", "points": pts[:8]})
    if which in ("all", "bullet_list"):
        render("bullet_list", {"title": "Преимущества решения",
                               "intro": "Платформа закрывает весь цикл работы с данными.",
                               "bullets": ["Единый контур безопасности",
                                           "Управляемые сервисы данных и ML",
                                           "Гибкое масштабирование под нагрузку",
                                           "Поддержка 24/7 с выделенным менеджером"]})
    if which in ("all", "table_zebra"):
        render("table_zebra", {"title": "Сравнение тарифов", "accent_col": 3,
                               "headers": ["Параметр", "Базовый", "Бизнес", "Enterprise"],
                               "rows": [["vCPU", "8", "32", "128"],
                                        ["RAM, ГБ", "16", "64", "256"],
                                        ["SLA", "99.5%", "99.9%", "99.95%"],
                                        ["Поддержка", "8/5", "24/7", "24/7 + TAM"]]})
    if which in ("all", "chart_columns"):
        render("chart_columns", {"title": "Рост потребления ресурсов",
                                 "categories": ["Q1", "Q2", "Q3", "Q4"],
                                 "series": [{"name": "Compute", "values": [4, 5, 7, 9]},
                                            {"name": "Storage", "values": [2, 3, 4, 6]},
                                            {"name": "Network", "values": [1, 2, 3, 4]}]})
    if which in ("all", "roadmap_timeline"):
        render("roadmap_timeline", {"title": "Дорожная карта 2024",
                                    "milestones": [{"label": "Q1", "text": "Запуск MVP", "accent": True},
                                                   {"label": "Q2", "text": "Бета-доступ"},
                                                   {"label": "Q3", "text": "GA релиз", "accent": True},
                                                   {"label": "Q4", "text": "Экосистема"}]})

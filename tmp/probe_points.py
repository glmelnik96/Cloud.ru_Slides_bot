from pptx import Presentation
from pptx.util import Emu
from renderers.designer import primitives as P
from worker import skill_bridge
skill_bridge.install()
prs = Presentation()
prs.slide_width = Emu(1280 * P.EMU_PER_PX)
prs.slide_height = Emu(720 * P.EMU_PER_PX)
s = prs.slides.add_slide(prs.slide_layouts[6])
P.background(s, "white")
P.title_block(s, "Три направления развития", (40, 40, 1000, 110))
cells = [
    ("Инфраструктура", "Масштабируемые вычисления и хранение в едином контуре."),
    ("Платформа", "Управляемые сервисы данных, ML и контейнеров."),
    ("Экосистема", "Маркетплейс решений и партнёрская сеть."),
]
gap = 30
cw = (1280 - 2 * 40 - gap * 2) / 3
for i, (head, body) in enumerate(cells):
    P.point_item(s, head, body, (40 + i * (cw + gap), 220, cw, 360))
prs.save("tmp/probe_points.pptx")
print("saved")

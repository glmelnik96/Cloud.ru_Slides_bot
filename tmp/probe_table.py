from pptx import Presentation
from pptx.util import Emu
from renderers.designer import primitives as P
from worker import skill_bridge
skill_bridge.install()
prs = Presentation()
prs.slide_width = Emu(1280 * P.EMU_PER_PX); prs.slide_height = Emu(720 * P.EMU_PER_PX)
s = prs.slides.add_slide(prs.slide_layouts[6])
P.background(s, "white")
P.title_block(s, "Сравнение тарифов", (40, 40, 1000, 100))
headers = ["Параметр", "Базовый", "Бизнес", "Enterprise"]
rows = [["vCPU", "8", "32", "128"], ["RAM, ГБ", "16", "64", "256"],
        ["SLA", "99.5%", "99.9%", "99.95%"], ["Поддержка", "8/5", "24/7", "24/7 + TAM"]]
P.table_block(s, headers, rows, (40, 180, 1200, 440), accent_col=3)
prs.save("tmp/probe_table.pptx"); print("saved")

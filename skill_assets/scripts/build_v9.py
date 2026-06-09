#!/usr/bin/env python3
"""
build_v7.py — build_v6 + поддержка ДУБЛЕЙ donor.

Ключевое улучшение vs v6:
- Если donor 13 нужен 3 раза → клонируется 3 раза (не один и не игнорируется)
- Через XML deepcopy slide part-а в presentation
- Сохраняет правильный порядок слайдов

Plan:
{
  "slides": [
    {"clone_from_slide": 13, "slots": {...}},
    {"clone_from_slide": 12, "slots": {...}},
    {"clone_from_slide": 13, "slots": {...}}  ← ДУБЛЬ donor 13!
  ]
}

Usage:
    python3 build_v7.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]
"""
import sys
import json
import os
import re
import copy
import io
import threading
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

# ---------------------------------------------------------------------------
# Process-level template bytes cache (perf B3).
#
# Celery prefork workers are long-lived: the same Presentation template is
# loaded from disk on every deck but the bytes never change inside a worker
# lifetime (Docker image is immutable). Cache raw bytes keyed by
# (resolved_path, mtime_ns, size); invalidate when either changes.
# The Presentation OBJECT is never cached because build() mutates it.
# ---------------------------------------------------------------------------
_TEMPLATE_BYTES_CACHE: dict[tuple, bytes] = {}
_TEMPLATE_BYTES_LOCK = threading.Lock()


def _read_bytes_cached(path: str) -> bytes:
    """Return raw bytes of *path* from a process-level cache.

    Cache key is ``(resolved_path_str, st_mtime_ns, st_size)`` so any
    on-disk change immediately invalidates the entry.  The caller must
    construct ``Presentation(io.BytesIO(bytes_))`` — never reuse the
    Presentation object itself across calls.
    """
    p = Path(path).resolve()
    stat = p.stat()
    key = (str(p), stat.st_mtime_ns, stat.st_size)
    cached = _TEMPLATE_BYTES_CACHE.get(key)
    if cached is not None:
        return cached
    with _TEMPLATE_BYTES_LOCK:
        # Double-checked locking: another thread may have populated it.
        cached = _TEMPLATE_BYTES_CACHE.get(key)
        if cached is not None:
            return cached
        data = p.read_bytes()
        # Evict stale entries for the same path (old mtime/size).
        stale = [k for k in list(_TEMPLATE_BYTES_CACHE) if k[0] == str(p)]
        for k in stale:
            _TEMPLATE_BYTES_CACHE.pop(k, None)
        _TEMPLATE_BYTES_CACHE[key] = data
        return data
from build_v5 import (
    load_donor_map, get_text_frame_by_shape_idx, replace_text_with_style,
    clear_text_frame
)
from kpi_renderer import (
    render_kpi, clean_slide_to_blank,
    BLANK_DONOR_WHITE, BLANK_DONOR_DARK
)
from image_renderer import render_image_native
try:
    from chart_engine import render_chart
    CHART_AVAILABLE = True
except ImportError:
    CHART_AVAILABLE = False

try:
    from chart_native_pptx import render_chart_pptx_slide
    CHART_NATIVE_PPTX_AVAILABLE = True
except ImportError:
    CHART_NATIVE_PPTX_AVAILABLE = False

try:
    from flow_renderer import render_flow_diagram_slide
    FLOW_RENDERER_AVAILABLE = True
except ImportError:
    FLOW_RENDERER_AVAILABLE = False

try:
    from table_renderer import render_table_native
    TABLE_RENDERER_AVAILABLE = True
except ImportError:
    TABLE_RENDERER_AVAILABLE = False

try:
    from infographic_renderer import (
        render_infographic_shapes,
        clear_donor_body_slots,
        clear_donor_non_title_text,
    )
    INFOGRAPHIC_RENDERER_AVAILABLE = True
except ImportError:
    INFOGRAPHIC_RENDERER_AVAILABLE = False

try:
    from bullet_splitter import split_slot_if_body
    BULLET_SPLITTER_AVAILABLE = True
except ImportError:
    BULLET_SPLITTER_AVAILABLE = False
    def split_slot_if_body(_name, text):  # noqa: D401 — no-op fallback
        return text

try:
    import textfit as _textfit
    import font_resolver as _font_resolver
    GEOFIT_AVAILABLE = True
except ImportError:
    GEOFIT_AVAILABLE = False

EMU_PER_PX = 9525

# Geometry-aware subtitle fit (D-fix 2026-06-06): some title donors (e.g.
# donor 5 "title_white_with_3d") carry a tiny subtitle box (≈338×54px) anchored
# near the slide bottom, designed for a short "speaker/caption" string. When the
# real subtitle is a full sentence (≤ donor max_chars but far too long for that
# box) it wraps to 4-5 lines and clips off the bottom edge. char-based autofit
# never fires because len(text) ≤ max_chars. This guard measures the actual box
# geometry and (a) widens a narrow bottom-region box into the clear area and
# (b) shrinks the font until the wrapped text fits the available vertical space.
_SUBTITLE_SAFE_BOTTOM_PX = 690   # leave room for footer copyright strip
_SUBTITLE_SAFE_RIGHT_PX = 1245
_SUBTITLE_MIN_PT = 16
_SUBTITLE_BOTTOM_REGION_PX = 540  # only widen boxes anchored below this (clear)


def _fit_subtitle_box(shape, base_pt, text):
    """Resize/shrink a subtitle shape so multi-line text does not clip.

    Mutates ``shape`` width when it is a narrow box in the bottom region.
    Returns the (possibly shrunk) font size in pt, or ``None`` to leave the
    donor size untouched.
    """
    import math
    try:
        top_px = shape.top / EMU_PER_PX
        left_px = shape.left / EMU_PER_PX
        width_px = shape.width / EMU_PER_PX
    except Exception:
        return None
    base_pt = float(base_pt or 20)
    txt = str(text or "")
    if not txt.strip():
        return None

    # (a) Widen a narrow box that sits in the clear bottom region.
    if width_px < 800 and top_px >= _SUBTITLE_BOTTOM_REGION_PX:
        new_width_px = max(width_px, _SUBTITLE_SAFE_RIGHT_PX - left_px)
        try:
            shape.width = Emu(int(new_width_px * EMU_PER_PX))
            width_px = new_width_px
        except Exception:
            pass
    try:
        shape.text_frame.word_wrap = True
    except Exception:
        pass

    avail_h = (_SUBTITLE_SAFE_BOTTOM_PX - top_px) * 0.90  # 10% breathing room
    if avail_h <= 0:
        return None

    # Longest explicit line drives wrapping; also honour hard newlines.
    segments = [seg for seg in txt.replace("\v", "\n").split("\n")] or [txt]

    def _lines_at(pt):
        font_px = pt * 1.3333
        char_w = font_px * 0.52  # conservative for SB Sans Display
        line_h = font_px * 1.30
        cpl = max(1, int(width_px / char_w))
        total = 0
        for seg in segments:
            total += max(1, math.ceil(len(seg) / cpl))
        return total * line_h

    pt = base_pt
    while pt > _SUBTITLE_MIN_PT and _lines_at(pt) > avail_h:
        pt -= 1
    if pt < base_pt:
        return int(round(pt))
    return None


_TITLE_MIDWORD_MIN_PT = 40


def _fit_title_no_midword_break(shape, base_pt, text):
    """Shrink a title font so its longest word fits the box width on one line.

    Divider donors carry very large title fonts (e.g. donor 13 = 96pt in a
    1024px box). A single long word ("ТЕХНИЧЕСКИЙ") then exceeds the box width
    and LibreOffice breaks it mid-word ("ТЕХНИЧЕСКИ"/"Й"). char-based autofit
    misses this because total length is small. Returns a reduced pt or None.
    """
    try:
        width_px = shape.width / EMU_PER_PX
    except Exception:
        return None
    base_pt = float(base_pt or 60)
    words = [w for w in str(text or "").split() if w]
    if not words:
        return None
    longest = max(words, key=len)
    # All-caps Cyrillic/Latin glyphs are much wider than mixed case.
    cased = [c for c in longest if c.isalpha()]
    is_upper = bool(cased) and longest == longest.upper()
    factor = 0.80 if is_upper else 0.58
    target_w = width_px * 0.95
    pt = base_pt
    while pt > _TITLE_MIDWORD_MIN_PT:
        word_w = len(longest) * pt * 1.3333 * factor
        if word_w <= target_w:
            break
        pt -= 2
    if pt < base_pt:
        return int(round(pt))
    return None


# --- Geometric fitter (textfit + font_resolver) -------------------------------
# Per-slot fit policy. Numbers are single big glyphs (no wrap); title/subtitle
# may wrap to 2-3 lines and centre vertically when short; body stays
# top-anchored multi-line. Floors keep text legible rather than micro-shrinking.
_GEO_MIN_PT = {"title": 28.0, "subtitle": 14.0, "number": 24.0}
_GEO_MIN_PT_BODY = 12.0


def _slot_kind(slot_name):
    if slot_name in ("title", "subtitle", "number"):
        return slot_name
    if slot_name == "body" or (slot_name[:4] == "body" and slot_name[4:].isdigit()):
        return "body"
    return "other"


# Column-body slots (donor 28: col1_body / col2_body). They are kind="other"
# (NOT line-balanced / NOT vertically centred — columns of different length
# would misalign, see the body-anchor comment near the slot loop), so the
# legacy char-count shrink and the body geo-balance never fire for them. A
# pathological column (e.g. the "ДЕЙСТВИЯ В ОФИСЕ" right column) can therefore
# overflow off-slide. _fit_column_body is the per-column safety net.
_COL_BODY_RE = re.compile(r"^col(\d+)?_body$")


def _fit_column_body(shape, base_pt, text, bold=False):
    """Shrink→truncate fit for a column-body slot so its text physically fits
    the column box, MIRRORING flow_renderer._fit_card_body but WITHOUT enabling
    vertical-centering/line-balancing (columns must stay top-anchored or they
    misalign). Returns ``(size_pt, text)`` — the size never grows above
    ``base_pt`` and, as a last resort, the text is truncated on a word boundary
    with an ellipsis so nothing clips off-slide.

    No-op (returns ``(base_pt, text)``) when geofit/fonts are unavailable or the
    box geometry is degenerate — never renders worse than before.
    """
    if not GEOFIT_AVAILABLE or not str(text or "").strip():
        return base_pt, text
    try:
        box_w_px = shape.width / EMU_PER_PX
        box_h_px = shape.height / EMU_PER_PX
    except Exception:
        return base_pt, text
    if box_w_px <= 0 or box_h_px <= 0:
        return base_pt, text
    try:
        from flow_renderer import _fit_card_body
    except Exception:
        return base_pt, text
    try:
        return _fit_card_body(str(text), box_w_px, box_h_px, float(base_pt), bold)
    except Exception:
        return base_pt, text


def _typeface_of(tf, override):
    """Run typeface + bold for a text frame: override wins, else the donor's
    first-run rPr, else (None, False)."""
    family = (override or {}).get("font")
    bold = (override or {}).get("bold")
    if family is None or bold is None:
        rPr = None
        try:
            p = tf._txBody.find(qn("a:p"))
            r = p.find(qn("a:r")) if p is not None else None
            rPr = r.find(qn("a:rPr")) if r is not None else None
        except Exception:
            rPr = None
        if rPr is not None:
            if family is None:
                latin = rPr.find(qn("a:latin"))
                family = latin.get("typeface") if latin is not None else None
            if bold is None:
                bold = rPr.get("b") == "1"
    return family, bool(bold)


def _geo_fit_slot(shape, tf, slot_name, base_pt, text, override):
    """Geometric fit for one slot using real font metrics.

    Side effect: enables word-wrap on non-number slots (donor wrap normaliser —
    e.g. donor 5's title carries wrap="none" which forces a single clipped
    line). Returns (size_pt, anchor_middle) or None to fall back to the legacy
    char-count fitters (Pillow/font unavailable)."""
    if not GEOFIT_AVAILABLE:
        return None
    kind = _slot_kind(slot_name)
    family, bold = _typeface_of(tf, override)
    font_path = _font_resolver.resolve(family, bold)
    if not font_path:
        return None
    try:
        box_w = shape.width
        box_h = shape.height
    except Exception:
        return None
    if kind != "number":
        try:
            tf.word_wrap = True
        except Exception:
            pass
    res = _textfit.fit_text(
        text,
        box_w_emu=box_w,
        box_h_emu=box_h,
        font_path=font_path,
        base_pt=float(base_pt or 20),
        min_pt=_GEO_MIN_PT.get(kind, _GEO_MIN_PT_BODY),
        wrap=(kind != "number"),
        balance=(kind in ("title", "subtitle", "body")),
    )
    if res is None:
        return None
    return res.size_pt, res.anchor_middle


# C (2026-06-07): bullet + spacing for multi-item body lists.
# Donor body placeholders carry NO pPr at all (no bullet, no spcBef), so a
# distributor/​bullet_splitter list of paragraphs stacks with zero gap and reads
# as one solid wall (taxonomy defect C; live: 5d s4, 9b s5). When a body slot
# holds ≥2 non-empty paragraphs we render it as a proper list: a "•" glyph with
# a hanging indent + inter-item spacing (отбивка). Single-paragraph bodies are
# left untouched — a conclusion is a paragraph, not a list.
_BULLET_CHAR = "•"


def _apply_body_bullets(tf, size_pt):
    """Turn a body text-frame's paragraphs into a bulleted list (in place).

    Skips when there are <2 non-empty paragraphs, or when the donor already
    carries bullet formatting. Scales the hanging indent and spacing to the
    effective font size so it stays proportional after a geofit shrink."""
    try:
        txBody = tf._txBody
    except Exception:
        return
    paras = txBody.findall(qn("a:p"))
    filled = [p for p in paras if (p.find(qn("a:r")) is not None
              and "".join(t.text or "" for t in p.iter(qn("a:t"))).strip())]
    if len(filled) < 2:
        return
    sz = float(size_pt or 18)
    mar_l = int(round(sz * 9525))             # ~1em hanging indent (EMU)
    spc_pts = int(round(sz * 100 * 0.35))     # отбивка ≈ 0.35×font (centi-pts)
    for i, p in enumerate(filled):
        pPr = p.find(qn("a:pPr"))
        if pPr is None:
            pPr = etree.Element(qn("a:pPr"))
            p.insert(0, pPr)
        # Donor already bulleted? leave it alone entirely.
        if (pPr.find(qn("a:buChar")) is not None
                or pPr.find(qn("a:buAutoNum")) is not None):
            continue
        pPr.set("marL", str(mar_l))
        pPr.set("indent", str(-mar_l))
        # Inter-item spacing on every item except the first (no gap above head).
        for old in pPr.findall(qn("a:spcBef")):
            pPr.remove(old)
        if i > 0:
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
            spcPts.set("val", str(spc_pts))
        # Bullet glyph. buFont=Arial for a reliable round glyph across renderers;
        # the bullet inherits the run colour (graphite body text) by default.
        for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
            for old in pPr.findall(qn(tag)):
                pPr.remove(old)
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "Arial")
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", _BULLET_CHAR)


def clone_slide(prs, src_slide):
    """Глубоко копирует slide-part и регистрирует его в presentation.
    Возвращает новый Slide (последний в prs.slides)."""
    from pptx.opc.constants import CONTENT_TYPE as CT
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.constants import RELATIONSHIP_TARGET_MODE as RTM
    from pptx.opc.packuri import PackURI
    from pptx.opc.package import _Relationship
    from pptx.parts.slide import SlidePart

    src_part = src_slide.part
    src_xml = src_part.blob

    # Подбираем уникальное имя slideN.xml
    package = prs.part.package
    existing_partnames = {str(p.partname) for p in package.iter_parts()}
    next_idx = 1
    while f"/ppt/slides/slide{next_idx}.xml" in existing_partnames:
        next_idx += 1
    new_partname = PackURI(f"/ppt/slides/slide{next_idx}.xml")

    # Создаём новый part — используем тот же content_type как у src
    new_part = SlidePart.load(
        partname=new_partname,
        content_type=src_part.content_type,
        blob=src_xml,
        package=package,
    )

    # Копируем relationships, СОХРАНЯЯ оригинальные rId (fix 2026-06-02).
    # Почему не relate_to: он ПЕРЕНУМЕРОВЫВАЕТ rId по порядку обхода, а XML слайда
    # (скопирован блобом) ссылается на ИСХОДНЫЕ rId. В итоге blip-картинка с
    # r:embed="rId3" начинала указывать на slideLayout вместо изображения →
    # PowerPoint «не смог прочитать часть содержимого» и удалял картинку.
    # Сохраняя rId, держим ссылки валидными. notesSlide-связь НЕ копируем: иначе
    # донорский notesSlide ссылается назад на ОРИГИНАЛ-слайд → тот остаётся
    # «сиротой» (в пакете, но не в sldIdLst) → PowerPoint repair. Без неё оригинал
    # недостижим и корректно отбрасывается при save. Заметки донора слайду не нужны.
    # (Зависит от внутренностей python-pptx 1.0.2: _rels / _Relationship.)
    dst_rels = new_part.rels
    for rel in src_part.rels.values():
        if rel.reltype == RT.NOTES_SLIDE:
            continue
        if rel.is_external:
            dst_rels._rels[rel.rId] = _Relationship(
                dst_rels._base_uri, rel.rId, rel.reltype,
                target_mode=RTM.EXTERNAL, target=rel.target_ref)
        else:
            dst_rels._rels[rel.rId] = _Relationship(
                dst_rels._base_uri, rel.rId, rel.reltype,
                target_mode=RTM.INTERNAL, target=rel.target_part)

    # Регистрируем slide в presentation через relationship
    rId = prs.part.relate_to(new_part, RT.SLIDE)

    # Добавляем sldId в sldIdLst
    sldIdLst = prs.slides._sldIdLst
    existing_ids = [int(el.attrib["id"]) for el in sldIdLst if "id" in el.attrib]
    next_id = max(existing_ids) + 1 if existing_ids else 256
    new_sldId = etree.SubElement(sldIdLst, qn("p:sldId"))
    new_sldId.set("id", str(next_id))
    new_sldId.set(qn("r:id"), rId)

    return prs.slides[-1]


def strip_residual_markdown(prs) -> int:
    """Final whole-deck markdown / control-char strip — runs AFTER emphasis.

    apply_kpi_emphasis consumes intentional ``**…**`` emphasis, BUT it SKIPS
    title_like shapes (first run >= 28pt) and skip-type slides. Donors 21/22
    carry a 32pt body → title_like → their ``**`` is never stripped and leaks
    as literal asterisks (session 81673 s5: ``**ССM (Cloud Certificate
    Manager)**``). The per-slot chokepoint wrote that text with
    ``strip_markdown=False`` on purpose so emphasis could see the markers; by
    the time this pass runs all intentional emphasis is already applied, so any
    remaining ``**`` / lone-emphasis ``*`` is a leak and blanket-stripping is
    safe.

    Covers every run on every shape (including title_like bodies emphasis
    skipped), table_native cells and the table subtitle / before-after textboxes
    (all written with ``strip_markdown=False``), plus grouped sub-shapes.
    Idempotent: a second strip_markdown pass over already-clean text is a no-op.

    MUST run AFTER apply_kpi_emphasis — never before, or it would eat the
    markers that drive bolding.

    Returns: number of runs whose text was changed.
    """
    from text_sanitize import sanitize_text

    changed = 0

    def _strip_runs_in_text_frame(text_frame):
        nonlocal changed
        for para in text_frame.paragraphs:
            for run in para.runs:
                cleaned = sanitize_text(run.text, strip_markdown=True)
                if cleaned != run.text:
                    run.text = cleaned
                    changed += 1

    def _strip_shape(shape):
        if shape.has_text_frame:
            _strip_runs_in_text_frame(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _strip_runs_in_text_frame(cell.text_frame)
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for sub in shape.shapes:
                _strip_shape(sub)

    for slide in prs.slides:
        for shape in slide.shapes:
            _strip_shape(shape)
    return changed


def build(plan_path, template_path, output_path, donor_map_path):
    plan = json.load(open(plan_path, encoding="utf-8"))
    p = Presentation(io.BytesIO(_read_bytes_cached(template_path)))
    donors = load_donor_map(donor_map_path)

    # === STEP 1: Собираем нужные donor_nums и клонируем все слайды plan-а ===
    # Сохраняем reference на оригинальные donor slides ДО любых модификаций
    original_slides = list(p.slides)
    donor_originals = {}  # {donor_num: original_slide}
    for ps in plan["slides"]:
        n = ps.get("clone_from_slide")
        if n and n not in donor_originals:
            if 1 <= n <= len(original_slides):
                donor_originals[n] = original_slides[n - 1]
            else:
                print(f"WARN: donor {n} вне диапазона (1..{len(original_slides)})", file=sys.stderr)

    # Клонируем КАЖДЫЙ слайд из plan (включая дубли) → новые slides в конце
    # Для slide_type=="kpi_native": используем blank donor (slide 30/22 шаблона)
    cloned_for_plan = []
    for ps in plan["slides"]:
        slide_type = ps.get("slide_type")
        if slide_type in ("kpi_native", "image_native", "chart_native", "chart_pptx_native", "flow_diagram_native", "table_native"):
            dark = ps.get("dark", False)
            blank_idx = (BLANK_DONOR_DARK if dark else BLANK_DONOR_WHITE)
            if 1 <= blank_idx <= len(original_slides):
                new_slide = clone_slide(p, original_slides[blank_idx - 1])
                cloned_for_plan.append(new_slide)
                continue
        n = ps.get("clone_from_slide")
        if not n or n not in donor_originals:
            cloned_for_plan.append(None)
            continue
        new_slide = clone_slide(p, donor_originals[n])
        cloned_for_plan.append(new_slide)

    # === STEP 2: Удаляем все ОРИГИНАЛЬНЫЕ слайды (101+ template слайдов), оставляем только клоны ===
    sldIdLst = p.slides._sldIdLst
    n_originals = len(original_slides)
    # Первые n_originals элементов — это оригиналы. Удаляем их.
    all_sldIds = list(sldIdLst)
    for sldId in all_sldIds[:n_originals]:
        rId = sldId.attrib[qn('r:id')]
        try:
            p.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)

    # === STEP 3: Заполняем text + pictures для каждого clone ===
    pictures_inserted = 0
    for plan_slide, actual in zip(plan["slides"], cloned_for_plan):
        if actual is None:
            continue

        # === NATIVE RENDERS: build shapes from scratch on clean canvas ===
        slide_type = plan_slide.get("slide_type")
        if slide_type == "kpi_native":
            kpi_config = plan_slide.get("kpi", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_kpi(actual, kpi_config, dark=dark)
            continue
        if slide_type == "image_native":
            image_config = plan_slide.get("image", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_image_native(actual, image_config, dark=dark)
            continue
        if slide_type == "chart_pptx_native":
            if not CHART_NATIVE_PPTX_AVAILABLE:
                print("WARN: chart_native_pptx модуль недоступен — chart_pptx_native пропущен",
                      file=sys.stderr)
                continue
            chart_config = plan_slide.get("chart", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_chart_pptx_slide(actual, chart_config, dark=dark)
            continue
        if slide_type == "flow_diagram_native":
            if not FLOW_RENDERER_AVAILABLE:
                print("WARN: flow_renderer модуль недоступен — flow_diagram_native пропущен",
                      file=sys.stderr)
                continue
            flow_config = plan_slide.get("flow", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_flow_diagram_slide(actual, flow_config, dark=dark)
            continue
        if slide_type == "table_native":
            if not TABLE_RENDERER_AVAILABLE:
                print("WARN: table_renderer модуль недоступен — table_native пропущен",
                      file=sys.stderr)
                continue
            table_config = plan_slide.get("table", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_table_native(actual, table_config, dark=dark)
            continue
        if slide_type == "chart_native":
            if not CHART_AVAILABLE:
                print("WARN: matplotlib не установлен — chart_native пропущен", file=sys.stderr)
                continue
            chart_config = plan_slide.get("chart", {})
            dark = plan_slide.get("dark", False)
            # Render chart to PNG
            chart_png = plan_slide.get("chart_output_png",
                                        f"pptx-skill/output/_chart_slide_{id(plan_slide)}.png")
            render_chart(chart_config, chart_png, dpi=150)
            # Pass to image_native renderer (wide_zone for charts)
            clean_slide_to_blank(actual)
            render_image_native(actual, {
                "title": chart_config.get("slide_title", chart_config.get("title", "")),
                "image_path": chart_png,
                "caption": chart_config.get("caption", "")
            }, dark=dark, wide_zone=True)
            continue

        src_num = plan_slide.get("clone_from_slide")
        if src_num is None:
            continue
        donor_def = donors.get(src_num)

        # === STEP 3a: PRE-CLEANUP (PNG-stripping) ===
        # Источники remove_idx:
        #   1. donor_def.remove_before_fill — всегда удалять
        #   2. plan_slide.remove_shapes — ad-hoc per slide
        #   3. donor_def.remove_if_not_used — удалять если slot пустой
        #      (формат: {slot_name: [shape_idx]} в slot.shape_idx_when_unused — упрощённо мапим)
        #   4. donor_def.remove_if_user_provides_table — удалять если plan имеет table_data
        if donor_def is not None:
            # P0-2 (2026-06-05): donor 53 / 54 mark their PNG-stub in
            # ``remove_before_fill`` so a generated table can take its
            # place. BUT when no table_data is supplied (Agent 03 didn't
            # produce one — live run4.slide4 "DNS Resolvers" was empty),
            # stripping the stub yields a blank slide. Keep the stub for
            # fixed_png_content donors when neither table nor infographic
            # is available — the placeholder PNG is preferable to nothing.
            base_remove = list(donor_def.get("remove_before_fill", []))
            dtype_pre = donor_def.get("donor_type")
            has_replacement = bool(
                plan_slide.get("table_data")
                or plan_slide.get("infographic")
            )
            if dtype_pre == "fixed_png_content" and not has_replacement:
                if base_remove:
                    print(
                        f"WARN: donor {src_num} fixed_png_content без "
                        f"table_data/infographic — оставляю PNG-stub "
                        f"(remove_before_fill={base_remove} suppressed)",
                        file=sys.stderr,
                    )
                base_remove = []
            remove_idx_list = base_remove
            remove_idx_list += list(plan_slide.get("remove_shapes", []))

            # remove_if_user_provides_table: например donor 53 имеет PNG-таблицу-заглушку
            if plan_slide.get("table_data"):
                remove_idx_list += list(donor_def.get("remove_if_user_provides_table", []))

            # remove_if_not_used: парсим формат {slot_name: shape_idx_to_strip}
            # Если slot не указан в plan_slide.slots — добавить shape_idx в remove
            remove_when_unused = donor_def.get("remove_if_not_used", {}) or {}
            slots_filled_now = plan_slide.get("slots", {}) or {}
            if isinstance(remove_when_unused, dict):
                for slot_name, idx_to_strip in remove_when_unused.items():
                    if slot_name not in slots_filled_now:
                        if isinstance(idx_to_strip, list):
                            remove_idx_list += idx_to_strip
                        else:
                            remove_idx_list.append(idx_to_strip)

            # remove_if_slot_empty: {slot_name: shape_idx}. Удаляем декоративный
            # shape слота, когда slot ПРИСУТСТВУЕТ в plan, но его значение пустое.
            # remove_if_not_used (выше) срабатывает только когда ключа нет вовсе;
            # Distributor же часто кладёт caption="" — донор 62 тогда оставлял
            # пустую белую карточку (dl2 slide-19). Здесь чистим и этот случай.
            remove_when_empty = donor_def.get("remove_if_slot_empty", {}) or {}
            if isinstance(remove_when_empty, dict):
                for slot_name, idx_to_strip in remove_when_empty.items():
                    val = slots_filled_now.get(slot_name)
                    if not (isinstance(val, str) and val.strip()):
                        if isinstance(idx_to_strip, list):
                            remove_idx_list += idx_to_strip
                        else:
                            remove_idx_list.append(idx_to_strip)

            # WARN если donor_type=fixed_png_content и нет ни remove_before_fill, ни overrides
            dtype = donor_def.get("donor_type")
            if dtype == "fixed_png_content" and not remove_idx_list:
                print(
                    f"WARN: donor {src_num} is 'fixed_png_content' но без remove_before_fill — "
                    f"PNG-заглушка может перекрыть контент",
                    file=sys.stderr,
                )

            if remove_idx_list:
                spTree = actual.shapes._spTree
                shape_elements = list(spTree)
                content_tags = ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp')
                content_shapes = [el for el in shape_elements
                                  if el.tag.split('}')[-1] in content_tags]
                for idx in sorted(set(remove_idx_list), reverse=True):
                    if 0 <= idx < len(content_shapes):
                        spTree.remove(content_shapes[idx])

        # TEXT slots
        if donor_def is not None:
            slot_defs = donor_def.get("slots", {})
            slots_filled = plan_slide.get("slots", {})
            styles_override = plan_slide.get("slot_styles_override", {})

            for slot_name, new_text in slots_filled.items():
                if slot_name not in slot_defs:
                    print(f"WARN: slot '{slot_name}' undefined for donor {src_num}", file=sys.stderr)
                    continue
                slot_cfg = slot_defs[slot_name]
                shape_idx = slot_cfg["shape_idx"]
                tf = get_text_frame_by_shape_idx(actual, shape_idx)
                if tf is None:
                    continue
                # D7 fix (2026-06-05): wall-of-text safety net. If a body
                # slot landed with a single 300+-char paragraph (distributor
                # didn't split), break at sentence boundaries so the donor's
                # bullet styling actually renders it as a list.
                if isinstance(new_text, str):
                    new_text = split_slot_if_body(slot_name, new_text)
                override = styles_override.get(slot_name)
                # D9 fix (2026-06-05): cover title overflow. When the text
                # noticeably exceeds the slot's safe_max_chars, proactively
                # shrink the font size so it fits — donor 4 title (60pt,
                # safe_max_chars=55) overflowed when the brief topic was
                # 70+ chars. We avoid relying on renderer-side shrink-to-fit
                # (normAutofit) because LibreOffice's autofit support is
                # inconsistent across versions used by render_png.
                txt_str = str(new_text or "")
                safe_max = slot_cfg.get("safe_max_chars") or slot_cfg.get("max_chars")
                base_size = (override or {}).get("size_pt") or slot_cfg.get("size_pt")
                # Geometric fitter (primary): measure the rendered text against
                # the box's real geometry, shrink to fit width+height, and
                # request vertical centring when title/subtitle text underfills.
                # Also normalises donor wrap="none" so titles wrap to 2-3 lines
                # as designed (donor 5) instead of clipping on one line. Falls
                # back to the legacy char-count guards below if Pillow/fonts are
                # unavailable, so it can never render worse than before.
                geo_anchor_middle = False
                geo_done = False
                if txt_str.strip():
                    try:
                        slot_shape = list(actual.shapes)[shape_idx]
                    except Exception:
                        slot_shape = None
                    if slot_shape is not None:
                        geo = _geo_fit_slot(
                            slot_shape, tf, slot_name, base_size, txt_str, override
                        )
                        if geo is not None:
                            geo_done = True
                            geo_size, geo_anchor_middle = geo
                            if base_size and geo_size < float(base_size):
                                override = dict(override or {})
                                override["size_pt"] = geo_size
                                print(
                                    f"geofit: slot={slot_name} donor={src_num} "
                                    f"len={len(txt_str)} size_pt {base_size}→{geo_size} "
                                    f"anchor_mid={geo_anchor_middle}",
                                    file=sys.stderr,
                                )

                # Legacy heuristic fitters — fallback only when geofit is off.
                if not geo_done:
                    if (safe_max and base_size and txt_str
                            and len(txt_str) > int(safe_max)):
                        # Linear shrink with 0.70 floor (below that titles become
                        # unreadable; better to let it clip than render at 8pt).
                        scale = max(0.70, float(safe_max) / float(len(txt_str)))
                        shrunk_pt = max(14, int(round(float(base_size) * scale)))
                        if shrunk_pt < int(base_size):
                            override = dict(override or {})
                            override["size_pt"] = shrunk_pt
                            print(
                                f"autofit: slot={slot_name} donor={src_num} "
                                f"len={len(txt_str)} safe_max={safe_max} "
                                f"size_pt {base_size}→{shrunk_pt}",
                                file=sys.stderr,
                            )
                    # Title mid-word-break guard: shrink large divider/title fonts
                    # so the longest word fits the box width on a single line.
                    if slot_name == "title" and txt_str.strip():
                        try:
                            ttl_shape = list(actual.shapes)[shape_idx]
                        except Exception:
                            ttl_shape = None
                        if ttl_shape is not None:
                            ttl_base = (override or {}).get("size_pt") or base_size or 60
                            fitted_t = _fit_title_no_midword_break(
                                ttl_shape, ttl_base, txt_str
                            )
                            if fitted_t is not None and fitted_t < float(ttl_base):
                                override = dict(override or {})
                                override["size_pt"] = fitted_t
                                print(
                                    f"autofit: slot=title donor={src_num} "
                                    f"len={len(txt_str)} midword-fit "
                                    f"size_pt {ttl_base}→{fitted_t}",
                                    file=sys.stderr,
                                )
                    # Subtitle overflow guard: geometry-aware fit for title-slide
                    # subtitle boxes that are too small/low for a full sentence.
                    if slot_name == "subtitle" and txt_str.strip():
                        try:
                            sub_shape = list(actual.shapes)[shape_idx]
                        except Exception:
                            sub_shape = None
                        if sub_shape is not None:
                            sub_base = (override or {}).get("size_pt") or base_size or 20
                            fitted = _fit_subtitle_box(sub_shape, sub_base, txt_str)
                            if fitted is not None and fitted < float(sub_base):
                                override = dict(override or {})
                                override["size_pt"] = fitted
                                print(
                                    f"autofit: slot=subtitle donor={src_num} "
                                    f"len={len(txt_str)} geom-fit "
                                    f"size_pt {sub_base}→{fitted}",
                                    file=sys.stderr,
                                )
                # Column-body overflow net: col1_body/col2_body are kind=
                # "other" → never line-balanced and (donor 28) carry no
                # safe_max_chars, so neither the geo body-balance nor the legacy
                # char-shrink bounds them. A pathological column overflows
                # off-slide (the "ДЕЙСТВИЯ В ОФИСЕ" right column). Mirror the
                # card-body fit: shrink the font toward the min, then truncate on
                # a word boundary with an ellipsis as a last resort. Runs AFTER
                # geofit so it tightens whatever size was already chosen. Does
                # NOT touch vertical-anchoring (columns stay top-anchored).
                if _COL_BODY_RE.match(slot_name) and txt_str.strip():
                    try:
                        col_shape = list(actual.shapes)[shape_idx]
                    except Exception:
                        col_shape = None
                    if col_shape is not None:
                        col_base = (override or {}).get("size_pt") or base_size or 20
                        _, col_bold = _typeface_of(tf, override)
                        fit_pt, fit_txt = _fit_column_body(
                            col_shape, col_base, str(new_text or ""), col_bold
                        )
                        if fit_pt < float(col_base) or fit_txt != str(new_text or ""):
                            if fit_pt < float(col_base):
                                override = dict(override or {})
                                override["size_pt"] = fit_pt
                            new_text = fit_txt
                            print(
                                f"autofit: slot={slot_name} donor={src_num} "
                                f"len={len(txt_str)} column-fit "
                                f"size_pt {col_base}→{fit_pt} "
                                f"truncated={fit_txt != str(txt_str)}",
                                file=sys.stderr,
                            )
                replace_text_with_style(tf, new_text, override)
                # Geo soft-balance: centre short title/subtitle text vertically
                # so it doesn't stick to the top of an oversized donor box.
                if geo_anchor_middle and slot_name in ("title", "subtitle"):
                    try:
                        from pptx.enum.text import MSO_ANCHOR
                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    except Exception:
                        pass
                # Body-text слоты: вертикальный якорь.
                # Донорские body-плейсхолдеры шаблона якорятся ПО НИЗУ
                # (vanchor=BOTTOM, template slide 21/22) → короткий текст тонет к
                # низу бокса (dl1 slide-5/6). Поэтому по умолчанию форсируем TOP.
                # B (2026-06-07): если текст НЕДОзаполняет крупный body-бокс
                # (geo_anchor_middle: высота блока < 55% бокса), центрируем по
                # вертикали — иначе контент липнет к верху и слайд top-heavy
                # (низ 55-60% пустой). Геометрию не двигаем, только якорь.
                # Колоночные body (col1_body…) — kind="other", сюда не попадают:
                # их центрировать нельзя, иначе колонки разной длины разъедутся.
                if slot_name == "body" or (
                    slot_name[:4] == "body" and slot_name[4:].isdigit()
                ):
                    try:
                        from pptx.enum.text import MSO_ANCHOR
                        tf.vertical_anchor = (
                            MSO_ANCHOR.MIDDLE if geo_anchor_middle
                            else MSO_ANCHOR.TOP
                        )
                    except Exception:
                        pass
                    # C: multi-item body → bulleted list with отбивка so the
                    # points don't merge into a wall (donors carry no pPr).
                    eff_size = (override or {}).get("size_pt") or base_size or 18
                    _apply_body_bullets(tf, eff_size)

            # Очистить незаполненные обязательные слоты
            for slot_name, slot_def in slot_defs.items():
                if slot_name in slots_filled:
                    continue
                if slot_def.get("optional"):
                    continue
                shape_idx = slot_def["shape_idx"]
                tf = get_text_frame_by_shape_idx(actual, shape_idx)
                if tf is not None:
                    clear_text_frame(tf)

        # INFOGRAPHIC native_block (Agent 06): инжектим shape-список (rounded_rect+text)
        # с абсолютным позиционированием поверх клона донора. Когда инфографика
        # есть, body-слоты донора чистим — иначе старый шаблонный текст
        # «просвечивает» между блоками сравнения. Title остаётся (он не
        # дублируется в shape-списке инфографикa).
        info_block = plan_slide.get("infographic") or {}
        info_shapes = info_block.get("shapes") or []
        info_type = (info_block.get("type") or "").lower()
        # B5 (2026-06-05): live a337cc86 slide 12 showed donor 34's three
        # native columns AND Agent 06's 9-shape matrix rendered together.
        # If the donor itself IS a structural multicolumn/matrix layout,
        # the donor's slots already cover the same visual job — skip the
        # Agent 06 overlay entirely so we don't double-render.
        donor_cat = (donor_def.get("category") or "").lower() if donor_def else ""
        _STRUCTURAL_DONOR_CATS = (
            "content_2col", "content_3col_subtitle", "content_4subtitles",
            "content_6subtitles", "content_8subtitles", "content_4block",
        )
        _OVERLAY_TYPES = ("comparison", "matrix", "process", "flow", "tree")
        donor_is_structural = donor_cat in _STRUCTURAL_DONOR_CATS
        donor_already_structural = (
            donor_is_structural and info_type in _OVERLAY_TYPES
        )

        # F1+F2 (2026-06-05): post-run7 visual review (eb6c4ceec3024bd9)
        # showed donor mock decoration ("Подзаголовок в две строки 20pt")
        # leaking through whenever the distributor produced a structural
        # donor with only the title slot filled and B5 dropped the overlay.
        # Two situations to disambiguate:
        #   Case A — distributor filled real body slots: drop overlay AND
        #       clear non-slot decoration only (keep the filled slot text).
        #   Case B — distributor filled only title/caption: keep the overlay
        #       (otherwise the slide has no content at all) AND clear ALL
        #       non-title text under the overlay.
        filled_body_slots_count = 0
        filled_slot_shape_indices: set[int] = set()
        if donor_def is not None:
            _slot_defs_local = donor_def.get("slots", {})
            for _slot_name, _slot_val in (plan_slide.get("slots") or {}).items():
                if _slot_name in ("title", "caption"):
                    continue
                if not str(_slot_val or "").strip():
                    continue
                if _slot_name not in _slot_defs_local:
                    continue
                filled_body_slots_count += 1
                _idx = _slot_defs_local[_slot_name].get("shape_idx")
                if isinstance(_idx, int):
                    filled_slot_shape_indices.add(_idx)

        case_a_drop_overlay = bool(
            info_shapes and donor_already_structural
            and filled_body_slots_count >= 2
        )
        if case_a_drop_overlay:
            print(
                f"infographic: SKIP donor={src_num} cat={donor_cat} "
                f"info_type={info_type} filled_body_slots={filled_body_slots_count} "
                "(skip-overlay rule B5/Case-A)", file=sys.stderr,
            )
            info_shapes = []
        elif info_shapes and donor_already_structural:
            # Case B: structural donor but distributor underfilled — keep
            # the overlay so the slide has actual content. Cleanup below
            # will wipe donor mock decoration before the overlay paints.
            print(
                f"infographic: KEEP donor={src_num} cat={donor_cat} "
                f"info_type={info_type} filled_body_slots={filled_body_slots_count} "
                "(Case-B: donor underfilled, overlay carries content)",
                file=sys.stderr,
            )

        # F1b (2026-06-05): live run8c (6c4e33c898824936) slide 7 used donor
        # 33 (content_6subtitles) with only `title` filled and NO Agent 06
        # overlay. Donor 33's slot map declares only `title`; its template
        # has 6 mock sub-headers ("Подзаголовок в две строки 20pt" + dup
        # body) which weren't slot-mapped, so neither the slot-fill loop
        # nor the F1 paths cleared them — 6 identical mock cells leaked.
        # When a structural donor is paired with a sparse fill AND no
        # overlay, wipe non-title decoration so we get a clean title-only
        # frame instead of duplicated mock content.
        sparse_structural_no_overlay = (
            donor_is_structural
            and filled_body_slots_count < 2
            and not info_shapes
        )
        if sparse_structural_no_overlay:
            print(
                f"infographic: WIPE donor={src_num} cat={donor_cat} "
                f"filled_body_slots={filled_body_slots_count} "
                "(F1b: sparse structural, no overlay — clearing mock decoration)",
                file=sys.stderr,
            )

        needs_cleanup = (
            bool(info_shapes) or case_a_drop_overlay or sparse_structural_no_overlay
        )
        if needs_cleanup and INFOGRAPHIC_RENDERER_AVAILABLE:
            try:
                # D1+D8 (2026-06-05): clear ALL non-title donor text before
                # injecting infographic shapes. Donors often have pre-labeled
                # boxes (process steps, comparison cells) whose labels aren't
                # in donor_def.slots — the old slot-only cleanup left them in
                # place, causing visual overlap with Agent 06's new boxes
                # (run1.slide7 verified). clear_donor_body_slots is now a
                # weaker layer behind the full-slide pass; we keep calling it
                # for the count.
                cleared = (
                    clear_donor_body_slots(actual, donor_def)
                    if donor_def and info_shapes else 0
                )
                if case_a_drop_overlay or sparse_structural_no_overlay:
                    # Preserve filled-slot text so we don't wipe what the
                    # distributor put into sub*/body* slots (Case A: ≥2
                    # body slots; F1b sparse: ≤1 body slot still preserved).
                    cleared_all = clear_donor_non_title_text(
                        actual, preserve_shape_idx=filled_slot_shape_indices,
                    )
                else:
                    cleared_all = clear_donor_non_title_text(actual)
                added = (
                    render_infographic_shapes(actual, info_shapes)
                    if info_shapes else 0
                )
                if added or cleared or cleared_all:
                    print(
                        f"infographic: slide donor={src_num} type={info_block.get('type')} "
                        f"shapes_added={added}/{len(info_shapes)} "
                        f"donor_slots_cleared={cleared} non_title_cleared={cleared_all} "
                        f"case_a_drop_overlay={case_a_drop_overlay} "
                        f"sparse_no_overlay={sparse_structural_no_overlay}",
                        file=sys.stderr,
                    )
            except Exception as e:  # noqa: BLE001 — never fail the build
                print(f"WARN: infographic render failed (donor {src_num}): {e}",
                      file=sys.stderr)
        elif info_shapes and not INFOGRAPHIC_RENDERER_AVAILABLE:
            print("WARN: infographic shapes present but infographic_renderer "
                  "module unavailable — skipping", file=sys.stderr)

        # PICTURES (вставляются ПОВЕРХ donor shapes)
        for pic in plan_slide.get("pictures", []):
            file_path = pic.get("file")
            if not file_path or not os.path.exists(file_path):
                print(f"WARN: image not found: {file_path}", file=sys.stderr)
                continue
            try:
                actual.shapes.add_picture(
                    file_path,
                    Emu(pic.get("left_px", 0) * EMU_PER_PX),
                    Emu(pic.get("top_px", 0) * EMU_PER_PX),
                    Emu(pic.get("width_px", 100) * EMU_PER_PX),
                    Emu(pic.get("height_px", 100) * EMU_PER_PX),
                )
                pictures_inserted += 1
            except Exception as e:
                print(f"WARN: insert_picture failed: {e}", file=sys.stderr)

        # TABLES (v8: fill_existing если donor имеет встроенную таблицу с брендовым стилем!)
        table_data = plan_slide.get("table_data")
        # D6 fix (2026-06-05): degenerate "tables" (one column, one row, or
        # missing cell content) render as a thin sliver — visual verifier
        # rejected slide as «table_native but no rows». Validate shape first;
        # when too thin, drop the table_data and let the body slot carry the
        # content as plain bullets.
        if table_data:
            is_degenerate = (
                not isinstance(table_data, list)
                or len(table_data) < 2
                or not any(isinstance(r, list) and len(r) >= 2 for r in table_data)
            )
            if is_degenerate:
                print(
                    f"WARN: table_data degenerate (rows={len(table_data) if isinstance(table_data, list) else 0}); "
                    f"skipping table render — content should already be in body slot",
                    file=sys.stderr,
                )
                table_data = None
        if table_data:
            try:
                # Найти существующую таблицу в donor (если есть)
                existing_table = None
                for sh in actual.shapes:
                    if sh.has_table:
                        existing_table = sh.table
                        break

                rows_needed = len(table_data)
                cols_needed = max(len(r) for r in table_data) if rows_needed else 1

                if existing_table:
                    # Donor уже имеет таблицу с брендовым стилем — заполняем её!
                    table_rows = len(existing_table.rows)
                    table_cols = len(existing_table.columns)
                    for r_idx, row_data in enumerate(table_data):
                        if r_idx >= table_rows:
                            break
                        for c_idx in range(table_cols):
                            cell = existing_table.cell(r_idx, c_idx)
                            if c_idx < len(row_data):
                                cell.text = str(row_data[c_idx])
                            else:
                                # Лишние колонки очищаем
                                cell.text = ""
                    # Очистить лишние строки если наш data короче
                    for r_idx in range(rows_needed, table_rows):
                        for c_idx in range(table_cols):
                            existing_table.cell(r_idx, c_idx).text = ""
                else:
                    # Donor не имеет таблицы — добавляем новую
                    left = Emu(35 * EMU_PER_PX)
                    top = Emu(120 * EMU_PER_PX)
                    width = Emu(1210 * EMU_PER_PX)
                    height = Emu(min(550, rows_needed * 50) * EMU_PER_PX)
                    tbl_shape = actual.shapes.add_table(rows_needed, cols_needed, left, top, width, height)
                    tbl = tbl_shape.table
                    for r_idx, row_data in enumerate(table_data):
                        for c_idx, cell_text in enumerate(row_data):
                            if c_idx >= cols_needed:
                                continue
                            tbl.cell(r_idx, c_idx).text = str(cell_text)
            except Exception as e:
                print(f"WARN: table fill failed: {e}", file=sys.stderr)

    # === FINAL: canonical enforcement над ВСЕМИ слайдами (для clone-based, где
    # native-фиксы не действуют). БЕЗОПАСНОЕ:
    #   - цвет: зелёный/белый текст → #222222 (кроме тёмного фона) [Problem #2]
    #   - вес: bold → SemiBold [Problem #3]
    #   - размер <12 → 12
    #   - заголовок контент-слайда → штатный TITLE-placeholder (35,38)/20pt
    #     SemiBold CAPS, СЕМАНТИЧЕСКИ (Вариант A) — не «угадывая по позиции»;
    #     титульные/divider и вертикальные/огромные заголовки не трогаются.
    # Bump до 16pt НЕ включаем — он даёт overflow на плотных/код-боксах. ===
    try:
        from enforce_canonical import enforce_canonical_slide, slide_is_dark
        enf_total = {}
        for slide in p.slides:
            st = enforce_canonical_slide(
                slide, dark=slide_is_dark(slide),
                min_pt=12, bump_from=None, bump_to=None, normalize_header=True)
            for k, v in st.items():
                enf_total[k] = enf_total.get(k, 0) + v
        if any(enf_total.values()):
            print(f"enforce_canonical: {enf_total}", file=sys.stderr)
    except Exception as e:
        print(f"WARN: enforce_canonical pass skipped: {e}", file=sys.stderr)

    # === FINAL: KPI emphasis (T2.2) ===
    # Detect digit-heavy runs in body text and bold+green them so 12pt
    # numbers buried in body actually catch the eye. Skips kpi_native
    # slides (render_kpi already styled them) and title-like runs.
    try:
        from kpi_emphasis import apply_kpi_emphasis
        emph_stats = apply_kpi_emphasis(p, plan_slides=plan["slides"])
        if emph_stats["total"]:
            print(f"kpi_emphasis: {emph_stats}", file=sys.stderr)
    except Exception as e:
        print(f"WARN: kpi_emphasis pass skipped: {e}", file=sys.stderr)

    # === FINAL: residual markdown / control-char strip over the WHOLE deck ===
    # apply_kpi_emphasis consumes intentional ``**…**`` emphasis, BUT it SKIPS
    # title_like shapes (first run >= 28pt) and skip-type slides. Donors 21/22
    # carry a 32pt body → title_like → their ``**`` is never stripped and leaks
    # as literal asterisks (session 81673 s5: ``**ССM (Cloud Certificate
    # Manager)**``). The per-slot chokepoint wrote that text with
    # strip_markdown=False on purpose (so emphasis could see the markers).
    # By NOW all intentional emphasis is already applied, so any remaining
    # ``**`` / lone-emphasis ``*`` is a leak and blanket-stripping is safe.
    # Runs AFTER emphasis (never before — that would break bolding) and covers
    # every run on every shape, including title_like bodies, table_native cells
    # and the table subtitle / before-after textboxes (all written with
    # strip_markdown=False). Idempotent: a second strip_markdown pass over text
    # whose control chars are already clean is a no-op.
    try:
        strip_residual_markdown(p)
    except Exception as e:
        print(f"WARN: markdown-strip pass skipped: {e}", file=sys.stderr)

    p.save(output_path)
    print(f"Saved {output_path}: {len(p.slides)} slides, {pictures_inserted} pictures inserted",
          file=sys.stderr)

    # Финальная структурная валидация (ловит orphan-слайды / битые blip-картинки /
    # dangling rId — причины «PowerPoint обнаружил проблему с содержимым»).
    try:
        from validate_deck import validate_pptx
        _problems = validate_pptx(output_path)
        if _problems:
            print("⚠️  DECK VALIDATION: %d проблем(ы) — PowerPoint может ругаться:"
                  % len(_problems), file=sys.stderr)
            for _p in _problems[:30]:
                print("     -", _p, file=sys.stderr)
        else:
            print("✅ DECK VALIDATION: структурно чисто", file=sys.stderr)
    except Exception as _e:
        print(f"(validate_deck пропущен: {_e})", file=sys.stderr)


def main():
    if len(sys.argv) < 4:
        print("Usage: build_v7.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]",
              file=sys.stderr)
        sys.exit(1)
    plan_p = sys.argv[1]
    tpl_p = sys.argv[2]
    out_p = sys.argv[3]
    donor_p = sys.argv[4] if len(sys.argv) > 4 else "pptx-skill/brand/donor-slot-map.yaml"
    build(plan_p, tpl_p, out_p, donor_p)


if __name__ == "__main__":
    main()

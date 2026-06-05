#!/usr/bin/env python3
"""
kpi_emphasis.py — финальный детерминистский пас: подсветка ключевых цифр.

Зачем: visual_verifier повторяет (2026-06-05 ран) что цифры в body-тексте
не выделены типографикой/цветом — «Ключевые цифры (275, 568 125 090)»,
«1,2 млн, 14,2 млн, 25 млн», «101 млн/мес, 100 млн разово». Классификатор
не промоутит multi-number бизнес-слайды в kpi_native (там бывает >3 числа
+ контекст), поэтому числа уходят 12pt regular в body — глаз их не ловит.

Этот пас находит число-токены В РАМКАХ runs body-плейсхолдеров и
разбивает run так, чтобы число стало bold + GREEN (#26D07C). Не трогает:
- title-плейсхолдеры (size_pt >= 28)
- табличные ячейки (table_renderer уже стилизует)
- runs, у которых size_pt не известен (могут быть из шаблонной графики)
- слайды с slide_type=kpi_native (там уже стилизовано render_kpi)

Безопасно для всех донорских лейаутов: работает на XML-уровне через lxml,
сохраняя оригинальный rPr через deepcopy.
"""
from __future__ import annotations

import re
import sys
from copy import deepcopy
from typing import Any

from lxml import etree
from pptx.oxml.ns import qn


# Cloud.ru brand green (canonical accent).
_GREEN_HEX = "26D07C"

# D2 fix (2026-06-05): green text on a green-filled box is invisible.
# Live run1.slide8 had `12.18` recolored green inside a green accent box,
# making the version number disappear. When the parent shape fills green
# we either skip the emphasis (preserve original colour) or fall back to
# graphite — both keep the number readable.
_GREEN_FILL_HEXES = {"26D07C", "00D97B", "1AB066", "1ABF6F", "22C993", "2DD27D"}
_GRAPHITE_HEX = "222222"

# Размер шрифта в hundredths of points; >= 2800 = >= 28pt — это заголовки,
# их не трогаем.
_TITLE_FONT_THRESHOLD_HPT = 2800

# Юнит-суффиксы, делающие число «значимым» даже если оно короткое.
_UNIT_PATTERN = (
    r"(?:млн|млрд|тыс|руб|долл|евро|usd|eur|rub|%|‰|\$|€|"
    r"раз[а-я]*|шт[а-я]*|чел[а-я]*|сек[а-я]*|мин[а-я]*|"
    r"мес[а-я]*|год[а-я]*|кв\.?|ед[а-я]*|"
    r"кг|мг|г|км|мм|см|м|тб|гб|мб)"
)

# Число: 1-3 цифры, опц. группы по 3 через пробел/неразр.пробел/запятую,
# опц. десятичная часть (",5" или ".5"), опц. юнит.
# Должно начинаться с границы (не внутри слова) и опц. иметь юнит/символ.
_NUMBER_RE = re.compile(
    r"(?<![\w])"                                   # не внутри слова/числа
    r"(?P<num>\d{1,3}(?:[ \u00a0\u202f]\d{3})*"
    r"(?:[.,]\d+)?"                                # десятичная
    r"|\d+(?:[.,]\d+)?)"                           # либо просто цифры
    r"(?:\s*(?P<unit>" + _UNIT_PATTERN + r"))?"
    r"(?![\w])",
    re.IGNORECASE,
)


def _qualifies(num: str, unit: str | None) -> bool:
    """Token достоин подсветки, если есть юнит ИЛИ >=3 цифр."""
    if unit:
        return True
    digits = sum(1 for c in num if c.isdigit())
    return digits >= 3


def _set_run_emphasis(rPr: etree._Element, *, color_hex: str = _GREEN_HEX) -> None:
    """Мутирует <a:rPr>: b='1', color=<color_hex> (default green).

    color_hex override lets the caller pick graphite when the parent shape
    fills green — green-on-green would be invisible (D2 fix).
    """
    rPr.set("b", "1")
    # Удаляем существующий solidFill (если был), вставляем новый.
    for sf in rPr.findall(qn("a:solidFill")):
        rPr.remove(sf)
    solid = etree.SubElement(rPr, qn("a:solidFill"))
    srgb = etree.SubElement(solid, qn("a:srgbClr"))
    srgb.set("val", color_hex)


def _shape_fill_hex(shape) -> str | None:
    """Return the shape's solid fill colour as an upper-case hex string,
    or ``None`` if the shape has no solid fill (background/pattern/scheme).

    Used to decide whether green KPI emphasis would be invisible against
    a green-tinted box (D2 fix).
    """
    try:
        fill = shape.fill
        if fill is None or fill.type is None:
            return None
        rgb = fill.fore_color.rgb
        if rgb is None:
            return None
        return str(rgb).upper()
    except Exception:  # noqa: BLE001
        return None


def _shape_bbox(shape) -> tuple[int, int, int, int] | None:
    """Return (left, top, right, bottom) in EMU for a shape, or ``None`` if
    coordinates are unavailable. Used by overlap detection."""
    try:
        left = int(shape.left or 0)
        top = int(shape.top or 0)
        w = int(shape.width or 0)
        h = int(shape.height or 0)
    except Exception:  # noqa: BLE001
        return None
    if w <= 0 or h <= 0:
        return None
    return (left, top, left + w, top + h)


def _collect_green_bboxes(slide) -> list[tuple[int, int, int, int]]:
    """Collect bounding boxes of every shape on the slide whose solid fill
    matches the brand-green palette. Used by P0-1 (2026-06-05): native
    Agent-06 infographics layer a text shape with ``fill=none`` *on top
    of* a green-filled ``rounded_rect``. The text shape itself reports
    no green fill, so the legacy D2 check (parent fill only) misses it
    and the KPI pass paints digits green-on-green → invisible.

    Returns a list of (left, top, right, bottom) EMU tuples.
    """
    boxes: list[tuple[int, int, int, int]] = []
    for sh in slide.shapes:
        fill_hex = _shape_fill_hex(sh)
        if not fill_hex or fill_hex not in _GREEN_FILL_HEXES:
            continue
        bbox = _shape_bbox(sh)
        if bbox is not None:
            boxes.append(bbox)
    return boxes


def _bbox_center(bbox: tuple[int, int, int, int]) -> tuple[float, float]:
    return ((bbox[0] + bbox[2]) / 2.0, (bbox[1] + bbox[3]) / 2.0)


def _point_in_bbox(point: tuple[float, float],
                   bbox: tuple[int, int, int, int]) -> bool:
    x, y = point
    return bbox[0] <= x <= bbox[2] and bbox[1] <= y <= bbox[3]


def _sits_on_green(shape, green_boxes: list[tuple[int, int, int, int]]) -> bool:
    """True if shape's geometric centre lies inside any green-filled
    bounding box on the slide. Cheap proxy for "the shape visually sits
    inside a green block"."""
    if not green_boxes:
        return False
    bbox = _shape_bbox(shape)
    if bbox is None:
        return False
    center = _bbox_center(bbox)
    return any(_point_in_bbox(center, gb) for gb in green_boxes)


def _emphasize_paragraph(p_el: etree._Element, *,
                         emphasis_hex: str = _GREEN_HEX) -> int:
    """Разбивает runs параграфа по KPI-токенам, делает их bold+<emphasis_hex>.

    ``emphasis_hex`` lets the caller override the colour when the parent
    shape's fill would make green invisible (D2 fix).

    Возвращает количество подсвеченных токенов.
    """
    runs = p_el.findall(qn("a:r"))
    if not runs:
        return 0
    emphasized = 0
    for r in list(runs):
        t_el = r.find(qn("a:t"))
        if t_el is None or not t_el.text:
            continue
        text = t_el.text

        # Skip headings — большой шрифт исключаем.
        rPr = r.find(qn("a:rPr"))
        if rPr is not None:
            sz = rPr.get("sz")
            if sz and sz.isdigit() and int(sz) >= _TITLE_FONT_THRESHOLD_HPT:
                continue

        matches = [m for m in _NUMBER_RE.finditer(text)
                   if _qualifies(m.group("num"), m.group("unit"))]
        if not matches:
            continue

        # Build replacement run sequence: [pre, kpi, mid, kpi, post, ...]
        parent = r.getparent()
        insert_idx = list(parent).index(r)

        cursor = 0
        new_runs: list[etree._Element] = []
        for m in matches:
            start = m.start()
            end = m.end()
            # Тeкст до — копия r с обычным rPr.
            if start > cursor:
                pre = deepcopy(r)
                pre_t = pre.find(qn("a:t"))
                pre_t.text = text[cursor:start]
                new_runs.append(pre)
            # Сам KPI токен — копия с emphasis.
            kpi = deepcopy(r)
            kpi_t = kpi.find(qn("a:t"))
            kpi_t.text = text[start:end]
            kpi_rPr = kpi.find(qn("a:rPr"))
            if kpi_rPr is None:
                kpi_rPr = etree.SubElement(kpi, qn("a:rPr"))
                # Move rPr to be first child as required by OOXML schema.
                kpi.insert(0, kpi_rPr)
            _set_run_emphasis(kpi_rPr, color_hex=emphasis_hex)
            new_runs.append(kpi)
            emphasized += 1
            cursor = end

        # Хвост после последнего match.
        if cursor < len(text):
            tail = deepcopy(r)
            tail_t = tail.find(qn("a:t"))
            tail_t.text = text[cursor:]
            new_runs.append(tail)

        # Замена исходного run-а на новую последовательность.
        parent.remove(r)
        for i, nr in enumerate(new_runs):
            parent.insert(insert_idx + i, nr)

    return emphasized


# D2 (2026-06-06): inline key-phrase emphasis. Agent 03 marks ONE key phrase
# per body paragraph with **…**. We strip the markup and (on light slides)
# bold+green the span — matching template slide 28. Guardrails: ≤1 phrase per
# paragraph, ≤6 words; on dark slides emphasis_hex=None → strip-only.
_PHRASE_RE = re.compile(r"\*\*(?P<phrase>[^*]+?)\*\*")
_PHRASE_MAX_WORDS = 6


def _strip_markers_in_run(r: etree._Element) -> None:
    """Remove any literal ** markers from a run's text in place."""
    t_el = r.find(qn("a:t"))
    if t_el is not None and t_el.text and "**" in t_el.text:
        t_el.text = t_el.text.replace("**", "")


def _emphasize_phrase_paragraph(p_el: etree._Element, *,
                                emphasis_hex: str | None) -> int:
    """Strip the first **…** span in the paragraph; if ``emphasis_hex`` is set
    and the phrase is ≤_PHRASE_MAX_WORDS words, bold+colour it. All remaining
    **…** markers in the paragraph are stripped to plain text. Returns 1 if a
    phrase was emphasized, else 0."""
    runs = p_el.findall(qn("a:r"))
    if not runs:
        return 0
    emphasized = 0
    for r in list(runs):
        if emphasized:  # already used our one-per-paragraph budget
            _strip_markers_in_run(r)
            continue
        t_el = r.find(qn("a:t"))
        if t_el is None or not t_el.text or "**" not in t_el.text:
            continue
        text = t_el.text
        m = _PHRASE_RE.search(text)
        if m is None:
            _strip_markers_in_run(r)
            continue
        phrase = m.group("phrase")
        n_words = len([w for w in phrase.split() if w])
        do_color = emphasis_hex is not None and 1 <= n_words <= _PHRASE_MAX_WORDS
        parent = r.getparent()
        insert_idx = list(parent).index(r)
        new_runs: list[etree._Element] = []
        # pre
        if m.start() > 0:
            pre = deepcopy(r)
            pre.find(qn("a:t")).text = text[:m.start()]
            new_runs.append(pre)
        # phrase
        ph = deepcopy(r)
        ph.find(qn("a:t")).text = phrase
        if do_color:
            ph_rPr = ph.find(qn("a:rPr"))
            if ph_rPr is None:
                ph_rPr = etree.SubElement(ph, qn("a:rPr"))
                ph.insert(0, ph_rPr)
            _set_run_emphasis(ph_rPr, color_hex=emphasis_hex)
            emphasized = 1
        new_runs.append(ph)
        # tail — strip any further ** markers here (one-phrase budget per
        # paragraph; the tail run isn't revisited by the outer loop).
        if m.end() < len(text):
            tail = deepcopy(r)
            tail.find(qn("a:t")).text = text[m.end():].replace("**", "")
            new_runs.append(tail)
        parent.remove(r)
        for i, nr in enumerate(new_runs):
            parent.insert(insert_idx + i, nr)
    return emphasized


def emphasize_phrases_in_slide(slide, *, emphasis_hex: str | None) -> int:
    """Apply D2 phrase emphasis to every non-title body shape. Returns the
    number of phrases emphasized on the slide."""
    total = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if _shape_is_title_like(shape):
            continue
        txBody = shape.text_frame._txBody
        for p_el in txBody.findall(qn("a:p")):
            try:
                total += _emphasize_phrase_paragraph(p_el, emphasis_hex=emphasis_hex)
            except Exception as e:  # noqa: BLE001 — never fail the build
                print(f"WARN: phrase emphasis paragraph failed: {e}",
                      file=sys.stderr)
    return total


def _shape_is_title_like(shape) -> bool:
    """Эвристика: shape — это title, если есть placeholder type=TITLE
    либо если первый run >= 28pt."""
    try:
        ph = shape.placeholder_format
        if ph is not None:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if ph.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                           PP_PLACEHOLDER.VERTICAL_TITLE):
                return True
    except (ValueError, AttributeError):
        pass
    # Heuristic by font size of the first run.
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    for p in tf.paragraphs[:1]:
        for run in p.runs[:1]:
            try:
                if run.font.size is not None and run.font.size.pt >= 28:
                    return True
            except Exception:  # noqa: BLE001
                pass
    return False


def emphasize_kpi_in_slide(slide) -> int:
    """Применяет KPI-emphasis ко всем body-shape-ам слайда.

    Возвращает количество подсвеченных токенов на слайде.
    """
    total = 0
    # P0-1 (2026-06-05): pre-compute green bboxes once per slide so we
    # can detect text shapes that visually sit ON TOP of a green-filled
    # rect (Agent 06 native infographic pattern — see live run4.slide8
    # "v1.12.17" lost its digits to green-on-green emphasis).
    green_boxes = _collect_green_bboxes(slide)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if _shape_is_title_like(shape):
            continue
        # D2 fix: if the shape fills green, emphasising in green produces
        # invisible text — fall back to graphite for readable contrast.
        # P0-1 extension: also check overlap with separately-positioned
        # green shapes (Agent 06 layers text shapes over filled rects).
        fill_hex = _shape_fill_hex(shape)
        if fill_hex and fill_hex in _GREEN_FILL_HEXES:
            emphasis_hex = _GRAPHITE_HEX
        elif _sits_on_green(shape, green_boxes):
            emphasis_hex = _GRAPHITE_HEX
        else:
            emphasis_hex = _GREEN_HEX
        # Walk paragraphs at XML level (python-pptx API возвращает Paragraph,
        # но нам нужен прямой доступ к <a:p>).
        txBody = shape.text_frame._txBody
        for p_el in txBody.findall(qn("a:p")):
            try:
                total += _emphasize_paragraph(p_el, emphasis_hex=emphasis_hex)
            except Exception as e:  # noqa: BLE001 — never fail the build
                print(f"WARN: kpi emphasis paragraph failed: {e}", file=sys.stderr)
    return total


def apply_kpi_emphasis(prs, *, skip_slide_types: set[str] | None = None,
                       plan_slides: list[dict[str, Any]] | None = None) -> dict[str, int]:
    """Финальный пас по презентации.

    Args:
        prs: pptx.Presentation.
        skip_slide_types: набор slide_type, для которых пропускаем emphasize
            (например 'kpi_native' — render_kpi уже всё сделал).
        plan_slides: parallel list к prs.slides (для соответствия slide_type).

    Returns:
        {'total': X, 'slides_touched': Y}
    """
    skip_types = skip_slide_types or {"kpi_native"}
    total = 0
    touched = 0
    slides = list(prs.slides)
    phrases = 0
    for idx, slide in enumerate(slides):
        plan = (plan_slides[idx] or {}) if plan_slides and idx < len(plan_slides) else {}
        st = plan.get("slide_type")
        # Skip native slides where the renderer has already styled numbers.
        if st in skip_types:
            continue
        # D2: phrase emphasis — green on light, strip-only on dark.
        dark = bool(plan.get("dark", False))
        phrases += emphasize_phrases_in_slide(
            slide, emphasis_hex=(None if dark else _GREEN_HEX))
        # Numeric auto-green pass (unchanged).
        n = emphasize_kpi_in_slide(slide)
        if n:
            total += n
            touched += 1
    return {"total": total, "slides_touched": touched, "phrases": phrases}

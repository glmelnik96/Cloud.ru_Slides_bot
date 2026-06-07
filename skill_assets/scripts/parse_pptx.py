#!/usr/bin/env python3
"""
parse_pptx.py — извлекает структуру .pptx в JSON для агентов.

Usage:
    python3 parse_pptx.py <input.pptx> [output.json]

Output:
    JSON со структурой:
    {
      "file": "...",
      "slide_count": N,
      "slide_size": {"width_emu": ..., "height_emu": ...},
      "slides": [
        {
          "num": 1,
          "layout_name": "...",
          "layout_idx_in_master": <int|null>,
          "title": "...",
          "body": ["...", "..."],
          "text_runs": [...],
          "images": [{"name": "...", "left_emu": ..., "top_emu": ..., "width_emu": ..., "height_emu": ...}],
          "shapes_count": N,
          "tables_count": N
        }
      ]
    }
"""
import sys, json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from visual_kind import classify_visual_kind


# XL_CHART_TYPE (python-pptx enum) → our ChartConfig.type literal.
# Keyed by enum .name so we don't import the (large) enum; unknown families
# fall back to "bar" (the safest universal render).
def _chart_type_of(xl_name: str) -> str:
    n = (xl_name or "").upper()
    if n.startswith("PIE") or n.startswith("DOUGHNUT"):
        return "pie"
    if n.startswith("LINE") or n.startswith("XY"):
        return "line"
    if n.startswith("AREA"):
        if "100" in n:
            return "area_100"
        return "area_stacked"
    # COLUMN_* / BAR_* and anything else → bar family
    if "STACKED" in n:
        return "bar_stacked"
    return "bar"


def _num(v):
    """Coerce a chart cell to float; missing/blank → 0.0."""
    try:
        return float(v)
    except (TypeError, ValueError):
        return 0.0


def _extract_chart(shape):
    """Pull a native PPTX chart object into a ChartConfig-shaped dict.

    Returns None when the shape is not a chart or carries no plottable data.
    Categories come from the first plot; each series yields {name, data:[float]}.
    """
    try:
        if not getattr(shape, "has_chart", False):
            return None
        chart = shape.chart
    except Exception:
        return None
    try:
        ctype = _chart_type_of(getattr(chart.chart_type, "name", str(chart.chart_type)))
    except Exception:
        ctype = "bar"
    try:
        plots = list(chart.plots)
        cats = [str(c) for c in plots[0].categories] if plots else []
    except Exception:
        cats = []
    series = []
    try:
        for s in chart.series:
            series.append({
                "name": str(getattr(s, "name", "") or ""),
                "data": [_num(v) for v in s.values],
            })
    except Exception:
        pass
    if not series:
        return None
    title = ""
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip()
    except Exception:
        pass
    return {
        "type": ctype, "title": title, "caption": "",
        "x": cats, "series": series, "accent_idx": 0,
    }


def _walk_shapes(shapes, depth=0):
    """Flatten the shape tree, recursing into GROUP shapes.

    Returns a list of leaf-shape dicts:
      {shape_type, text, left, top, w, h, depth}
    Text is "" for non-text shapes. Geometry is in EMU (None if absent).
    A shape that raises on attribute access is skipped, not fatal.
    """
    out = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                out += _walk_shapes(shape.shapes, depth + 1)
                continue
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
            out.append({
                "shape_type": shape.shape_type,
                "text": text,
                "left": shape.left, "top": shape.top,
                "w": shape.width, "h": shape.height,
                "depth": depth,
            })
        except Exception:
            continue
    return out


def parse(input_path):
    p = Presentation(input_path)
    result = {
        "file": input_path,
        "slide_count": len(p.slides),
        "slide_size": {"width_emu": p.slide_width, "height_emu": p.slide_height,
                       "width_px_at96": int(p.slide_width / 9525), "height_px_at96": int(p.slide_height / 9525)},
        "slides": [],
    }

    # Build layout idx lookup (SlideLayout is not hashable — use id())
    layout_to_idx = {id(lay): i for i, lay in enumerate(p.slide_layouts)}

    for snum, slide in enumerate(p.slides, start=1):
        sdata = {
            "num": snum,
            "layout_name": slide.slide_layout.name,
            "layout_idx_in_master": layout_to_idx.get(id(slide.slide_layout)),
            # Carried so classify_visual_kind can measure an image as a fraction
            # of the actual slide (dominant-raster recovery on hybrid slides).
            "slide_size": result["slide_size"],
            "title": None,
            "body": [],
            "text_runs": [],
            "images": [],
            "shapes_count": 0,
            "tables_count": 0,
            "tables": [],
            "charts": [],
        }

        for shape in slide.shapes:
            sdata["shapes_count"] += 1

            # Title placeholder
            if shape.has_text_frame and shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                txt = shape.text_frame.text.strip()
                if not txt:
                    continue
                # Title placeholder type values: TITLE=13/CTR_TITLE=15
                if str(ph_type) in ("TITLE (13)", "CENTER_TITLE (15)") or "TITLE" in str(ph_type).upper():
                    if not sdata["title"]:
                        sdata["title"] = txt
                    else:
                        sdata["body"].append(txt)
                else:
                    sdata["body"].append(txt)
                continue

            # Generic text frames (not placeholder)
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    sdata["text_runs"].append(txt)
                    if not sdata["title"] and len(txt) < 100:
                        sdata["title"] = txt
                    else:
                        sdata["body"].append(txt)

            # Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sdata["images"].append({
                    "name": shape.name,
                    "left_emu": shape.left, "top_emu": shape.top,
                    "width_emu": shape.width, "height_emu": shape.height,
                })

            # Native charts — extract the plotted data so downstream agents
            # can regenerate a branded chart instead of dropping it (defect D:
            # parse_pptx previously kept only PICTURE shapes, so PPTX chart
            # objects vanished → "0 pictures inserted").
            chart_cfg = _extract_chart(shape)
            if chart_cfg is not None:
                sdata["charts"].append(chart_cfg)
                head = chart_cfg["title"] or "Диаграмма"
                snames = ", ".join(s["name"] for s in chart_cfg["series"] if s["name"])
                sdata["body"].append(
                    f"{head}: {chart_cfg['type']} "
                    f"({len(chart_cfg['x'])} категорий"
                    + (f", ряды: {snames}" if snames else "") + ")"
                )

            # Tables — extract real cell text so downstream agents can render
            # the actual data (previously only the count was kept, so table
            # content was silently dropped → donor-53 PNG-stub placeholder).
            if shape.has_table:
                sdata["tables_count"] += 1
                tbl = shape.table
                grid = []
                merged = False
                for row in tbl.rows:
                    cells = list(row.cells)
                    row_vals = [c.text.strip() for c in cells]
                    grid.append(row_vals)
                    if any(getattr(c, "is_spanned", False)
                           or getattr(c, "is_merge_origin", False) for c in cells):
                        merged = True
                widths = {len(r) for r in grid}
                regular = (
                    len(grid) >= 2
                    and len(widths) == 1
                    and next(iter(widths)) >= 2
                    and not merged
                )
                sdata["tables"].append({
                    "headers": grid[0] if grid else [],
                    "rows": grid[1:] if len(grid) > 1 else [],
                    "regular": regular,
                })
                # Surface the table content as text so vision/LLM nodes that
                # read `body` also see it (pipe-delimited, one row per line).
                for r in grid:
                    line = " | ".join(v for v in r if v)
                    if line:
                        sdata["body"].append(line)

        # --- group recovery: pull text + pictures buried inside GROUP shapes
        leaves = _walk_shapes(slide.shapes)
        group_nodes = []
        order = 0
        buried_text = []
        for lf in leaves:
            if lf["depth"] >= 1 and lf["text"]:
                order += 1
                group_nodes.append({
                    "text": lf["text"],
                    "left": lf["left"], "top": lf["top"],
                    "w": lf["w"], "h": lf["h"],
                    "order": order,
                })
                buried_text.append(lf["text"])
            # pictures inside groups were missed by the top-level loop
            if lf["depth"] >= 1 and lf["shape_type"] == MSO_SHAPE_TYPE.PICTURE:
                sdata["images"].append({
                    "name": "group_pic",
                    "left_emu": lf["left"], "top_emu": lf["top"],
                    "width_emu": lf["w"], "height_emu": lf["h"],
                })
        sdata["group_nodes"] = group_nodes
        # classify BEFORE surfacing buried text into body, otherwise a structured
        # diagram slide (no native title/body) would look like a normal text slide.
        sdata["visual_kind"] = classify_visual_kind(sdata)
        # now surface buried text to body so vision/LLM can see it
        sdata["body"].extend(buried_text)

        result["slides"].append(sdata)

    return result


def main():
    if len(sys.argv) < 2:
        print("Usage: parse_pptx.py <input.pptx> [output.json]", file=sys.stderr)
        sys.exit(1)
    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    data = parse(input_path)
    out = json.dumps(data, ensure_ascii=False, indent=2)

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(out)
        print(f"Wrote {output_path} ({data['slide_count']} slides)", file=sys.stderr)
    else:
        print(out)


if __name__ == "__main__":
    main()

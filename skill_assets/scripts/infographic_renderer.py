#!/usr/bin/env python3
"""
infographic_renderer.py — рендер native_block shapes от infographic_maker.

Зачем: Agent 06 (infographic_maker) производит список абсолютно
позиционированных shapes (rounded_rect + text) — описание визуального
сравнения / списка / процесса / etc. До этого build_v9 хендлил только
``slide_type in (kpi_native, table_native, ...)`` и игнорировал
``plan_slide.infographic.shapes`` — 9 шейпов на слайд молча выкидывались,
и слайд рендерился как «голый текст» в донорском body-плейсхолдере.
Это блокер plan_compliance в visual_verifier (2026-06-05 live run).

Shape spec (как эмиттит Agent 06):
    {
      "type": "rounded_rect" | "text",
      "left_emu", "top_emu", "width_emu", "height_emu",   # int EMU
      "fill_color": "#RRGGBB" | "none",
      "stroke_color": "#RRGGBB" | "none",
      "stroke_width_pt": float,
      "text": str,                                         # может быть ""
      "font": str,                                         # "SB Sans Display"…
      "font_size_pt": int,
      "font_color": "#RRGGBB"
    }
"""
from __future__ import annotations

import sys
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from text_sanitize import sanitize_text

try:
    import textfit as _textfit
    import font_resolver as _font_resolver
    GEOFIT_AVAILABLE = True
except ImportError:
    GEOFIT_AVAILABLE = False


_GRAPHITE = RGBColor(0x22, 0x22, 0x22)

# Text-frame insets applied in _set_textframe (EMU); subtract from the box so
# the geometric fitter measures the actual usable area.
_TF_MARGIN_X_EMU = 60000
_TF_MARGIN_Y_EMU = 20000
# Floor for infographic card text — below this it's unreadable; better to let a
# pathological string clip than render at 6pt.
_INFOGRAPHIC_MIN_PT = 10.0


def _fit_infographic_size(spec: dict[str, Any], width_emu: int, height_emu: int):
    """Shrink a card text's font so it fits its box height/width.

    Agent 06 sizes card-title and body boxes independently; a long card title
    wraps to 3 lines and overflows downward into the body box below it (the
    "title/body overlap" defect). Measuring against the real box height and
    shrinking keeps each text contained. Returns the (possibly reduced)
    font_size_pt, or the original on any failure."""
    base = spec.get("font_size_pt")
    if not GEOFIT_AVAILABLE or not base:
        return base
    text = spec.get("text") or ""
    if not text.strip() or width_emu <= 0 or height_emu <= 0:
        return base
    font_path = _font_resolver.resolve(spec.get("font"), False)
    if not font_path:
        return base
    res = _textfit.fit_text(
        text,
        box_w_emu=max(1, width_emu - 2 * _TF_MARGIN_X_EMU),
        box_h_emu=max(1, height_emu - 2 * _TF_MARGIN_Y_EMU),
        font_path=font_path,
        base_pt=float(base),
        min_pt=_INFOGRAPHIC_MIN_PT,
        wrap=True,
        balance=False,
    )
    if res is None:
        return base
    return res.size_pt


def _parse_hex(value: Any) -> RGBColor | None:
    """'#26D07C' / '26D07C' → RGBColor; 'none' / '' / None → None."""
    if not value or not isinstance(value, str):
        return None
    s = value.strip().lstrip("#")
    if not s or s.lower() == "none":
        return None
    if len(s) != 6:
        return None
    try:
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return None


def _apply_fill(shape, fill_color: Any) -> None:
    color = _parse_hex(fill_color)
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def _apply_line(shape, stroke_color: Any, width_pt: Any) -> None:
    color = _parse_hex(stroke_color)
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    try:
        w = float(width_pt or 0)
    except (TypeError, ValueError):
        w = 0.0
    if w > 0:
        shape.line.width = Pt(w)


def _set_textframe(tf, text: str, font: str | None, size_pt: Any,
                   color_hex: Any, *, bold: bool = False) -> None:
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # python-pptx adds a default empty paragraph; reuse it.
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    # Clear any pre-existing runs (cheap — paragraph just created).
    for r in list(p.runs):
        r.text = ""
    run = p.add_run() if not p.runs else p.runs[0]
    # strip_markdown=False: infographic overlays sit on regular slides that the
    # apply_kpi_emphasis phrase pass still visits, which consumes ``**…**``
    # itself — only remove control chars here (fixes _X000B_), leave ** to it.
    run.text = sanitize_text(text or "", strip_markdown=False)
    if font:
        run.font.name = font
    try:
        size = int(size_pt) if size_pt is not None else None
    except (TypeError, ValueError):
        size = None
    if size and size > 0:
        run.font.size = Pt(size)
    color = _parse_hex(color_hex) or _GRAPHITE
    run.font.color.rgb = color
    if bold:
        run.font.bold = True


def _emu_bounds(spec: dict[str, Any]) -> tuple[Emu, Emu, Emu, Emu]:
    left = Emu(int(spec.get("left_emu", 0) or 0))
    top = Emu(int(spec.get("top_emu", 0) or 0))
    width = Emu(int(spec.get("width_emu", 0) or 0))
    height = Emu(int(spec.get("height_emu", 0) or 0))
    return left, top, width, height


def _add_filled_shape(slide, spec: dict[str, Any], mso_shape: int,
                     *, rounded_adjust: float | None = None) -> None:
    """Generic add-shape with fill/line/text dispatch.

    Used by rectangle / rounded_rect / circle / arrow — they share the
    same fill+line+optional-text contract; only the MSO_SHAPE constant
    differs. ``rounded_adjust`` overrides the default corner-radius
    fraction (rounded_rect only).
    """
    left, top, width, height = _emu_bounds(spec)
    shape = slide.shapes.add_shape(mso_shape, left, top, width, height)
    if rounded_adjust is not None:
        try:
            shape.adjustments[0] = rounded_adjust
        except Exception:  # noqa: BLE001
            pass
    _apply_fill(shape, spec.get("fill_color"))
    _apply_line(shape, spec.get("stroke_color"), spec.get("stroke_width_pt"))
    text = spec.get("text") or ""
    if text:
        _set_textframe(
            shape.text_frame, text,
            font=spec.get("font"),
            size_pt=_fit_infographic_size(spec, int(width), int(height)),
            color_hex=spec.get("font_color"),
            bold=False,
        )
        tf = shape.text_frame
        tf.margin_left = Emu(60000)
        tf.margin_right = Emu(60000)
        tf.margin_top = Emu(20000)
        tf.margin_bottom = Emu(20000)


def _add_rounded_rect(slide, spec: dict[str, Any]) -> None:
    _add_filled_shape(slide, spec, MSO_SHAPE.ROUNDED_RECTANGLE,
                     rounded_adjust=0.10)


def _add_rectangle(slide, spec: dict[str, Any]) -> None:
    _add_filled_shape(slide, spec, MSO_SHAPE.RECTANGLE)


def _add_circle(slide, spec: dict[str, Any]) -> None:
    _add_filled_shape(slide, spec, MSO_SHAPE.OVAL)


def _add_arrow(slide, spec: dict[str, Any]) -> None:
    """Default to RIGHT_ARROW (process flow l→r). Agent 06 doesn't specify
    direction, but its `process` infographic_type uses horizontal step
    layouts where arrows go left→right between boxes."""
    _add_filled_shape(slide, spec, MSO_SHAPE.RIGHT_ARROW)


def _add_line(slide, spec: dict[str, Any]) -> None:
    """Render as a thin rectangle (height = stroke_width). python-pptx
    Connector API is brittle for absolute positions in this codebase,
    and Agent 06 emits line shapes with full bounding boxes anyway."""
    left, top, width, height = _emu_bounds(spec)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    # Lines have no fill — colour comes from stroke_color (use as fill).
    stroke_hex = spec.get("stroke_color") or spec.get("fill_color")
    color = _parse_hex(stroke_hex)
    if color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    shape.line.fill.background()


def _add_textbox(slide, spec: dict[str, Any]) -> None:
    left = Emu(int(spec.get("left_emu", 0) or 0))
    top = Emu(int(spec.get("top_emu", 0) or 0))
    width = Emu(int(spec.get("width_emu", 0) or 0))
    height = Emu(int(spec.get("height_emu", 0) or 0))
    box = slide.shapes.add_textbox(left, top, width, height)
    _apply_fill(box, spec.get("fill_color"))
    # Textboxes generally render lines as nuisance; only draw if explicit.
    _apply_line(box, spec.get("stroke_color"), spec.get("stroke_width_pt"))
    _set_textframe(
        box.text_frame, spec.get("text") or "",
        font=spec.get("font"),
        size_pt=_fit_infographic_size(spec, int(width), int(height)),
        color_hex=spec.get("font_color"),
    )


_HANDLERS = {
    "rounded_rect": _add_rounded_rect,
    "rectangle":    _add_rectangle,
    "circle":       _add_circle,
    "arrow":        _add_arrow,
    "line":         _add_line,
    "text":         _add_textbox,
}


# P1-1 (2026-06-05): Agent 06's `process` infographic regularly overshoots
# the safe-area on 4-block layouts. Live run3.slide2 produced four 330-px
# blocks starting at x=30 — total reach 1410 px on a 1280-px canvas, so
# the last block ("Ожидание") clipped 160 px past the right edge. The
# prompt now states block_width formulas but the LLM still emits raw
# numbers. We post-validate here so visual quality doesn't depend on
# perfect prompt adherence.
#
# Canvas + safe-area are duplicated from llm/prompts/_shared.py — both
# modules ship in the skill bundle and we keep imports local to avoid
# a runtime dep on the prompts package.
_EMU_PER_PX = 9525
_SAFE_AREA_PX = {"left": 30, "right": 1250, "top": 140, "bottom": 660}
_SAFE_AREA_EMU = {
    k: v * _EMU_PER_PX for k, v in _SAFE_AREA_PX.items()
}


# Task 5 (2026-06-07): process/flow item cap. A `process`/`flow`
# infographic is a horizontal row of N cards (blocks) with N-1 arrows
# between them. The cards must fit inside the safe-area width:
#   block_width = (safe_w_px - (N-1)*gap) / N   with gap = 60 px
# safe_w = 1250-30 = 1220 px. At N=8 → block_width = (1220-420)/8 = 100 px,
# still wide enough for a short label at the 10pt floor. At N=9 the block
# drops to 82 px (and the deck's recovered-body path produced 10-11 cards
# in the live failure, clipping the bottom cards off-slide). So 8 is the
# realistic layout capacity. Cards beyond the cap are NOT dropped — their
# text is merged into the last shown card so no source word is clipped.
#
# `flow` shares the same horizontal step layout family as `process`; both
# are capped. `comparison`/`matrix`/`tree`/`chart_*`/`none` are different
# layouts and are left untouched.
_PROCESS_MAX_ITEMS = 8
# Shape types that act as a "card" (carry a step label) vs. connectors.
# Scope boundary: standalone ``type:"text"`` shapes are NOT counted as
# cards, so a label-only step does not consume cap budget.
_CARD_SHAPE_TYPES = ("rounded_rect", "rectangle", "circle")
_CONNECTOR_SHAPE_TYPES = ("arrow", "line")
_CAPPED_INFOGRAPHIC_TYPES = ("process", "flow")


def cap_process_items(
    infographic_type: str | None,
    shapes: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    """Cap a process/flow infographic's cards to the layout capacity.

    Called at the feed/distribution point (assemble_plan_node) BEFORE the
    spec reaches the renderer, so the renderer always sees ≤ cap cards.

    A card is a block shape (``rounded_rect``/``rectangle``/``circle``)
    carrying step text; ``arrow``/``line`` shapes are connectors. When the
    card count exceeds ``_PROCESS_MAX_ITEMS`` we keep the first
    ``cap-1`` cards verbatim and merge the text of ALL overflow cards into
    the ``cap``-th (last shown) card — every word is preserved, nothing is
    silently clipped. Orphan connectors that would have pointed at dropped
    cards are removed so no arrows dangle.

    Returns a NEW list when capping occurs; the original list (unchanged)
    when the type isn't a step layout or the card count is ≤ cap.
    Card text-merge mutates only the surviving last card's ``text``.
    """
    if infographic_type not in _CAPPED_INFOGRAPHIC_TYPES:
        return shapes
    if not shapes:
        return shapes

    # Indices of card shapes, in render order.
    card_idxs = [
        i for i, s in enumerate(shapes)
        if isinstance(s, dict) and s.get("type") in _CARD_SHAPE_TYPES
    ]
    if len(card_idxs) <= _PROCESS_MAX_ITEMS:
        return shapes

    keep_idxs = card_idxs[:_PROCESS_MAX_ITEMS]
    overflow_idxs = card_idxs[_PROCESS_MAX_ITEMS:]
    last_keep_idx = keep_idxs[-1]

    # Merge every overflow card's text into the last kept card so no word
    # is clipped. Join with " · " (the deck already uses mid-dot separators
    # for compacted steps) and skip empties / duplicates of the carrier.
    last_card = shapes[last_keep_idx]
    parts: list[str] = []
    base = (last_card.get("text") or "").strip()
    if base:
        parts.append(base)
    for oi in overflow_idxs:
        t = (shapes[oi].get("text") or "").strip()
        if t and t not in parts:
            parts.append(t)
    merged_text = " · ".join(parts)

    overflow_set = set(overflow_idxs)
    # Connectors are only meaningful between kept cards. The arrow that
    # would have linked the last kept card to the first dropped card now
    # points into empty space, so drop any arrow/line that starts at or
    # beyond the last kept card's left edge (that arrow and every later
    # one). Earlier arrows sit between two surviving cards and are kept.
    # NOTE: this assumes shapes are emitted monotonically left-to-right,
    # which is the Agent 06 process/flow pattern (cards laid out in a row).
    # A MISSING or non-numeric ``left_emu`` on the last kept card means we
    # have no positional anchor — leave ``drop_after_emu=None`` so NO
    # connector is pruned (better to keep an arrow than prune them all by
    # treating an absent position as 0).
    raw_left = last_card.get("left_emu")
    if raw_left is None:
        drop_after_emu = None
    else:
        try:
            drop_after_emu = int(raw_left)
        except (TypeError, ValueError):
            drop_after_emu = None

    result: list[dict[str, Any]] = []
    for i, s in enumerate(shapes):
        if i in overflow_set:
            continue  # overflow card text already merged
        if i == last_keep_idx:
            # Carry the merged text on the surviving last card.
            new_card = dict(s)
            new_card["text"] = merged_text
            result.append(new_card)
            continue
        # Prune orphan connectors that would point past the last kept card.
        if (
            isinstance(s, dict)
            and s.get("type") in _CONNECTOR_SHAPE_TYPES
            and drop_after_emu is not None
        ):
            try:
                left = int(s.get("left_emu", 0) or 0)
            except (TypeError, ValueError):
                left = -1
            if left >= drop_after_emu:
                continue
        result.append(s)
    return result


# If horizontal bounding span occupies less than this fraction of the
# safe-area, treat as undersize (GLM-5.1 hallucination from 2026-06-05
# run 29e189bb where 3 columns spanned 179 px on a 1220-px safe area —
# 14.7% utilization). Anything below 50% is almost certainly a wrong-unit
# emit; expand to fill.
_UNDERSCALE_THRESHOLD = 0.50
# When expanding, leave a small margin so corners breathe.
_UPSCALE_TARGET_FRAC = 0.95


def _clamp_shapes_to_safe_area(shapes: list[dict[str, Any]]) -> int:
    """Rescale a list of Agent 06 shape specs so they fit the safe-area
    properly. Returns count of shapes mutated.

    Two failure modes from live runs:
      - **Overshoot** (run3.slide2, 2026-06-05): bounding span > safe
        width → scale down + shift left.
      - **Underscale** (run 29e189bb, 2026-06-05): bounding span <
        50% of safe width → GLM hallucinated wrong-unit values; scale
        up to ~95% of safe width and shift to safe.left.

    Strategy: compute horizontal bounding span over rectangle-class
    shapes. Decide overshoot vs underscale vs ok. Apply a single affine
    transform (scale + shift) to every shape's ``left_emu``,
    ``width_emu``. Underscale ALSO scales Y (top_emu / height_emu)
    because the same hallucinated unit affects both axes.

    Mutates the spec dicts in place.
    """
    if not shapes:
        return 0
    spans = []
    for spec in shapes:
        if not isinstance(spec, dict):
            continue
        try:
            l = int(spec.get("left_emu", 0) or 0)
            w = int(spec.get("width_emu", 0) or 0)
        except (TypeError, ValueError):
            continue
        if w <= 0:
            continue
        spans.append((l, l + w, spec))
    if not spans:
        return 0
    min_left = min(s[0] for s in spans)
    max_right = max(s[1] for s in spans)
    safe_left = _SAFE_AREA_EMU["left"]
    safe_right = _SAFE_AREA_EMU["right"]
    safe_top = _SAFE_AREA_EMU["top"]
    safe_bottom = _SAFE_AREA_EMU["bottom"]
    safe_w = safe_right - safe_left
    safe_h = safe_bottom - safe_top
    current_w = max_right - min_left
    if current_w <= 0 or safe_w <= 0:
        return 0

    # ── classify mode ─────────────────────────────────────────────────
    inside = (min_left >= safe_left and max_right <= safe_right)
    overshoot = max_right > safe_right or min_left < safe_left
    underscale = inside and (current_w < safe_w * _UNDERSCALE_THRESHOLD)
    if not overshoot and not underscale:
        return 0  # already well-placed

    mutated = 0
    if underscale:
        # GLM-5.1 hallucination: shapes are ~5-10x too small. Scale both
        # axes by the same factor so circles stay circular, then shift to
        # safe-top-left. Also scale font_size_pt proportionally so text
        # doesn't dominate tiny boxes after upscale.
        scale = (safe_w * _UPSCALE_TARGET_FRAC) / current_w
        # Compute Y span across all shapes (use full list, including
        # arrows whose width might be 0 if vertical).
        ys = []
        for spec in shapes:
            if not isinstance(spec, dict):
                continue
            try:
                t = int(spec.get("top_emu", 0) or 0)
                h = int(spec.get("height_emu", 0) or 0)
            except (TypeError, ValueError):
                continue
            if h <= 0:
                continue
            ys.append((t, t + h, spec))
        min_top = min((y[0] for y in ys), default=safe_top)
        max_bot = max((y[1] for y in ys), default=safe_top)
        current_h = max(1, max_bot - min_top)
        # Cap Y scale so it doesn't overshoot safe-height.
        scale_y = min(scale, (safe_h * _UPSCALE_TARGET_FRAC) / current_h)
        for spec in shapes:
            if not isinstance(spec, dict):
                continue
            try:
                l = int(spec.get("left_emu", 0) or 0)
                w = int(spec.get("width_emu", 0) or 0)
                t = int(spec.get("top_emu", 0) or 0)
                h = int(spec.get("height_emu", 0) or 0)
            except (TypeError, ValueError):
                continue
            new_left = int(round(safe_left + (l - min_left) * scale))
            new_width = max(1, int(round(w * scale))) if w > 0 else w
            new_top = int(round(safe_top + (t - min_top) * scale_y))
            new_height = max(1, int(round(h * scale_y))) if h > 0 else h
            changed = False
            if new_left != l:
                spec["left_emu"] = new_left; changed = True
            if new_width != w and w > 0:
                spec["width_emu"] = new_width; changed = True
            if new_top != t:
                spec["top_emu"] = new_top; changed = True
            if new_height != h and h > 0:
                spec["height_emu"] = new_height; changed = True
            # Scale text up proportionally too, but cap at 24pt.
            try:
                size_pt = float(spec.get("font_size_pt") or 0)
            except (TypeError, ValueError):
                size_pt = 0.0
            if size_pt > 0:
                # Use sqrt-ish scaling so text doesn't explode on 10x boxes.
                new_size = min(24, max(10, int(round(size_pt * min(scale, 2.5)))))
                if new_size != int(size_pt):
                    spec["font_size_pt"] = new_size
                    changed = True
            if changed:
                mutated += 1
        if mutated:
            print(
                f"infographic upscale: span {current_w//_EMU_PER_PX}px → "
                f"{int(current_w*scale)//_EMU_PER_PX}px (scale={scale:.2f}x, "
                f"scale_y={scale_y:.2f}x), shapes mutated={mutated}",
                file=sys.stderr,
            )
        return mutated

    # ── overshoot path (legacy) ──────────────────────────────────────
    scale = min(1.0, safe_w / current_w)
    for left_emu, right_emu, spec in spans:
        new_left = int(round(safe_left + (left_emu - min_left) * scale))
        new_width = max(1, int(round((right_emu - left_emu) * scale)))
        if new_left != spec.get("left_emu") or new_width != spec.get("width_emu"):
            spec["left_emu"] = new_left
            spec["width_emu"] = new_width
            mutated += 1
    if mutated:
        print(
            f"infographic clamp: span {current_w//_EMU_PER_PX}px "
            f"→ safe-area {safe_w//_EMU_PER_PX}px (scale={scale:.3f}), "
            f"shapes mutated={mutated}",
            file=sys.stderr,
        )
    return mutated


def render_infographic_shapes(slide, shapes: list[dict[str, Any]]) -> int:
    """Inject Agent 06 shape specs onto an existing (cloned) donor slide.

    Returns the count of shapes successfully added. Unknown types are
    logged and skipped — we never raise, because losing one shape is
    better than failing the entire build.
    """
    if not shapes:
        return 0
    # P1-1: clamp the Agent 06 spec list to safe-area before rendering.
    # Cheap, idempotent — never expands shapes that already fit.
    try:
        _clamp_shapes_to_safe_area(shapes)
    except Exception as e:  # noqa: BLE001 — never fail the build
        print(f"WARN: clamp_shapes_to_safe_area failed: {e}", file=sys.stderr)
    added = 0
    for i, spec in enumerate(shapes):
        if not isinstance(spec, dict):
            print(f"WARN: infographic shape #{i} is not a dict ({type(spec).__name__})",
                  file=sys.stderr)
            continue
        kind = spec.get("type")
        handler = _HANDLERS.get(kind)
        if handler is None:
            print(f"WARN: infographic shape #{i} unknown type {kind!r}",
                  file=sys.stderr)
            continue
        try:
            handler(slide, spec)
            added += 1
        except Exception as e:  # noqa: BLE001 — never fail the build
            print(f"WARN: infographic shape #{i} ({kind}) failed: {e}",
                  file=sys.stderr)
    return added


# Slots that infographic shapes are expected to replace on the donor.
# Title stays (the infographic_maker doesn't reproduce slide titles).
# Anything else with text → cleared so the body doesn't bleed through
# under translucent infographic blocks.
_BODY_SLOT_KEYWORDS = (
    "body", "content", "caption", "subtitle", "description", "desc",
    "text", "list", "card", "col1", "col2", "col3", "col4", "col5", "col6",
    "annotation", "left", "right", "top", "bottom", "center",
)


def clear_donor_body_slots(slide, donor_def: dict[str, Any] | None) -> int:
    """Empty text frames of donor body-like slots — leaves title intact.

    Called when infographic shapes are about to be injected so the
    donor's pre-rendered body text doesn't show through between/around
    the new blocks. Returns count of slots cleared.
    """
    if not donor_def:
        return 0
    slots = donor_def.get("slots") or {}
    if not isinstance(slots, dict):
        return 0

    # Need lazy import — build_v5 sits next to this file and pulls pptx
    # at import time, which is fine inside the package but not in
    # isolated tests that don't ship build_v5.
    from build_v5 import get_text_frame_by_shape_idx, clear_text_frame

    cleared = 0
    for slot_name, slot_def in slots.items():
        name_low = (slot_name or "").lower()
        # Keep title — infographic_maker doesn't reproduce slide titles.
        if "title" in name_low and "subtitle" not in name_low:
            continue
        if not any(kw in name_low for kw in _BODY_SLOT_KEYWORDS):
            continue
        if not isinstance(slot_def, dict):
            continue
        idx = slot_def.get("shape_idx")
        if idx is None:
            continue
        tf = get_text_frame_by_shape_idx(slide, idx)
        if tf is None:
            continue
        try:
            clear_text_frame(tf)
            cleared += 1
        except Exception as e:  # noqa: BLE001
            print(f"WARN: clear donor slot {slot_name!r} (idx={idx}) failed: {e}",
                  file=sys.stderr)
    return cleared


# D1+D8 fix (2026-06-05): when injecting Agent 06's infographic shapes onto
# a donor template, the donor often has its own pre-labeled boxes (process
# steps, comparison cells) whose text isn't in the donor slot map — it lives
# as plain text inside shapes that decorate the layout. The previous
# `clear_donor_body_slots` only touched slot-mapped shapes, so donor labels
# bled through and overlapped Agent 06's new labels (run1.slide7: "Recorder"
# overlap "Хранение данных"). This is a more aggressive pass: clear *all*
# non-title text on the slide, regardless of slot mapping.

# B2 (2026-06-05): tightened. Live run a337cc86 slides 7/9 showed donor 33's
# decoration sub-headers "Подзаголовок в две строки" bleeding through behind
# infographic boxes even though clear_donor_non_title_text reported
# non_title_cleared=6. Those sub-headers are 18-24pt — they sneaked past the
# old `>= 18.0` threshold which treated them as titles.
#
# New rule: ONLY the slide TITLE placeholder is preserved. Every other text
# shape (regardless of font size) gets cleared when infographic shapes are
# about to be painted on top. The actual slide title is uniquely identified
# by its placeholder type; donor decoration "sub-titles" are NOT placeholder
# titles even if they look big.
def _shape_is_title_like(shape) -> bool:
    """True only for the slide's TITLE placeholder. Donor decoration shapes
    with large fonts are NOT treated as titles — those caused live bleed."""
    try:
        ph = shape.placeholder_format
        if ph is not None:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if ph.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                           PP_PLACEHOLDER.VERTICAL_TITLE):
                return True
    except (ValueError, AttributeError):
        pass
    return False


def clear_donor_non_title_text(slide, preserve_shape_idx=None) -> int:
    """Strip text from every non-title shape on the slide.

    Run before injecting infographic shapes so donor decoration labels
    (process-step names, comparison cells, etc.) don't overlap the new
    Agent 06 boxes. Title remains; everything else loses its text but
    keeps its visual shape (fill/stroke). Recurses into groups and
    table cells so the cleanup is complete.

    F1 (2026-06-05): added ``preserve_shape_idx`` so Case A in build_v9
    (structural donor with real filled body slots, overlay dropped by B5)
    can wipe non-slot decoration without erasing the just-filled slot
    text. Indices are top-level shape positions in ``slide.shapes``;
    only top-level shapes consult the preserve set — nested group
    children are never slot-mapped and always cleared.

    Returns count of shapes whose text was cleared.
    """
    cleared = 0
    preserve = set(preserve_shape_idx or ())
    try:
        from build_v5 import clear_text_frame
    except ImportError:
        return 0

    def _process(sh):
        nonlocal cleared
        # Groups: recurse (no preservation; slots are never inside groups).
        try:
            if sh.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                for child in sh.shapes:
                    _process(child)
                return
        except Exception:  # noqa: BLE001
            pass
        # Tables: blank every cell text frame (table itself is decoration).
        if getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    try:
                        clear_text_frame(cell.text_frame)
                        cleared += 1
                    except Exception:  # noqa: BLE001
                        pass
            return
        if not getattr(sh, "has_text_frame", False):
            return
        if _shape_is_title_like(sh):
            return
        text = (sh.text_frame.text or "").strip()
        if not text:
            return
        try:
            clear_text_frame(sh.text_frame)
            cleared += 1
        except Exception as e:  # noqa: BLE001
            print(f"WARN: clear_donor_non_title_text failed: {e}",
                  file=sys.stderr)

    for idx, sh in enumerate(slide.shapes):
        if idx in preserve:
            continue
        _process(sh)
    return cleared

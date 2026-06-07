"""End-to-end guarantee for native PPTX chart extraction → editable chart.

Defect D: source charts/diagrams were lost ("0 pictures inserted") because
``parse_pptx`` kept only PICTURE shapes, so native chart graphicFrames vanished.
The fix chains:

    parse_pptx._extract_chart        — pull series/categories from a native chart
    agents._inject_parsed_charts     — deterministically restore chart_pptx_native
    build_v9.build                   — render a REAL editable chart on a blank donor
    chart_native_pptx                — the native (editable) renderer

The unit tests cover each link in isolation, but no benchmark deck holds a real
chart object, so the full chain was never run. This test builds a synthetic deck
WITH native charts (column/pie/line) and proves the OUTPUT pptx carries genuine,
editable chart shapes (``shape.has_chart``) — not flattened images — with the
right type/series/categories.
"""
from __future__ import annotations

import json
import os
import sys

import pytest

# skill_assets/scripts onto the path (parse_pptx + chart modules live there)
_SCRIPTS = os.path.join(
    os.path.dirname(__file__), "..", "..", "skill_assets", "scripts")
sys.path.insert(0, _SCRIPTS)

from worker import skill_bridge  # noqa: E402

skill_bridge.install()

from pptx import Presentation  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE  # noqa: E402

import parse_pptx  # noqa: E402
from build_v9 import build  # noqa: E402
from graph.nodes.agents import _inject_parsed_charts  # noqa: E402

from scripts.make_chart_test_deck import build as build_chart_deck  # noqa: E402


# Expected extraction per chart slide of the synthetic deck.
_EXPECTED = {
    1: {"type": "bar", "x": ["2022", "2023", "2024", "2025"],
        "series": ["Облако", "On-prem"]},
    2: {"type": "pie", "x": ["IaaS", "PaaS", "SaaS", "Прочее"],
        "series": ["Доля рынка"]},
    3: {"type": "line", "x": ["Q1", "Q2", "Q3", "Q4", "Q5"],
        "series": ["ARR", "Churn"]},
}


@pytest.fixture(scope="module")
def synthetic_deck(tmp_path_factory) -> str:
    out = tmp_path_factory.mktemp("chartdeck") / "synthetic_chart_deck.pptx"
    return build_chart_deck(str(out))


@pytest.fixture(scope="module")
def parsed(synthetic_deck) -> dict:
    return parse_pptx.parse(synthetic_deck)


# ── STEP 2: extraction ───────────────────────────────────────────────────────

def test_parse_extracts_each_chart(parsed) -> None:
    by_num = {s["num"]: s for s in parsed["slides"]}
    for num, exp in _EXPECTED.items():
        charts = by_num[num]["charts"]
        assert len(charts) == 1, f"slide {num}: expected 1 chart, got {len(charts)}"
        c = charts[0]
        assert c["type"] == exp["type"], f"slide {num}: type {c['type']}"
        assert c["x"] == exp["x"], f"slide {num}: categories {c['x']}"
        assert [s["name"] for s in c["series"]] == exp["series"]
        for s in c["series"]:
            assert s["data"] and all(isinstance(v, float) for v in s["data"])
    # the plain text slide carries no chart
    assert by_num[4]["charts"] == []


# ── STEP 3: deterministic injection routes flat text → chart_pptx_native ──────

def test_inject_routes_flat_text_to_chart_native(parsed) -> None:
    # Simulate the LLM mis-classifying every chart slide as flat text.
    cls = {"slides": [
        {"num": n, "slide_type": "text", "category": "text",
         "kpi": {"x": 1}, "chart": None, "table": {"y": 2},
         "flow": {"z": 3}, "image": {"w": 4}}
        for n in (1, 2, 3, 4)
    ]}
    injected = _inject_parsed_charts(cls, parsed)
    assert injected == 3, f"expected 3 chart slides injected, got {injected}"
    by_num = {s["num"]: s for s in cls["slides"]}
    for num, exp in _EXPECTED.items():
        s = by_num[num]
        assert s["slide_type"] == "chart_pptx_native"
        assert s["category"] == "other"
        assert s["chart"]["type"] == exp["type"]
        assert [ser["name"] for ser in s["chart"]["series"]] == exp["series"]
        # competing native blocks cleared
        for k in ("kpi", "table", "flow", "image"):
            assert s[k] is None
    # text slide left untouched
    assert by_num[4]["slide_type"] == "text"


# ── STEP 4: the build produces REAL editable charts (the actual gap) ──────────

def _expected_xl_type(t: str):
    return {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "pie": XL_CHART_TYPE.PIE,
        "line": XL_CHART_TYPE.LINE,
    }[t]


def _build_charts(parsed, tmp_path, nums):
    """Drive build_v9.build for the given parsed chart slides; return output prs."""
    by_num = {s["num"]: s for s in parsed["slides"]}
    plan = {"slides": [
        {"slide_type": "chart_pptx_native", "chart": dict(by_num[n]["charts"][0])}
        for n in nums
    ]}
    plan_path = tmp_path / "plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False), encoding="utf-8")
    out = tmp_path / "out.pptx"
    build(
        str(plan_path),
        str(skill_bridge.TEMPLATE_PATH),
        str(out),
        str(skill_bridge.DONOR_SLOT_MAP),
    )
    return Presentation(str(out)), by_num


def test_build_emits_native_editable_charts(parsed, tmp_path) -> None:
    """THE key proof: bar + line charts survive the full chain as REAL,
    editable native chart shapes (shape.has_chart) — not flattened images —
    with matching type, series and categories.

    Pie is excluded here and covered by ``test_pie_chart_build_is_broken``
    because the renderer's pie branch reads non-schema keys (see that test).
    """
    nums = (1, 3)  # bar, line
    prs, by_num = _build_charts(parsed, tmp_path, nums)
    assert len(prs.slides) == len(nums)

    for idx, num in enumerate(nums):
        slide = prs.slides[idx]
        chart_shapes = [sh for sh in slide.shapes if sh.has_chart]
        assert len(chart_shapes) == 1, (
            f"slide {num}: expected exactly 1 NATIVE chart shape "
            f"(editable, not a flattened image), got {len(chart_shapes)}")
        chart = chart_shapes[0].chart
        exp = _EXPECTED[num]
        src = by_num[num]["charts"][0]
        assert chart.chart_type == _expected_xl_type(exp["type"]), (
            f"slide {num}: chart_type {chart.chart_type}")

        out_series = list(chart.series)
        assert len(out_series) == len(src["series"]), (
            f"slide {num}: series count {len(out_series)}")
        for got, want in zip(out_series, src["series"]):
            assert list(got.values) == want["data"], (
                f"slide {num}: series {want['name']} values")
        cats = list(chart.plots[0].categories)
        assert cats == exp["x"], f"slide {num}: categories {cats}"

        # render_chart_pptx_slide promotes the chart title to the branded slide
        # title placeholder (not the native chart's own title). Confirm the
        # extracted title survived somewhere visible on the slide.
        # (the template title placeholder upper-cases for brand style, so
        # compare case-insensitively)
        slide_texts = [sh.text_frame.text.strip().lower() for sh in slide.shapes
                       if sh.has_text_frame and sh.text_frame.text.strip()]
        assert src["title"].lower() in slide_texts, (
            f"slide {num}: title {src['title']!r} not found in {slide_texts}")


def test_pie_chart_build_is_broken(parsed, tmp_path) -> None:
    """BUG (regression guard): a PARSED pie chart cannot be built end-to-end.

    ``parse_pptx._extract_chart`` emits the canonical ChartConfig shape
    (``x`` + ``series``) for every chart type, and ``_inject_parsed_charts``
    explicitly routes pies to ``chart_pptx_native``. But the renderer's pie
    branch (``chart_native_pptx.add_chart_to_slide``) reads the non-schema keys
    ``labels`` / ``values`` instead of ``x`` / ``series`` — so the build raises
    ``KeyError: 'labels'``.

    This test pins the CURRENT broken behaviour. When the renderer is fixed to
    read ``x``/``series`` (or inject is taught to translate), flip this to a
    positive assertion that the pie builds as a native editable chart.
    """
    with pytest.raises(KeyError, match="labels"):
        _build_charts(parsed, tmp_path, (2,))

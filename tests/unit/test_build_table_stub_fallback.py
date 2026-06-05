"""P0-2 (2026-06-05): donor 53 (fixed_png_content table donor) keeps its
PNG-table-stub when no ``table_data`` is supplied.

Live run4.slide4 ("DNS Resolvers – access logs") was empty because:
1. donor 53 has ``remove_before_fill=[0]`` — the PNG stub is always stripped.
2. Pipeline didn't produce ``table_data`` for that slide (Agent 02/03 path).
3. Result: title + empty caption, no body, no table, nothing else — a
   visually blank slide on a layout that promised a table.

The fix in build_v9 suppresses ``remove_before_fill`` when the donor is
``fixed_png_content`` and neither ``table_data`` nor ``infographic`` is
present, so the stub stays as a visual placeholder. These tests pin
that behaviour and also assert the regular "stub-stripped-with-data"
path still works.
"""
from __future__ import annotations

import json

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402

from build_v9 import build  # noqa: E402


@pytest.fixture
def tmp_workdir(tmp_path):
    return tmp_path


def _build_donor53(tmp_workdir, *, include_table_data: bool):
    plan_slide = {
        "clone_from_slide": 53,
        "slots": {
            "title": "DNS Resolvers – access logs",
            "caption": "",
        },
    }
    if include_table_data:
        plan_slide["table_data"] = [
            ["Resolver", "Latency", "QPS"],
            ["8.8.8.8", "9 ms", "1.2M"],
            ["1.1.1.1", "11 ms", "0.9M"],
        ]
    plan = {"slides": [plan_slide]}
    plan_path = tmp_workdir / "plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False), encoding="utf-8")
    out = tmp_workdir / "out.pptx"
    build(
        str(plan_path),
        str(skill_bridge.TEMPLATE_PATH),
        str(out),
        str(skill_bridge.DONOR_SLOT_MAP),
    )
    return out


def _count_tables(slide) -> int:
    return sum(1 for sh in slide.shapes if getattr(sh, "has_table", False))


def _count_content_shapes(slide) -> int:
    """Number of body-class shapes (excludes placeholder TITLE).

    Used as a proxy for "the slide isn't blank" — placeholder title alone
    doesn't count as content here.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    return sum(
        1 for sh in slide.shapes
        if sh.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER
    )


def test_donor53_without_table_data_keeps_stub(tmp_workdir) -> None:
    """No table_data → donor's table-stub must remain so the slide isn't
    a blank "title only" frame. (Live run4.slide4 regression.)"""
    out = _build_donor53(tmp_workdir, include_table_data=False)
    prs = Presentation(out)
    slide = prs.slides[0]
    # The donor's table-stub is the visual placeholder we want to preserve.
    assert _count_tables(slide) >= 1, (
        "donor table-stub should be preserved when no replacement supplied"
    )
    # Sanity: more than just placeholder title remains.
    assert _count_content_shapes(slide) >= 1


def test_donor53_with_table_data_strips_stub_and_adds_table(tmp_workdir) -> None:
    """Sanity: regular path — stub IS stripped when table_data is provided.

    The donor's stub IS itself a table, so the count stays >=1 but the
    cell contents shift from the donor's example data to ours.
    """
    out = _build_donor53(tmp_workdir, include_table_data=True)
    prs = Presentation(out)
    slide = prs.slides[0]
    tables = [sh for sh in slide.shapes if getattr(sh, "has_table", False)]
    assert tables, "expected at least one table"
    # At least one table cell contains a value we supplied — the donor
    # stub didn't include "8.8.8.8".
    found_user_value = False
    for tbl_shape in tables:
        for row in tbl_shape.table.rows:
            for cell in row.cells:
                if "8.8.8.8" in (cell.text_frame.text or ""):
                    found_user_value = True
                    break
    assert found_user_value, (
        "expected user-supplied table_data to be present after stub removal"
    )

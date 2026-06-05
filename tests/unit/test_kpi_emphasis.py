"""T2.2: KPI emphasis post-pass — numbers in body text become bold+green.

Live run 2026-06-05 had 3 blockers from visual_verifier complaining that
key metrics weren't typographically highlighted (275, 568 125 090,
1,2 млн, 14,2 млн, 25 млн, 101 млн/мес, 100 млн разово). These tests pin
the regex matcher and the run-split mechanics.
"""
from __future__ import annotations

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from pptx.util import Pt  # noqa: E402

from kpi_emphasis import (  # noqa: E402  (path injected by skill_bridge)
    _NUMBER_RE,
    _qualifies,
    apply_kpi_emphasis,
    emphasize_kpi_in_slide,
)


@pytest.fixture
def blank_slide():
    prs = Presentation(str(skill_bridge.TEMPLATE_PATH))
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs, prs.slides.add_slide(layout)


# ─── Regex / qualifies ────────────────────────────────────────────────────────

@pytest.mark.parametrize("text, expected_tokens", [
    ("Выручка 1,2 млн в Q1", ["1,2 млн"]),
    ("275 контактов с добавлением", ["275"]),
    ("Сумма 568 125 090 руб", ["568 125 090 руб"]),
    ("101 млн/мес, 100 млн разово", ["101 млн", "100 млн"]),
    ("14,2 млн и 25 млн рублей", ["14,2 млн", "25 млн"]),
    # Edge: solo small number → not qualified
    ("в 5 раз быстрее", ["5 раз"]),  # has unit
    ("на странице 7", []),            # solo "7" — too short, no unit
    ("два или три проекта", []),     # no digits
])
def test_number_regex_matches(text: str, expected_tokens: list[str]) -> None:
    found = [m.group(0).rstrip() for m in _NUMBER_RE.finditer(text)
             if _qualifies(m.group("num"), m.group("unit"))]
    assert found == expected_tokens


def test_qualifies_threshold() -> None:
    assert _qualifies("275", None) is True       # 3 digits
    assert _qualifies("99", None) is False       # 2 digits, no unit
    assert _qualifies("99", "млн") is True       # has unit
    assert _qualifies("1,2", "млн") is True      # decimal + unit


# ─── Slide-level: run splitting + emphasis ────────────────────────────────────

def _add_textbox(slide, text: str, size_pt: int = 14):
    box = slide.shapes.add_textbox(Pt(50), Pt(50), Pt(400), Pt(80))
    tf = box.text_frame
    tf.text = text
    # Force a font size so the title-skip heuristic doesn't trigger.
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size_pt)
    return box


def _count_emphasized_runs(slide) -> int:
    """Count <a:r> nodes whose rPr has b='1' AND solidFill srgbClr='26D07C'."""
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txBody = shape.text_frame._txBody
        for p in txBody.findall(qn("a:p")):
            for r in p.findall(qn("a:r")):
                rPr = r.find(qn("a:rPr"))
                if rPr is None:
                    continue
                if rPr.get("b") != "1":
                    continue
                sf = rPr.find(qn("a:solidFill"))
                if sf is None:
                    continue
                srgb = sf.find(qn("a:srgbClr"))
                if srgb is not None and srgb.get("val", "").upper() == "26D07C":
                    n += 1
    return n


def test_emphasize_single_token(blank_slide) -> None:
    prs, slide = blank_slide
    _add_textbox(slide, "Выручка 1,2 млн в Q1", size_pt=14)
    n = emphasize_kpi_in_slide(slide)
    assert n == 1
    assert _count_emphasized_runs(slide) == 1


def test_emphasize_multiple_tokens(blank_slide) -> None:
    prs, slide = blank_slide
    _add_textbox(slide, "101 млн/мес, 100 млн разово", size_pt=14)
    n = emphasize_kpi_in_slide(slide)
    assert n == 2
    assert _count_emphasized_runs(slide) == 2


def test_emphasize_skips_title_size(blank_slide) -> None:
    """Runs with size_pt >= 28 are treated as titles and left alone."""
    prs, slide = blank_slide
    _add_textbox(slide, "275 контактов", size_pt=32)
    n = emphasize_kpi_in_slide(slide)
    assert n == 0


def test_emphasize_skips_short_number(blank_slide) -> None:
    """Bare 2-digit number with no unit is not emphasized."""
    prs, slide = blank_slide
    _add_textbox(slide, "часть 12 из проекта", size_pt=14)
    n = emphasize_kpi_in_slide(slide)
    assert n == 0


def test_apply_kpi_emphasis_skips_kpi_native(blank_slide) -> None:
    """Slides marked as kpi_native are skipped — render_kpi handles them."""
    prs, slide = blank_slide
    _add_textbox(slide, "275 контактов", size_pt=14)
    stats = apply_kpi_emphasis(
        prs,
        plan_slides=[{"slide_type": "kpi_native"}] * len(prs.slides),
    )
    assert stats["total"] == 0


def test_apply_kpi_emphasis_processes_when_no_plan(blank_slide) -> None:
    prs, slide = blank_slide
    _add_textbox(slide, "275 контактов", size_pt=14)
    stats = apply_kpi_emphasis(prs)
    assert stats["total"] >= 1


def _count_emphasized_runs_with_color(slide, hex_upper: str) -> int:
    """Like _count_emphasized_runs but parameterised on the expected colour."""
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txBody = shape.text_frame._txBody
        for p in txBody.findall(qn("a:p")):
            for r in p.findall(qn("a:r")):
                rPr = r.find(qn("a:rPr"))
                if rPr is None or rPr.get("b") != "1":
                    continue
                sf = rPr.find(qn("a:solidFill"))
                if sf is None:
                    continue
                srgb = sf.find(qn("a:srgbClr"))
                if srgb is not None and srgb.get("val", "").upper() == hex_upper:
                    n += 1
    return n


def test_emphasize_falls_back_to_graphite_on_green_box(blank_slide) -> None:
    """D2 fix: green KPI on a green-filled rect is invisible. The pass must
    switch to graphite (#222222) when the parent shape fills brand green —
    live run1.slide8 had `12.18` disappear inside an accent box."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs, slide = blank_slide
    # Add a green-filled rounded rect with KPI-shaped text in it.
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Emu(100000), Emu(100000),
                                   Emu(2000000), Emu(800000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x26, 0xD0, 0x7C)  # #26D07C
    tf = shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "v1.12.18 (май)"
    run.font.size = Pt(14)

    n = emphasize_kpi_in_slide(slide)
    assert n >= 1, "KPI token should still be emphasized"
    # Should be GRAPHITE, not GREEN — otherwise it disappears.
    assert _count_emphasized_runs_with_color(slide, "222222") >= 1
    assert _count_emphasized_runs_with_color(slide, "26D07C") == 0


def test_emphasize_falls_back_to_graphite_when_overlapping_green(blank_slide) -> None:
    """P0-1 (2026-06-05): Agent 06 native infographics layer a *text shape*
    (fill=none) ON TOP of a separate green-filled rounded_rect. The text
    shape itself isn't green so the legacy D2 check missed it — emphasis
    painted digits green-on-green, hiding them (live run4.slide8 lost
    "12.17" inside the middle accent block). The overlap detector must
    catch this and switch to graphite."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs, slide = blank_slide
    # 1. Green-filled backing rect (Agent 06 rounded_rect).
    green = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Emu(4286250), Emu(3143250),
                                   Emu(3190350), Emu(1143000))
    green.fill.solid()
    green.fill.fore_color.rgb = RGBColor(0x26, 0xD0, 0x7C)
    # No text on the backing rect (matches Agent 06 plan).

    # 2. Text shape sitting on top — its own fill is "none".
    text_box = slide.shapes.add_textbox(
        Emu(4381500), Emu(3238500), Emu(2990850), Emu(476250)
    )
    text_box.fill.background()  # no fill — relies on the green rect behind.
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "v1.12.17 (апрель)"
    run.font.size = Pt(14)

    n = emphasize_kpi_in_slide(slide)
    assert n >= 1, "KPI token in overlay text must still be emphasized"
    # Must be GRAPHITE because the underlying green rect would make
    # green-on-green invisible.
    assert _count_emphasized_runs_with_color(slide, "222222") >= 1, (
        "overlap-with-green detection should have picked graphite"
    )
    assert _count_emphasized_runs_with_color(slide, "26D07C") == 0, (
        "no run should still be coloured green"
    )


def test_emphasize_uses_green_when_no_overlap(blank_slide) -> None:
    """Sanity: a text shape that doesn't sit on a green rect still gets
    the regular green emphasis (regression guard for the P0-1 detector
    erroneously flagging unrelated shapes)."""
    prs, slide = blank_slide
    _add_textbox(slide, "Выручка 1,2 млн в Q1", size_pt=14)
    n = emphasize_kpi_in_slide(slide)
    assert n >= 1
    assert _count_emphasized_runs_with_color(slide, "26D07C") >= 1


def test_emphasize_preserves_surrounding_text(blank_slide) -> None:
    """Run-splitting keeps the non-number text in separate plain runs."""
    prs, slide = blank_slide
    _add_textbox(slide, "Выручка 1,2 млн в Q1 2026", size_pt=14)
    emphasize_kpi_in_slide(slide)
    # Concatenated text of all runs must still be the original sentence
    # (modulo regex whitespace handling).
    box = slide.shapes[-1]
    full = box.text_frame.text
    assert "Выручка" in full
    assert "1,2 млн" in full
    assert "Q1 2026" in full

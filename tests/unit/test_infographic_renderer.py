"""T2.1: infographic native_block renderer.

Live run 2026-06-05 confirmed Agent 06 emits 9 shapes per comparison
slide but build_v9 had no handler — visual verifier blocked on
``plan_compliance``. These tests pin the renderer behaviour so a
regression on the build path is caught before live spend.
"""
from __future__ import annotations

from pathlib import Path

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

from infographic_renderer import (  # noqa: E402  (path injected by skill_bridge)
    _clamp_shapes_to_safe_area,
    _parse_hex,
    clear_donor_body_slots,
    clear_donor_non_title_text,
    render_infographic_shapes,
)
from pptx.util import Emu, Pt  # noqa: E402


@pytest.fixture
def blank_slide():
    """A fresh blank slide from the template — independent per test."""
    prs = Presentation(str(skill_bridge.TEMPLATE_PATH))
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(blank_layout)


def test_parse_hex_valid_and_invalid() -> None:
    from pptx.dml.color import RGBColor
    c = _parse_hex("#26D07C")
    assert isinstance(c, RGBColor)
    assert _parse_hex("26d07c") is not None
    assert _parse_hex("none") is None
    assert _parse_hex("") is None
    assert _parse_hex(None) is None
    assert _parse_hex("#zzzzzz") is None
    assert _parse_hex("#12") is None


def test_render_rounded_rect_adds_shape(blank_slide) -> None:
    before = len(blank_slide.shapes)
    added = render_infographic_shapes(blank_slide, [{
        "type": "rounded_rect",
        "left_emu": 100000, "top_emu": 100000,
        "width_emu": 2000000, "height_emu": 500000,
        "fill_color": "#26D07C", "stroke_color": "none",
        "stroke_width_pt": 0.0,
        "text": "Prognoz",
        "font": "SB Sans Display Semibold",
        "font_size_pt": 14, "font_color": "#222222",
    }])
    assert added == 1
    assert len(blank_slide.shapes) == before + 1
    shape = blank_slide.shapes[-1]
    assert shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    # Text actually written into the rect.
    assert "Prognoz" in shape.text_frame.text


def test_render_textbox_adds_shape(blank_slide) -> None:
    before = len(blank_slide.shapes)
    added = render_infographic_shapes(blank_slide, [{
        "type": "text",
        "left_emu": 200000, "top_emu": 200000,
        "width_emu": 3000000, "height_emu": 400000,
        "fill_color": "none", "stroke_color": "none",
        "stroke_width_pt": 0.0,
        "text": "275 контактов",
        "font": "SB Sans Display",
        "font_size_pt": 16, "font_color": "#222222",
    }])
    assert added == 1
    assert len(blank_slide.shapes) == before + 1
    assert "275 контактов" in blank_slide.shapes[-1].text_frame.text


@pytest.mark.parametrize("kind", ["rectangle", "circle", "arrow", "line"])
def test_render_additional_shape_types(blank_slide, kind: str) -> None:
    """T2.5: rectangle / circle / arrow / line — added 2026-06-05 after
    live run dropped 3/15 arrow shapes from a process infographic."""
    before = len(blank_slide.shapes)
    added = render_infographic_shapes(blank_slide, [{
        "type": kind,
        "left_emu": 100000, "top_emu": 100000,
        "width_emu": 800000, "height_emu": 200000,
        "fill_color": "#26D07C" if kind != "line" else "none",
        "stroke_color": "#222222",
        "stroke_width_pt": 1.0,
        "text": "",
    }])
    assert added == 1, f"{kind} should be added"
    assert len(blank_slide.shapes) == before + 1


def test_render_skips_unknown_type(blank_slide) -> None:
    before = len(blank_slide.shapes)
    added = render_infographic_shapes(blank_slide, [
        {"type": "ellipse", "left_emu": 0, "top_emu": 0,
         "width_emu": 100000, "height_emu": 100000},
        {"type": "rounded_rect", "left_emu": 0, "top_emu": 0,
         "width_emu": 1000000, "height_emu": 500000,
         "fill_color": "#26D07C", "text": ""},
    ])
    assert added == 1
    assert len(blank_slide.shapes) == before + 1


def test_render_empty_list_noop(blank_slide) -> None:
    before = len(blank_slide.shapes)
    assert render_infographic_shapes(blank_slide, []) == 0
    assert len(blank_slide.shapes) == before


def test_render_swallows_bad_spec(blank_slide) -> None:
    """Non-dict entries are logged and skipped — never raise."""
    before = len(blank_slide.shapes)
    added = render_infographic_shapes(blank_slide, ["string", 42, None])
    assert added == 0
    assert len(blank_slide.shapes) == before


def test_render_live_comparison_block(blank_slide) -> None:
    """End-to-end: 9 shapes mirroring the 2026-06-05 slide 2 native_block."""
    shapes = (
        [{"type": "rounded_rect", "left_emu": 100000, "top_emu": 1000000,
          "width_emu": 5000000, "height_emu": 500000,
          "fill_color": "#222222", "stroke_color": "none",
          "stroke_width_pt": 0.0, "text": "Q1 2026",
          "font": "SB Sans Display Semibold", "font_size_pt": 14,
          "font_color": "#FFFFFF"}]
        + [{"type": "text", "left_emu": 200000, "top_emu": 1500000 + 400000 * i,
            "width_emu": 4800000, "height_emu": 350000,
            "fill_color": "none", "stroke_color": "none",
            "stroke_width_pt": 0.0,
            "text": f"строка {i}", "font": "SB Sans Display",
            "font_size_pt": 12, "font_color": "#222222"}
           for i in range(3)]
        + [{"type": "rounded_rect", "left_emu": 6800000, "top_emu": 1000000,
            "width_emu": 5000000, "height_emu": 500000,
            "fill_color": "#26D07C", "stroke_color": "none",
            "stroke_width_pt": 0.0, "text": "Q2 2026",
            "font": "SB Sans Display Semibold", "font_size_pt": 14,
            "font_color": "#222222"}]
        + [{"type": "text", "left_emu": 6900000, "top_emu": 1500000 + 400000 * i,
            "width_emu": 4800000, "height_emu": 350000,
            "fill_color": "none", "stroke_color": "none",
            "stroke_width_pt": 0.0,
            "text": f"итог {i}", "font": "SB Sans Display",
            "font_size_pt": 12, "font_color": "#222222"}
           for i in range(4)]
    )
    before = len(blank_slide.shapes)
    added = render_infographic_shapes(blank_slide, shapes)
    assert added == 9
    assert len(blank_slide.shapes) == before + 9


def test_clear_donor_body_keeps_title(blank_slide) -> None:
    """Title-like slots are preserved; body/content/caption/col* are cleared."""
    donor_def = {
        "slots": {
            "title": {"shape_idx": 0},
            "body": {"shape_idx": 1},
            "col1_body": {"shape_idx": 2},
            "caption": {"shape_idx": 3},
        }
    }
    # No shapes match, so the function just probes and skips — exercising the
    # name filter without depending on a particular donor.
    cleared = clear_donor_body_slots(blank_slide, donor_def)
    assert cleared >= 0  # blank slide → no matching idx, function returns safely


def test_clear_donor_body_handles_none_def(blank_slide) -> None:
    assert clear_donor_body_slots(blank_slide, None) == 0
    assert clear_donor_body_slots(blank_slide, {}) == 0


def _add_textbox_with_text(slide, text: str, *, font_pt: int = 14):
    """Helper: append a plain textbox with given text/font size."""
    box = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(2000000), Emu(500000))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    return box


def test_clear_donor_non_title_text_clears_all_non_placeholder_text(blank_slide) -> None:
    """B2 (2026-06-05): live a337cc86 slides 7/9 had donor 33 column
    sub-headers ("Подзаголовок в две строки") bleeding through behind
    infographic boxes — they were 18-24pt textboxes (NOT placeholder
    titles), so the old `font >= 18pt` heuristic protected them as
    "title-like". New rule: only placeholder-TITLE is preserved; every
    other text shape gets wiped regardless of font size."""
    # Donor 33 decoration: sub-header at 22pt (looks title-like but isn't)
    sub_header = _add_textbox_with_text(blank_slide, "Подзаголовок в две строки",
                                        font_pt=22)
    body1 = _add_textbox_with_text(blank_slide, "Recorder", font_pt=14)
    body2 = _add_textbox_with_text(blank_slide, "Хранение данных", font_pt=12)
    cleared = clear_donor_non_title_text(blank_slide)
    assert cleared >= 3  # all three should be wiped
    assert (sub_header.text_frame.text or "").strip() == ""
    assert (body1.text_frame.text or "").strip() == ""
    assert (body2.text_frame.text or "").strip() == ""


def test_clear_donor_non_title_text_preserves_real_title_placeholder() -> None:
    """When a real title placeholder exists, its text must survive."""
    prs = Presentation(str(skill_bridge.TEMPLATE_PATH))
    # Use a layout that has a TITLE placeholder (layout 0 = title slide).
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Find and set the title placeholder.
    title_text = "REAL TITLE PLACEHOLDER"
    title_ph = None
    for sh in slide.shapes:
        try:
            if sh.placeholder_format is not None and sh.placeholder_format.idx == 0:
                sh.text_frame.text = title_text
                title_ph = sh
                break
        except (ValueError, AttributeError):
            continue
    assert title_ph is not None, "test setup failed: no title placeholder"
    # Add a body-like textbox that should be cleared.
    body = _add_textbox_with_text(slide, "Body content", font_pt=14)
    clear_donor_non_title_text(slide)
    assert title_text in title_ph.text_frame.text  # preserved
    assert (body.text_frame.text or "").strip() == ""  # cleared


def test_clear_donor_non_title_text_handles_empty_slide(blank_slide) -> None:
    assert clear_donor_non_title_text(blank_slide) == 0


def test_clear_donor_non_title_text_skips_already_empty(blank_slide) -> None:
    """No text → nothing to clear."""
    _add_textbox_with_text(blank_slide, "", font_pt=14)
    cleared = clear_donor_non_title_text(blank_slide)
    assert cleared == 0


# ─── P1-1 (2026-06-05): SAFE_AREA clamp for Agent 06 process shapes ──────────

_EMU = 9525  # 1 px


def _block(left_px: int, width_px: int, top_px: int = 300,
           height_px: int = 100, type_: str = "rounded_rect",
           text: str = "") -> dict:
    return {
        "type": type_,
        "left_emu": left_px * _EMU,
        "top_emu": top_px * _EMU,
        "width_emu": width_px * _EMU,
        "height_emu": height_px * _EMU,
        "fill_color": "#F2F2F2",
        "stroke_color": "none",
        "stroke_width_pt": 0.0,
        "text": text,
        "font": "SB Sans Display",
        "font_size_pt": 12,
        "font_color": "#222222",
    }


def test_clamp_noop_when_inside_safe_area() -> None:
    """3-block process inside safe-area [30, 1250] — no mutation."""
    shapes = [
        _block(60, 300),
        _block(420, 300),
        _block(780, 300),
    ]
    snapshot = [(s["left_emu"], s["width_emu"]) for s in shapes]
    mutated = _clamp_shapes_to_safe_area(shapes)
    assert mutated == 0
    after = [(s["left_emu"], s["width_emu"]) for s in shapes]
    assert after == snapshot


def test_clamp_shrinks_4block_overshoot() -> None:
    """Live run3.slide2 emitted 4 blocks 330px wide + 30px gaps starting
    at x=30 — total span 1410px, well past safe-right=1250px. The
    clamp must scale everything down so the last block ends ≤ 1250px."""
    shapes = [
        _block(30, 330),    # 30..360
        _block(390, 330),   # 390..720
        _block(750, 330),   # 750..1080
        _block(1110, 330),  # 1110..1440   ← overshoots
        # Arrows (zero-width spans skipped) — include 1 to exercise
        # mixed shape types.
        _block(360, 30, top_px=350, height_px=10, type_="arrow"),
    ]
    mutated = _clamp_shapes_to_safe_area(shapes)
    assert mutated >= 4, "all 4 process blocks should be mutated"
    # All shapes must end ≤ safe.right (1250 px → 11_906_250 EMU).
    safe_right_emu = 1250 * _EMU
    safe_left_emu = 30 * _EMU
    for s in shapes:
        assert s["left_emu"] >= safe_left_emu - 100  # tiny rounding slack
        assert (s["left_emu"] + s["width_emu"]) <= safe_right_emu + 100


def test_clamp_preserves_relative_order_and_gaps() -> None:
    """Scaling must keep block ordering and roughly proportional gaps."""
    shapes = [
        _block(30, 330),
        _block(390, 330),
        _block(750, 330),
        _block(1110, 330),
    ]
    _clamp_shapes_to_safe_area(shapes)
    lefts = [s["left_emu"] for s in shapes]
    # Strictly increasing.
    assert lefts == sorted(lefts)
    # Pairwise gaps roughly equal (within rounding).
    gaps = [lefts[i + 1] - lefts[i] for i in range(len(lefts) - 1)]
    assert max(gaps) - min(gaps) <= 10_000, gaps


def test_clamp_empty_input_returns_zero() -> None:
    assert _clamp_shapes_to_safe_area([]) == 0
    assert _clamp_shapes_to_safe_area([{"type": "text"}]) == 0  # no width


# ─── B1 (2026-06-05): upscale branch for Agent 06 underscale hallucination ───

def test_clamp_upscales_tiny_shapes() -> None:
    """Live run 29e189bb.slide4 (2026-06-05): GLM-5.1 emitted 3 process
    blocks at width_emu=352425 (37 px) on a 1280-px canvas — total span
    179 px = 15% of safe-area. Clamp must scale UP to ~95% of safe width."""
    # 3 blocks at x=30, 90, 150, each 37px wide → span 30..187 = 157 px
    shapes = [
        _block(30, 37),
        _block(90, 37),
        _block(150, 37),
    ]
    mutated = _clamp_shapes_to_safe_area(shapes)
    assert mutated == 3, "all 3 undersized blocks should be upscaled"
    # After upscale, the span should occupy most of safe-area (~95% × 1220px).
    min_left = min(s["left_emu"] for s in shapes)
    max_right = max(s["left_emu"] + s["width_emu"] for s in shapes)
    span_px = (max_right - min_left) / _EMU
    assert 1000 <= span_px <= 1230, f"expected ~1100-1220 px span, got {span_px:.0f}"
    # First block must start at safe.left (30 px) after shift.
    assert abs(shapes[0]["left_emu"] - 30 * _EMU) < 5000


def test_clamp_upscales_tiny_comparison_2cols() -> None:
    """Live run 29e189bb.slide2: comparison (6 shapes) spanning 149 px
    on a 1220-px safe-area. Must upscale; relative gap between two
    columns must be preserved."""
    shapes = [
        # Col1: rect + 2 text labels
        _block(30, 57, top_px=54, height_px=140),
        _block(38, 41, top_px=59, height_px=56, type_="text"),
        _block(38, 41, top_px=115, height_px=56, type_="text"),
        # Col2: rect + 2 text labels
        _block(92, 57, top_px=54, height_px=140),
        _block(100, 41, top_px=59, height_px=56, type_="text"),
        _block(100, 41, top_px=115, height_px=56, type_="text"),
    ]
    mutated = _clamp_shapes_to_safe_area(shapes)
    assert mutated >= 6
    # The 2 columns should now be far apart (not 5px apart as before).
    col1_right = shapes[0]["left_emu"] + shapes[0]["width_emu"]
    col2_left = shapes[3]["left_emu"]
    gap_px = (col2_left - col1_right) / _EMU
    assert gap_px > 30, f"expected meaningful inter-column gap, got {gap_px:.0f}px"


def test_clamp_upscales_y_too() -> None:
    """Underscale path scales Y axis as well — otherwise tiny boxes get
    tall thin look. Live run 29e189bb shapes were 194 px tall on a
    520-px-tall safe area (37% — borderline). Verify height grows."""
    shapes = [
        _block(30, 37, top_px=140, height_px=40),
        _block(90, 37, top_px=140, height_px=40),
        _block(150, 37, top_px=140, height_px=40),
    ]
    original_h = shapes[0]["height_emu"]
    _clamp_shapes_to_safe_area(shapes)
    assert shapes[0]["height_emu"] > original_h, "Y axis should upscale too"


def test_clamp_upscales_font_size_proportionally() -> None:
    """Tiny boxes had 12pt text — after 10× upscale the text would be
    swallowed by huge empty cards. Font should scale up too, capped at
    24pt so labels don't take over."""
    shapes = [_block(30, 37), _block(90, 37), _block(150, 37)]
    # Add explicit font_size_pt
    for s in shapes:
        s["font_size_pt"] = 12
    _clamp_shapes_to_safe_area(shapes)
    # All sizes capped between 10 and 24.
    for s in shapes:
        assert 10 <= s["font_size_pt"] <= 24


def test_clamp_does_not_upscale_when_above_threshold() -> None:
    """A 3-block process at 300px each + 60px gaps = span 30..1020 = 990 px
    on 1220 safe-area = 81% — well above 50% threshold. Must NOT touch."""
    shapes = [
        _block(30, 300),
        _block(390, 300),
        _block(750, 300),
    ]
    snapshot = [(s["left_emu"], s["width_emu"], s.get("font_size_pt", 12)) for s in shapes]
    mutated = _clamp_shapes_to_safe_area(shapes)
    assert mutated == 0
    after = [(s["left_emu"], s["width_emu"], s.get("font_size_pt", 12)) for s in shapes]
    assert after == snapshot

"""C (2026-06-07): bullet + spacing for multi-item body lists.

Taxonomy defect C: multi-point body content renders as plain paragraphs with no
bullet and no spacing (donor body placeholders carry no pPr), so the items merge
into one solid block (live: 5d s4, 9b s5). Fix: a body slot with >=2 paragraphs
becomes a real bulleted list (buChar + hanging indent + spcBef отбивка); a single
paragraph stays untouched.
"""
from __future__ import annotations

import json

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

from build_v9 import build  # noqa: E402


@pytest.fixture
def tmp_workdir(tmp_path):
    return tmp_path


def _body_paragraphs(out_pptx):
    """Donor 21 body lives on shape_idx=1 (per donor-slot-map.yaml)."""
    prs = Presentation(out_pptx)
    shapes = list(prs.slides[0].shapes)
    tb = shapes[1].text_frame._txBody
    return tb.findall(qn("a:p"))


def _build_body(tmp_workdir, body: str):
    plan = {"slides": [{"clone_from_slide": 21,
                        "slots": {"title": "Преимущества", "body": body}}]}
    plan_path = tmp_workdir / "plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False), encoding="utf-8")
    out = tmp_workdir / "out.pptx"
    build(str(plan_path), str(skill_bridge.TEMPLATE_PATH), str(out),
          str(skill_bridge.DONOR_SLOT_MAP))
    return out


def test_multi_item_body_gets_bullets(tmp_workdir) -> None:
    body = "Масштабируемость до тысяч ядер.\nГибкая тарификация.\nПоддержка 24/7."
    paras = _body_paragraphs(_build_body(tmp_workdir, body))
    bulleted = [p for p in paras if p.find(qn("a:pPr")) is not None
                and p.find(qn("a:pPr")).find(qn("a:buChar")) is not None]
    assert len(bulleted) == 3, "each list item should carry a bullet glyph"


def test_multi_item_body_has_spacing_between_items(tmp_workdir) -> None:
    body = "Первый пункт списка.\nВторой пункт списка.\nТретий пункт списка."
    paras = _body_paragraphs(_build_body(tmp_workdir, body))
    # First item: no spcBef (no gap above the head). Items 2+: spcBef present.
    spaced = [p for p in paras
              if p.find(qn("a:pPr")) is not None
              and p.find(qn("a:pPr")).find(qn("a:spcBef")) is not None]
    assert len(spaced) == 2, "items after the first should have отбивка"


def test_single_paragraph_body_not_bulleted(tmp_workdir) -> None:
    paras = _body_paragraphs(_build_body(tmp_workdir, "Один короткий вывод."))
    has_bullet = any(
        p.find(qn("a:pPr")) is not None
        and p.find(qn("a:pPr")).find(qn("a:buChar")) is not None
        for p in paras
    )
    assert not has_bullet, "a single conclusion is a paragraph, not a list"


def test_pPr_child_order_is_schema_valid(tmp_workdir) -> None:
    """spcBef must precede buFont/buChar in CT_TextParagraphProperties."""
    body = "Альфа пункт.\nБета пункт.\nГамма пункт."
    paras = _body_paragraphs(_build_body(tmp_workdir, body))
    order = {"spcBef": 0, "buFont": 1, "buChar": 2}
    for p in paras:
        pPr = p.find(qn("a:pPr"))
        if pPr is None:
            continue
        seen = [order[c.tag.split("}")[-1]] for c in pPr
                if c.tag.split("}")[-1] in order]
        assert seen == sorted(seen), f"pPr children out of schema order: {seen}"

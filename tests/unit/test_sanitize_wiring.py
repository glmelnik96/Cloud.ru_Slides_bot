"""Wiring tests for sanitize_text into the PPTX text chokepoints.

build_v5.replace_text_with_style is the main donor-text chokepoint (body AND
the donor TITLE route in build_v9). Its output later feeds kpi_emphasis, so it
strips control chars but KEEPS ** (kpi-safe). kpi_renderer.set_slide_title is a
generic title setter (chart/table/flow) and never carries intentional ** — it
strips both.
"""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

from build_v5 import replace_text_with_style  # noqa: E402
from kpi_renderer import set_slide_title  # noqa: E402


def _make_text_frame():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(914400), Emu(914400))
    tf = box.text_frame
    # Seed a run so replace_text_with_style has an rPr to clone.
    tf.paragraphs[0].add_run().text = "seed"
    return prs, slide, tf


def _all_text(tf):
    return "".join(t.text or "" for t in tf._txBody.iter(qn("a:t")))


def test_replace_text_strips_vertical_tab_in_title() -> None:
    """Artifact A: donor title with \x0b must not yield _X000B_ / raw \x0b."""
    _prs, _slide, tf = _make_text_frame()
    replace_text_with_style(tf, "ТЕХНИЧЕСКИЙ\x0bРАЗДЕЛ")
    text = _all_text(tf)
    assert "\x0b" not in text
    # build_v5 converts \x0b -> \n before sanitize, so it becomes a line break.
    lines = [p for p in tf._txBody.findall(qn("a:p"))]
    assert len(lines) == 2


def test_replace_text_strips_stray_control_char() -> None:
    _prs, _slide, tf = _make_text_frame()
    replace_text_with_style(tf, "clean\x07text")
    assert "\x07" not in _all_text(tf)


def test_replace_text_keeps_double_asterisk_for_kpi_emphasis() -> None:
    """build_v5 feeds kpi_emphasis later — ** MUST survive this chokepoint."""
    _prs, _slide, tf = _make_text_frame()
    replace_text_with_style(tf, "see **the key phrase** here")
    assert "**the key phrase**" in _all_text(tf)


def test_set_slide_title_strips_control_and_markdown() -> None:
    """Generic title setter: strip both control chars and ** (no intentional **)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_title(slide, "**Заголовок**\x0bтест")
    # find the title text
    title_tf = slide.shapes.title.text_frame
    text = "".join(t.text or "" for t in title_tf._txBody.iter(qn("a:t")))
    assert "**" not in text
    assert "\x0b" not in text
    assert "ЗАГОЛОВОК" in text  # uppercased by set_slide_title

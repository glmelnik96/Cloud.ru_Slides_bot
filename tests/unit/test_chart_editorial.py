"""D1: editorial single-series bar chart — per-bar green ramp + выноска overlay."""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402

from pptx.enum.chart import XL_CHART_TYPE  # noqa: E402
from pptx.util import Emu  # noqa: E402

from chart_native_pptx import (  # noqa: E402
    _green_ramp, add_chart_to_slide, is_editorial_eligible,
    render_chart_pptx_slide,
)


def _blank_slide():
    prs = Presentation(str(skill_bridge.TEMPLATE_PATH))
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def test_green_ramp_dark_to_bright_with_accent():
    ramp = _green_ramp(4, accent_idx=3)
    assert len(ramp) == 4
    # Accent bar is the brightest brand green (#26D07C).
    assert str(ramp[3]).upper() == "26D07C"
    # Earlier bars are darker (lower luminance) than the accent.
    def lum(c):
        h = str(c)
        return int(h[0:2], 16) + int(h[2:4], 16) + int(h[4:6], 16)
    assert lum(ramp[0]) < lum(ramp[3])


def test_editorial_eligible_only_single_series_bar():
    assert is_editorial_eligible(
        {"type": "bar", "style": "editorial",
         "series": [{"name": "a", "data": [1, 2, 3]}]}) is True
    # Multi-series → not eligible.
    assert is_editorial_eligible(
        {"type": "bar", "style": "editorial",
         "series": [{"name": "a", "data": [1]}, {"name": "b", "data": [2]}]}) is False
    # Non-bar → not eligible.
    assert is_editorial_eligible(
        {"type": "line", "style": "editorial",
         "series": [{"name": "a", "data": [1, 2]}]}) is False
    # No editorial style → not eligible.
    assert is_editorial_eligible(
        {"type": "bar", "series": [{"name": "a", "data": [1, 2]}]}) is False


def test_editorial_render_splits_bars_and_overlays_number():
    prs = Presentation(str(skill_bridge.TEMPLATE_PATH))
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    render_chart_pptx_slide(slide, {
        "title": "Динамика числа клиентов",
        "type": "bar", "style": "editorial",
        "x": [2023, 2024, 2025, 2026],
        "series": [{"name": "Клиенты", "data": [120, 340, 720, 1280]}],
        "accent_idx": 3,
    }, dark=True)
    # One native chart object with N single-value series (one per bar).
    charts = [s for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
    assert len(list(charts[0].chart.series)) == 4
    # A large выноска number textbox carrying the peak value.
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    assert any("1280" in t for t in texts)


def test_pie_builds_from_canonical_x_series():
    """Pie path reads the canonical ChartConfig shape (``x`` + ``series``)."""
    slide = _blank_slide()
    box = (Emu(457200), Emu(914400), Emu(8229600), Emu(4572000))
    add_chart_to_slide(slide, {
        "type": "pie",
        "title": "Доля рынка",
        "x": ["IaaS", "PaaS", "SaaS", "Прочее"],
        "series": [{"name": "Доля", "data": [40.0, 25.0, 20.0, 15.0]}],
        "accent_idx": 0,
    }, *box)
    charts = [s for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
    chart = charts[0].chart
    assert chart.chart_type == XL_CHART_TYPE.PIE
    out_series = list(chart.series)
    assert len(out_series) == 1
    assert list(out_series[0].values) == [40.0, 25.0, 20.0, 15.0]
    assert list(chart.plots[0].categories) == ["IaaS", "PaaS", "SaaS", "Прочее"]
    # Title still applied to the native chart object.
    assert chart.has_title


def test_pie_builds_from_legacy_labels_values():
    """Backward compat: legacy ``labels`` + ``values`` keys still build a pie."""
    slide = _blank_slide()
    box = (Emu(457200), Emu(914400), Emu(8229600), Emu(4572000))
    add_chart_to_slide(slide, {
        "type": "pie",
        "labels": ["A", "B", "C"],
        "values": [10.0, 30.0, 60.0],
        "accent_idx": 2,
    }, *box)
    charts = [s for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
    chart = charts[0].chart
    assert chart.chart_type == XL_CHART_TYPE.PIE
    out_series = list(chart.series)
    assert len(out_series) == 1
    assert list(out_series[0].values) == [10.0, 30.0, 60.0]
    assert list(chart.plots[0].categories) == ["A", "B", "C"]

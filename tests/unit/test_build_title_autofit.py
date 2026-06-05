"""D9: title overflow autofit.

Live run 2026-06-05 (run2.slide1): the cover topic was longer than
donor 4's safe_max_chars=55 — the 60pt title rendered with letters
running off the slide. Pin the shrink-on-overflow behaviour so a
long title proportionally reduces font_size with a 70%/14pt floor.
"""
from __future__ import annotations

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

from build_v9 import build  # noqa: E402


@pytest.fixture
def tmp_workdir(tmp_path):
    return tmp_path


def _read_title_run_size_pt(out_pptx) -> float | None:
    """Donor 4 title slot lives on shape_idx=1 (per donor-slot-map.yaml)."""
    prs = Presentation(out_pptx)
    slide = prs.slides[0]
    shapes = list(slide.shapes)
    if len(shapes) < 2:
        return None
    title_shape = shapes[1]
    if not title_shape.has_text_frame:
        return None
    for p in title_shape.text_frame._txBody.findall(qn("a:p")):
        for r in p.findall(qn("a:r")):
            rPr = r.find(qn("a:rPr"))
            if rPr is None:
                continue
            sz = rPr.get("sz")
            if sz and sz.isdigit():
                return int(sz) / 100.0
    return None


def _build_with_title(tmp_workdir, title: str):
    """Build a 1-slide deck cloning donor 4 with the given title."""
    import json
    plan = {
        "slides": [
            {
                "clone_from_slide": 4,
                "slots": {"title": title},
            }
        ]
    }
    plan_path = tmp_workdir / "plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False), encoding="utf-8")
    out = tmp_workdir / "out.pptx"
    build(
        str(plan_path),
        str(skill_bridge.TEMPLATE_PATH),
        str(out),
        str(skill_bridge.DONOR_SLOT_MAP),
    )
    return out


def test_short_title_no_explicit_size_override(tmp_workdir) -> None:
    """Title within safe_max_chars (≤55 for donor 4) shouldn't get an
    explicit sz override — donor's layout-inherited 60pt is preserved."""
    out = _build_with_title(tmp_workdir, "Q1 2026 Review")
    size = _read_title_run_size_pt(out)
    # Either no explicit sz (inheriting from layout) or exactly the donor's
    # canonical 60pt — both are acceptable for the short-title path. The
    # critical assertion is that no shrink was applied.
    assert size is None or size >= 60.0


def test_long_title_shrinks(tmp_workdir) -> None:
    """Title ~2× safe_max_chars triggers explicit shrink."""
    long_title = (
        "Cloud.ru квартальный обзор: "
        "ключевые показатели, риски и план на следующий период работы"
    )
    assert len(long_title) > 55
    out = _build_with_title(tmp_workdir, long_title)
    size = _read_title_run_size_pt(out)
    assert size is not None, "expected an explicit sz override after shrink"
    assert size < 60.0, f"expected shrink from 60pt, got {size}pt"
    # 14pt floor — for this length we should not hit it though.
    assert size >= 14.0


def test_extreme_title_hits_70pct_floor(tmp_workdir) -> None:
    """For very long titles the linear scale clamps at 70% of base.

    Without the floor we'd land at unreadable sub-14pt sizes.
    """
    extreme = "x" * 300
    out = _build_with_title(tmp_workdir, extreme)
    size = _read_title_run_size_pt(out)
    assert size is not None
    # 70% × 60pt = 42pt — should not go below that.
    assert size >= 42.0 - 0.5

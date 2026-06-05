"""B5 (2026-06-05): skip Agent 06 overlay when donor is itself a
structural multicolumn layout.

Live run a337cc86.slide12 used donor 34 (content_3col_subtitle) with an
Agent 06 ``matrix`` block. Both rendered together: the donor's 3 native
columns AND Agent 06's 9-shape matrix overlay, producing massive visual
overlap. The fix in build_v9 detects this combination and drops the
overlay; donor-native slots are used as-is.

These tests pin the rule:

* donor 34 + ``matrix`` → overlay shapes NOT injected (donor native only)
* donor 21 + ``matrix`` → overlay shapes ARE injected (non-structural donor)
"""
from __future__ import annotations

import json

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402

from build_v9 import build  # noqa: E402


@pytest.fixture
def tmp_workdir(tmp_path):
    return tmp_path


# Sentinel text that only the Agent 06 overlay shapes carry — if we see
# this text inside the rendered slide, the overlay was injected.
_OVERLAY_MARKER_A = "OVERLAY_CELL_A_UNIQUE"
_OVERLAY_MARKER_B = "OVERLAY_CELL_B_UNIQUE"


def _matrix_shapes():
    """Return a tiny but valid Agent 06 shape list. Two text cells with
    sentinel markers + one rounded_rect background each — coordinates
    are inside SAFE_AREA so the clamp doesn't reshape them."""
    return [
        {
            "type": "rounded_rect",
            "left_emu": 400000, "top_emu": 600000,
            "width_emu": 1500000, "height_emu": 800000,
            "fill_color": "#F2F2F2", "stroke_color": "none",
            "stroke_width_pt": 0.0, "text": "",
            "font": "SB Sans Display", "font_size_pt": 12,
            "font_color": "#222222",
        },
        {
            "type": "text",
            "left_emu": 500000, "top_emu": 700000,
            "width_emu": 1300000, "height_emu": 600000,
            "fill_color": "none", "stroke_color": "none",
            "stroke_width_pt": 0.0, "text": _OVERLAY_MARKER_A,
            "font": "SB Sans Display", "font_size_pt": 12,
            "font_color": "#222222",
        },
        {
            "type": "text",
            "left_emu": 2200000, "top_emu": 700000,
            "width_emu": 1300000, "height_emu": 600000,
            "fill_color": "none", "stroke_color": "none",
            "stroke_width_pt": 0.0, "text": _OVERLAY_MARKER_B,
            "font": "SB Sans Display", "font_size_pt": 12,
            "font_color": "#222222",
        },
    ]


def _build_one(tmp_workdir, plan_slide):
    plan = {"slides": [plan_slide]}
    plan_path = tmp_workdir / "plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False), encoding="utf-8")
    out = tmp_workdir / "out.pptx"
    build(
        str(plan_path),
        str(skill_bridge.TEMPLATE_PATH),
        str(out),
        str(skill_bridge.DONOR_SLOT_MAP),
    )
    return out


def _slide_text(slide) -> str:
    parts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            parts.append(sh.text_frame.text or "")
    return "\n".join(parts)


def test_structural_donor_skips_overlay(tmp_workdir) -> None:
    """donor 34 (content_3col_subtitle) + matrix → overlay NOT injected."""
    plan_slide = {
        "clone_from_slide": 34,
        "slots": {
            "title": "Three columns",
            "sub1": "A", "sub2": "B", "sub3": "C",
            "body1": "alpha", "body2": "beta", "body3": "gamma",
        },
        "infographic": {
            "type": "matrix",
            "shapes": _matrix_shapes(),
        },
    }
    out = _build_one(tmp_workdir, plan_slide)
    prs = Presentation(out)
    txt = _slide_text(prs.slides[0])
    # Overlay markers MUST be absent — B5 should have skipped the overlay.
    assert _OVERLAY_MARKER_A not in txt, (
        "Agent 06 overlay shape leaked into structural donor 34"
    )
    assert _OVERLAY_MARKER_B not in txt, (
        "Agent 06 overlay shape leaked into structural donor 34"
    )
    # Donor's native slot content should still be present.
    assert "alpha" in txt and "beta" in txt and "gamma" in txt


def test_nonstructural_donor_keeps_overlay(tmp_workdir) -> None:
    """donor 21 (content_text) + matrix → overlay IS injected.

    This guards against over-skipping: if we accidentally extend B5 to
    every donor, plain content donors will silently lose their Agent 06
    infographics.
    """
    plan_slide = {
        "clone_from_slide": 21,
        "slots": {
            "title": "Plain content",
            "body": "irrelevant body text",
        },
        "infographic": {
            "type": "matrix",
            "shapes": _matrix_shapes(),
        },
    }
    out = _build_one(tmp_workdir, plan_slide)
    prs = Presentation(out)
    txt = _slide_text(prs.slides[0])
    # Overlay markers MUST be present — non-structural donor still
    # needs the Agent 06 overlay.
    assert _OVERLAY_MARKER_A in txt, (
        "Agent 06 overlay shape missing on non-structural donor 21"
    )
    assert _OVERLAY_MARKER_B in txt, (
        "Agent 06 overlay shape missing on non-structural donor 21"
    )


# ---------------------------------------------------------------------------
# F1 (2026-06-05): Case A vs Case B disambiguation.
#
# Live run7 (eb6c4ceec3024bd9) showed that when the distributor produced a
# structural donor with only the title filled, B5 dropped the overlay AND
# nothing cleared the donor mock decoration ("Подзаголовок в две строки").
# The fix: drop overlay only when ≥2 body slots are filled; otherwise keep
# overlay AND clear donor mock decoration.
# ---------------------------------------------------------------------------


def test_case_b_structural_donor_underfilled_keeps_overlay(tmp_workdir) -> None:
    """Case B: structural donor 34 with ONLY title filled + matrix overlay
    → overlay MUST be rendered (otherwise the slide would be visually empty).
    """
    plan_slide = {
        "clone_from_slide": 34,
        "slots": {
            "title": "Underfilled structural donor",
            # sub1/sub2/sub3, body1/body2/body3 deliberately omitted —
            # this is what live run7 distributor produced.
        },
        "infographic": {
            "type": "matrix",
            "shapes": _matrix_shapes(),
        },
    }
    out = _build_one(tmp_workdir, plan_slide)
    prs = Presentation(out)
    txt = _slide_text(prs.slides[0])
    # Overlay markers MUST be present — Case B keeps overlay so the slide
    # has actual content.
    assert _OVERLAY_MARKER_A in txt, (
        "Case B: overlay must be kept when structural donor is underfilled"
    )
    assert _OVERLAY_MARKER_B in txt, (
        "Case B: overlay must be kept when structural donor is underfilled"
    )


def test_case_a_one_filled_body_slot_keeps_overlay(tmp_workdir) -> None:
    """Boundary: structural donor 34 with ONLY 1 body slot filled (below
    the Case-A threshold of ≥2) → overlay must be KEPT, not dropped.

    Guards against an off-by-one in the filled_body_slots threshold.
    """
    plan_slide = {
        "clone_from_slide": 34,
        "slots": {
            "title": "Single body slot",
            "body1": "lonely content",
            # sub*/body2/body3 deliberately omitted
        },
        "infographic": {
            "type": "matrix",
            "shapes": _matrix_shapes(),
        },
    }
    out = _build_one(tmp_workdir, plan_slide)
    prs = Presentation(out)
    txt = _slide_text(prs.slides[0])
    # Overlay markers MUST be present — only 1 body filled < 2 threshold.
    assert _OVERLAY_MARKER_A in txt, (
        "1 filled body slot is below Case-A threshold; overlay must remain"
    )


def test_case_a_preserves_filled_slot_text(tmp_workdir) -> None:
    """Case A: structural donor 34 with all 6 body slots filled + matrix
    overlay → overlay dropped (B5) AND filled slot text (alpha/beta/gamma)
    survives the non-title cleanup pass.

    Regression guard: an earlier draft of F1 wiped filled slots because
    `clear_donor_non_title_text` had no awareness of slot mapping.
    """
    plan_slide = {
        "clone_from_slide": 34,
        "slots": {
            "title": "Three columns",
            "sub1": "A", "sub2": "B", "sub3": "C",
            "body1": "alpha-CASE-A", "body2": "beta-CASE-A", "body3": "gamma-CASE-A",
        },
        "infographic": {
            "type": "matrix",
            "shapes": _matrix_shapes(),
        },
    }
    out = _build_one(tmp_workdir, plan_slide)
    prs = Presentation(out)
    txt = _slide_text(prs.slides[0])
    # Overlay dropped.
    assert _OVERLAY_MARKER_A not in txt
    assert _OVERLAY_MARKER_B not in txt
    # Filled slot content preserved through the F1 cleanup pass.
    assert "alpha-CASE-A" in txt, "Case A: filled body slot text was wiped"
    assert "beta-CASE-A" in txt, "Case A: filled body slot text was wiped"
    assert "gamma-CASE-A" in txt, "Case A: filled body slot text was wiped"

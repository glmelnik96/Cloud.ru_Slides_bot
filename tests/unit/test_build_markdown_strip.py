"""Residual markdown/control-char leak in title-like donor bodies (2026-06-07).

Spec review found a RESIDUAL LEAK that the per-slot chokepoint + apply_kpi_emphasis
do NOT close:

  * build_v5.replace_text_with_style writes donor body/title with
    ``strip_markdown=False`` so the ``**…**`` markers survive for the whole-deck
    ``apply_kpi_emphasis`` pass to consume.
  * BUT apply_kpi_emphasis SKIPS shapes whose first run is >= 28pt
    (``_shape_is_title_like``). Donors 21 (content_text_white) and 22
    (content_text_dark) define ``body: size_pt: 32`` — a 32pt body is title_like
    → skipped → its ``**`` is NEVER stripped → leaks as literal asterisks.

This reproduces session 81673 slide 5: ``**ССM (Cloud Certificate Manager)**``,
a single-paragraph conclusion body landing on donor 21 at 32pt.

The fix is a FINAL whole-deck markdown-strip pass that runs AFTER
apply_kpi_emphasis (so intentional emphasis is already consumed) over EVERY run
on EVERY shape — including title_like bodies that emphasis skipped.
"""
from __future__ import annotations

import json

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

from build_v9 import build, strip_residual_markdown  # noqa: E402


@pytest.fixture
def tmp_workdir(tmp_path):
    return tmp_path


def _all_run_texts(out_pptx) -> list[str]:
    """Every <a:t> text across every slide/shape/paragraph/run (and table cells)."""
    prs = Presentation(out_pptx)
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            sp = shape._element
            for t in sp.iter(qn("a:t")):
                texts.append(t.text or "")
    return texts


def _build(tmp_workdir, plan: dict):
    plan_path = tmp_workdir / "plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False), encoding="utf-8")
    out = tmp_workdir / "out.pptx"
    build(str(plan_path), str(skill_bridge.TEMPLATE_PATH), str(out),
          str(skill_bridge.DONOR_SLOT_MAP))
    return out


# ─── E2E: the original 81673 s5 bug ──────────────────────────────────────────

def test_donor21_32pt_body_strips_markdown_and_vtab(tmp_workdir) -> None:
    """Reproduce 81673 s5: a 32pt donor-21 body with ``**…**`` + a \\x0b in the
    title. After the full build path (which calls apply_kpi_emphasis), NO ``**``
    and NO \\x0b / ``_X000B_`` may appear in ANY <a:t>."""
    plan = {"slides": [{
        "clone_from_slide": 21,
        "slots": {
            "title": "Итог\x0bвторая строка",
            "body": "**ССM (Cloud Certificate Manager)**",
        },
    }]}
    out = _build(tmp_workdir, plan)
    texts = _all_run_texts(out)
    joined = "".join(texts)
    assert "**" not in joined, f"literal markdown bold leaked: {texts!r}"
    assert "\x0b" not in joined, f"vertical tab leaked: {texts!r}"
    assert "_X000B_" not in joined, f"escaped vtab leaked: {texts!r}"
    # The phrase content itself must survive (minus the markers).
    assert any("ССM (Cloud Certificate Manager)" in t for t in texts), texts


def test_donor22_32pt_dark_body_strips_markdown(tmp_workdir) -> None:
    """Donor 22 (dark, also 32pt body) — same title_like skip, must be clean."""
    plan = {"slides": [{
        "clone_from_slide": 22,
        "slots": {"title": "Вывод", "body": "**ключевой результат проекта**"},
    }]}
    out = _build(tmp_workdir, plan)
    joined = "".join(_all_run_texts(out))
    assert "**" not in joined
    assert any("ключевой результат проекта" in t for t in _all_run_texts(out))


# ─── Boundary: title_like bodies at exactly 28pt and 32pt come out clean ─────

@pytest.mark.parametrize("size_pt", [28, 32])
def test_title_like_body_boundary_strips(size_pt) -> None:
    """A title_like body (first run >= 28pt) is SKIPPED by apply_kpi_emphasis,
    so the final strip pass must still consume its ``**``. Verify the exact
    28pt boundary and the 32pt donor-21/22 size both come out clean."""
    from pptx import Presentation as _P
    from pptx.util import Pt

    prs = _P(str(skill_bridge.TEMPLATE_PATH))
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(Pt(50), Pt(50), Pt(400), Pt(80))
    tf = box.text_frame
    tf.text = "**ССM (Cloud Certificate Manager)**"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)

    n = strip_residual_markdown(prs)
    assert n >= 1
    full = box.text_frame.text
    assert "**" not in full, full
    assert "ССM (Cloud Certificate Manager)" in full


def test_strip_residual_markdown_is_idempotent() -> None:
    """Running the pass twice leaves text untouched the second time."""
    from pptx import Presentation as _P
    from pptx.util import Pt

    prs = _P(str(skill_bridge.TEMPLATE_PATH))
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(Pt(50), Pt(50), Pt(400), Pt(80))
    tf = box.text_frame
    tf.text = "**жирный** текст"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)

    assert strip_residual_markdown(prs) >= 1
    assert strip_residual_markdown(prs) == 0  # idempotent


# ─── Regression guard: emphasis on <28pt bodies must still bold ──────────────

def _emphasized_run_count(out_pptx) -> int:
    """Count bold runs (emphasis applied b='1') across the deck."""
    prs = Presentation(out_pptx)
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame._txBody.findall(qn("a:p")):
                for r in p.findall(qn("a:r")):
                    rPr = r.find(qn("a:rPr"))
                    if rPr is not None and rPr.get("b") == "1":
                        n += 1
    return n


def test_normal_body_still_gets_emphasized(tmp_workdir) -> None:
    """A normal <28pt body with ``**phrase**`` must still be BOLDED by
    apply_kpi_emphasis, and end up with NO literal ``**``. Donor 28 body is
    20pt → not title_like → emphasis runs."""
    plan = {"slides": [{
        "clone_from_slide": 28,
        "slots": {
            "title": "Преимущества",
            "col1_body": "Платформа даёт **полную изоляцию данных** для клиентов",
            "col2_body": "обычный текст без разметки",
        },
    }]}
    out = _build(tmp_workdir, plan)
    joined = "".join(_all_run_texts(out))
    assert "**" not in joined, "markers must be consumed by emphasis/strip"
    assert any("полную изоляцию данных" in t for t in _all_run_texts(out))
    assert _emphasized_run_count(out) >= 1, "emphasis must still bold <28pt body"

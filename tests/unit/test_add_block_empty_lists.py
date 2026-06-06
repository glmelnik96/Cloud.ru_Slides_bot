"""add_block must not crash when font_sizes/bolds arrive as empty lists.

Regression: dl2 build aborted with IndexError at flow_renderer add_block when a
block carried non-empty `lines` but an explicit empty `bolds`/`font_sizes` list
(the `is None` guard skipped them, then `bolds[-1]` indexed an empty list).
"""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402

from flow_renderer import add_block  # noqa: E402


def _blank_slide():
    prs = Presentation()
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def test_empty_bolds_and_font_sizes_do_not_crash():
    slide = _blank_slide()
    shape = add_block(slide, 100, 100, 400, 200,
                      ["Заголовок", "Тело"], font_sizes=[], bolds=[])
    # Two paragraphs rendered, defaults applied — no IndexError.
    assert len(shape.text_frame.paragraphs) == 2

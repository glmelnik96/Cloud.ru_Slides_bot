"""A2: long card TITLES must shrink-to-fit instead of clipping.

Two gaps closed:
- render_card_grid chip branch: title sits in one chip-height row next to the
  green chip and was rendered at full title_size with no fit.
- render_numbered_rows: title box height assumed ~1 line; a wrapped 2-3 line
  title overlapped the body below.
"""
from __future__ import annotations

import pytest

from worker import skill_bridge

skill_bridge.install()

from flow_renderer import (  # noqa: E402
    GEOFIT_AVAILABLE,
    render_card_grid,
    render_numbered_rows,
)

_LONG_TITLE = (
    "Очень длинный заголовок карточки который никак не помещается в одну "
    "строку рядом с чипом и обязан уменьшиться"
)


def _blank_slide():
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def _run_sizes_by_text(slide, needle: str) -> list[float]:
    """Font sizes (pt) of runs inside shapes whose text contains *needle*."""
    sizes = []
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        if needle not in shp.text_frame.text:
            continue
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    sizes.append(r.font.size.pt)
    return sizes


def test_card_grid_chip_title_shrinks():
    """Chip branch: a long title next to the chip must come out below the
    20pt default instead of clipping in the single chip-height row."""
    if not GEOFIT_AVAILABLE:
        pytest.skip("textfit/Pillow unavailable")
    slide = _blank_slide()
    render_card_grid(slide, {"cols": 2, "cards": [
        {"title": _LONG_TITLE, "text": "тело", "number": 1},
        {"title": "Коротко", "text": "тело", "number": 2},
    ]})
    long_sizes = _run_sizes_by_text(slide, "Очень длинный")
    short_sizes = _run_sizes_by_text(slide, "Коротко")
    assert long_sizes and all(s < 20.0 for s in long_sizes), (
        f"long chip title not shrunk: {long_sizes}"
    )
    # The short title must stay at the default.
    assert short_sizes and all(s == 20.0 for s in short_sizes)


def test_card_grid_chip_short_title_untouched():
    if not GEOFIT_AVAILABLE:
        pytest.skip("textfit/Pillow unavailable")
    slide = _blank_slide()
    render_card_grid(slide, {"cols": 2, "cards": [
        {"title": "Коротко", "text": "тело", "check": True},
        {"title": "Тоже коротко", "text": "тело", "check": True},
    ]})
    sizes = _run_sizes_by_text(slide, "оротко")
    assert sizes and all(s == 20.0 for s in sizes)


def test_numbered_rows_long_title_body_below_title():
    """numbered_rows: with a wrapped multi-line title (narrow 2-col layout)
    the title box must grow to its REAL line count and the body must start
    below it — the legacy code sized the title box for ~1 line."""
    if not GEOFIT_AVAILABLE:
        pytest.skip("textfit/Pillow unavailable")
    slide = _blank_slide()
    render_numbered_rows(slide, {"cols": 2, "rows": [
        {"title": _LONG_TITLE, "text": "тело строки"},
        {"title": "Б", "text": "т"},
        {"title": "В", "text": "т"},
        {"title": "Г", "text": "т"},
    ]})
    boxes = []
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = shp.text_frame.text
        if not txt:
            continue
        boxes.append((shp.top, shp.top + shp.height, shp.height, txt))
    title_box = next(b for b in boxes if "Очень длинный" in b[3])
    body_box = next(b for b in boxes if "тело строки" in b[3])
    short_title_box = next(b for b in boxes if b[3] == "Б")
    # The wrapped title box is at least 2 single-line heights tall.
    assert title_box[2] >= 2 * short_title_box[2], (
        f"title box did not grow: {title_box[2]} vs single-line {short_title_box[2]}"
    )
    assert body_box[0] >= title_box[1], (
        f"body top {body_box[0]} overlaps title bottom {title_box[1]}"
    )


def test_numbered_rows_short_title_layout_unchanged():
    """Single-line titles keep the legacy geometry (body at ~1.5 line heights)."""
    if not GEOFIT_AVAILABLE:
        pytest.skip("textfit/Pillow unavailable")
    slide = _blank_slide()
    render_numbered_rows(slide, {"cols": 1, "rows": [
        {"title": "Коротко", "text": "тело строки"},
    ]})
    sizes = _run_sizes_by_text(slide, "Коротко")
    # 17pt default title preserved (no shrink for a fitting title).
    assert sizes and all(s == 17.0 for s in sizes)

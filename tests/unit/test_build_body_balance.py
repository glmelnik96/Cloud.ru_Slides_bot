"""B (2026-06-07): vertical balance for underfilled body slots.

Taxonomy defect B: content sits in the top 35-40% of the slide, the bottom
55-60% is empty. Donor body boxes (21/22) are large; short text top-anchored
leaves the slide top-heavy. Fix: when the wrapped body block underfills the box
(textfit balance), anchor it MIDDLE instead of the forced TOP — geometry is not
moved, only the vertical anchor.
"""
from __future__ import annotations

import json

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402
from pptx.enum.text import MSO_ANCHOR  # noqa: E402

from build_v9 import build  # noqa: E402


@pytest.fixture
def tmp_workdir(tmp_path):
    return tmp_path


def _body_anchor(out_pptx):
    """Donor 21 body lives on shape_idx=1 (per donor-slot-map.yaml)."""
    prs = Presentation(out_pptx)
    shapes = list(prs.slides[0].shapes)
    if len(shapes) < 2 or not shapes[1].has_text_frame:
        return None
    return shapes[1].text_frame.vertical_anchor


def _build_with_body(tmp_workdir, body: str):
    plan = {"slides": [{"clone_from_slide": 21,
                        "slots": {"title": "Итоги", "body": body}}]}
    plan_path = tmp_workdir / "plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False), encoding="utf-8")
    out = tmp_workdir / "out.pptx"
    build(str(plan_path), str(skill_bridge.TEMPLATE_PATH), str(out),
          str(skill_bridge.DONOR_SLOT_MAP))
    return out


def test_short_body_centres_vertically(tmp_workdir) -> None:
    """A one-line conclusion underfills donor 21's large body box → MIDDLE."""
    out = _build_with_body(tmp_workdir, "Выручка выросла на 40%.")
    assert _body_anchor(out) == MSO_ANCHOR.MIDDLE


def test_full_body_stays_top(tmp_workdir) -> None:
    """A body that nearly fills the box keeps the forced TOP anchor — no
    centring, so multi-line text reads from the top as designed."""
    long_body = (
        "Платформа Cloud.ru показала уверенный рост по всем ключевым "
        "направлениям: облачная инфраструктура, управляемые сервисы и "
        "решения для искусственного интеллекта вышли на плановые объёмы, "
        "а число корпоративных клиентов выросло кратно за отчётный период."
    )
    out = _build_with_body(tmp_workdir, long_body)
    assert _body_anchor(out) == MSO_ANCHOR.TOP

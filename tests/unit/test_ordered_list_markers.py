"""D2 (2026-06-07): ordered-list markers (1./2./3.) must be preserved through
the flow renderer.

Defect: numeric list markers were not detected, so an ordered list either lost
its numbering (the square auto-bullet would replace the visible "N.") or merged
into a flat paragraph. Fix: ``_detect_bullet`` recognises numeric markers and
``add_block`` renders them as list items (hanging indent) while keeping the
literal "N." prefix as text — no square glyph, no double-numbering. Decimals
like "3.5" are NOT treated as markers.
"""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from flow_renderer import _detect_bullet, add_block  # noqa: E402


def test_detect_numeric_marker_kept_as_text():
    """A leading "1." is recognised as an ordered list item but the number is
    preserved in the clean text (not stripped, not replaced)."""
    is_bullet, clean = _detect_bullet("1. Первый пункт")
    # Not a square-glyph bullet (that would replace the visible number).
    assert is_bullet is False
    # Literal number preserved.
    assert clean == "1. Первый пункт"


def test_detect_paren_numeric_marker_kept_as_text():
    is_bullet, clean = _detect_bullet("2) Второй пункт")
    assert is_bullet is False
    assert clean == "2) Второй пункт"


def test_decimal_not_treated_as_marker():
    """A decimal like "3.5" at the start of a line is NOT an ordered marker."""
    is_bullet, clean = _detect_bullet("3.5 раза быстрее")
    assert is_bullet is False
    assert clean == "3.5 раза быстрее"


def _para_texts_and_indents(shape):
    out = []
    for p in shape.text_frame.paragraphs:
        txt = "".join(r.text for r in p.runs)
        pPr = p._p.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
        marL = pPr.get("marL") if pPr is not None else None
        out.append((txt, marL))
    return out


def test_ordered_list_preserved_through_render():
    """An ordered list survives parse→render: three numbered items render as
    three distinct list paragraphs with the numbers intact and no flattening."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    lines = ["1. Подготовка данных", "2. Обучение модели", "3. Валидация"]
    shape = add_block(slide, 50, 50, 400, 200, lines)
    rendered = _para_texts_and_indents(shape)
    texts = [t for t, _ in rendered]
    # Numbering intact, no merge, no double-numbering.
    assert texts == lines
    # Rendered as list items: hanging indent applied to numbered paragraphs.
    indents = [marL for _, marL in rendered]
    assert all(m is not None for m in indents), "ordered items get hanging indent"

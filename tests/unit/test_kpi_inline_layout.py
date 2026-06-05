"""D3: KPI inline layout — number left / description right, vertically centered."""
from __future__ import annotations

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402

from kpi_renderer import render_kpi, clean_slide_to_blank, BLANK_DONOR_WHITE  # noqa: E402


def _blank():
    prs = Presentation(str(skill_bridge.TEMPLATE_PATH))
    slide = list(prs.slides)[BLANK_DONOR_WHITE - 1]
    clean_slide_to_blank(slide)
    return prs, slide


def _textboxes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def test_desc_sits_right_of_number():
    prs, slide = _blank()
    render_kpi(slide, {"title": "ИТОГ", "numbers": [
        {"value": "84", "desc": "вовлечённость команды"}]})
    boxes = _textboxes(slide)
    num = next(b for b in boxes if b.text_frame.text == "84")
    desc = next(b for b in boxes if "вовлечённость" in b.text_frame.text)
    # Description box must start to the right of the number box.
    assert desc.left > num.left
    # And overlap the number vertically (inline, not stacked below).
    n_top, n_bot = num.top, num.top + num.height
    d_mid = desc.top + desc.height // 2
    assert n_top <= d_mid <= n_bot


def test_pct_is_enlarged():
    prs, slide = _blank()
    render_kpi(slide, {"title": "ИТОГ", "numbers": [
        {"value": "99", "desc": "аптайм", "pct": True}]})
    boxes = [s for s in slide.shapes if s.has_text_frame
             and s.text_frame.text.strip() == "%"]
    assert len(boxes) == 1
    pct = boxes[0]
    # Enlarged %: font ≈ 0.5× number height. For the single hero (150pt)
    # the % must be ≥ 80pt (was max(40, 199//3)=66).
    sz = pct.text_frame.paragraphs[0].runs[0].font.size.pt
    assert sz >= 80


@pytest.mark.parametrize("n", [1, 2, 3])
def test_columns_do_not_overlap(n):
    prs, slide = _blank()
    nums = [{"value": str(10 + i), "desc": f"метрика {i}"} for i in range(n)]
    render_kpi(slide, {"title": "T", "numbers": nums})
    # Number boxes, left-to-right, must not overlap horizontally.
    num_boxes = sorted(
        [s for s in slide.shapes if s.has_text_frame
         and s.text_frame.text.strip().isdigit()],
        key=lambda b: b.left)
    for a, b in zip(num_boxes, num_boxes[1:]):
        assert a.left + a.width <= b.left + 1

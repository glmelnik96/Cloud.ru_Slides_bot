"""Recursive group walk: nested text + pictures are recovered."""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from parse_pptx import _walk_shapes  # noqa: E402


def _deck_flat(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(914400), Emu(914400))
    tb.text_frame.text = "TOP"
    out = tmp_path / "flat.pptx"
    prs.save(out)
    return str(out)


def test_walk_shapes_flat_returns_leaves(tmp_path):
    path = _deck_flat(tmp_path)
    prs = Presentation(path)
    leaves = _walk_shapes(prs.slides[0].shapes)
    texts = [lf["text"] for lf in leaves if lf["text"]]
    assert "TOP" in texts
    assert all({"shape_type", "text", "left", "top", "w", "h", "depth"} <= set(lf) for lf in leaves)


def test_walk_shapes_recurses_into_group_fixture():
    import os
    fixture = os.path.join(os.path.dirname(__file__), "..", "fixtures", "grouped_diagram.pptx")
    prs = Presentation(fixture)
    leaves = _walk_shapes(prs.slides[0].shapes)
    texts = [lf["text"] for lf in leaves if lf["text"]]
    assert "Node A" in texts and "Node B" in texts and "Node C" in texts
    node_c = next(lf for lf in leaves if lf["text"] == "Node C")
    assert node_c["depth"] >= 2


def test_parse_emits_visual_kind_and_group_nodes(tmp_path):
    import os
    from parse_pptx import parse
    fixture = os.path.join(os.path.dirname(__file__), "..", "fixtures", "grouped_diagram.pptx")
    result = parse(fixture)
    s = result["slides"][0]
    assert "visual_kind" in s
    assert "group_nodes" in s
    # The fixture has 3 grouped text nodes and no normal text → structured.
    assert s["visual_kind"] == "structured"
    assert len(s["group_nodes"]) == 3
    assert all("order" in n and "text" in n for n in s["group_nodes"])

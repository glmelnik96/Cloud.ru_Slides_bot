"""D4 part 1: clean_slide_to_blank can preserve decorative pics outside the
content zone, while still removing pics that overlap content."""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from kpi_renderer import clean_slide_to_blank, CONTENT_ZONE_EMU  # noqa: E402

_PX = 9525


def _png(tmp_path):
    from PIL import Image
    p = tmp_path / "dot.png"
    Image.new("RGB", (10, 10), (0, 200, 120)).save(p)
    return str(p)


def _slide_with_pic(prs, left_px, top_px, w_px, h_px, png):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.add_picture(png, Emu(left_px * _PX), Emu(top_px * _PX),
                             Emu(w_px * _PX), Emu(h_px * _PX))
    return slide


def _count_pics(slide):
    return sum(1 for s in slide.shapes if s.shape_type == 13)  # PICTURE


def test_default_strips_all_pics(tmp_path):
    png = _png(tmp_path)
    prs = Presentation()
    slide = _slide_with_pic(prs, 1100, 20, 150, 150, png)  # corner decor
    clean_slide_to_blank(slide)  # keep_decor defaults False
    assert _count_pics(slide) == 0


def test_keep_decor_preserves_corner_pic(tmp_path):
    png = _png(tmp_path)
    prs = Presentation()
    # Top-right decor above the content zone (bottom edge 100px < zone top 120px).
    slide = _slide_with_pic(prs, 1100, 10, 140, 90, png)
    clean_slide_to_blank(slide, keep_decor=True)
    assert _count_pics(slide) == 1


def test_keep_decor_removes_content_zone_pic(tmp_path):
    png = _png(tmp_path)
    prs = Presentation()
    # Pic squarely inside the content zone — must be removed even with keep_decor.
    slide = _slide_with_pic(prs, 300, 300, 600, 300, png)
    clean_slide_to_blank(slide, keep_decor=True)
    assert _count_pics(slide) == 0

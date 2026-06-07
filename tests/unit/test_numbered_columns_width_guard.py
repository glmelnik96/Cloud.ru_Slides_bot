"""Task 6: numbered_columns enforces a minimum column width.

``render_numbered_columns`` lays out ``n`` open columns at ``cw = SAFE_W/n``.
With many columns ``cw`` shrinks until typical Russian words break mid-word.
The guard enforces ``_MIN_COL_W``: when ``n`` columns would each be narrower
than the minimum, the renderer falls back to ``render_numbered_rows`` (a 2-row
grid of stacked number+title+text) so words wrap on spaces, never inside a word.

Behaviour is asserted via the real rendered geometry (textbox left edges) — no
mocks — and by checking every item's number survives the fallback.
"""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()

from pptx import Presentation  # noqa: E402

from flow_renderer import (  # noqa: E402
    render_numbered_columns, _MIN_COL_W, SAFE_W, SAFE_LEFT,
)

_PX = 9525


def _blank_slide():
    prs = Presentation()
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def _col_width_for(n: int, gap: int = 32) -> int:
    """The exact width formula the renderer uses for ``n`` columns."""
    return int((SAFE_W - (n - 1) * gap) / n)


def _distinct_left_bases(slide, tol: int = 100) -> list[int]:
    """Cluster textbox left edges (px) into distinct column x-bases.

    A single layout column may emit several textboxes at slightly different
    x (e.g. a number at the base and its title indented by ~44px); ``tol``
    groups those into one column while keeping genuinely separate columns
    (hundreds of px apart) distinct.
    """
    lefts = sorted({round(s.left / _PX) for s in slide.shapes
                    if s.has_text_frame})
    bases: list[int] = []
    for x in lefts:
        if not bases or x - bases[-1] > tol:
            bases.append(x)
    return bases


def _all_text(slide) -> str:
    return "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def test_min_width_violation_falls_back_to_rows():
    """8 columns would each be ~123px (< _MIN_COL_W) → fall back to rows.

    The fallback must (a) NOT lay the items out in 8 narrow columns and
    (b) preserve every item.
    """
    n = 8
    assert _col_width_for(n) < _MIN_COL_W  # precondition: the guard must trigger
    cols = [{"title": f"Заголовок {i}",
             "text": "взаимодействие инфраструктура",
             "number": "%02d" % (i + 1)} for i in range(n)]
    slide = _blank_slide()
    render_numbered_columns(slide, {"columns": cols})

    # Fallback (numbered_rows, 2 cols) uses at most 2 column x-bases, never 8.
    bases = _distinct_left_bases(slide)
    assert len(bases) <= 2, f"expected rows fallback, got {len(bases)} columns"

    # Every item survived the fallback — nothing dropped.
    text = _all_text(slide)
    for i in range(n):
        assert ("%02d" % (i + 1)) in text
        assert f"Заголовок {i}" in text


def test_small_n_keeps_columns():
    """3 columns are ~382px wide (>= _MIN_COL_W) → columns unchanged."""
    n = 3
    assert _col_width_for(n) >= _MIN_COL_W  # precondition: no fallback
    cols = [{"title": f"Заголовок {i}",
             "text": "взаимодействие",
             "number": "%02d" % (i + 1)} for i in range(n)]
    slide = _blank_slide()
    render_numbered_columns(slide, {"columns": cols})

    # Columns layout: n distinct x-bases at SAFE_LEFT + i*(cw+gap).
    bases = _distinct_left_bases(slide)
    assert len(bases) == n, f"expected {n} columns, got {len(bases)}"
    assert bases[0] == SAFE_LEFT

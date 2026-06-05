"""D4 part 2: short card grids are vertically centered in the safe area."""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402

from flow_renderer import render_card_grid, SAFE_TOP  # noqa: E402

_PX = 9525


def _blank_slide():
    prs = Presentation()
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def _topmost_panel_top_px(slide):
    tops = [round(s.top / _PX) for s in slide.shapes
            if s.shape_type == 1]  # AUTO_SHAPE (panel rects)
    return min(tops) if tops else None


def test_two_cards_are_centered():
    slide = _blank_slide()
    render_card_grid(slide, {"cols": 2, "cards": [
        {"title": "A", "text": "short"}, {"title": "B", "text": "short"}]})
    top = _topmost_panel_top_px(slide)
    # A single-row grid (height ~one row) must be pushed below SAFE_TOP so the
    # remaining vertical slack is split above and below.
    assert top is not None
    assert top > SAFE_TOP + 20


def test_full_grid_starts_at_safe_top():
    slide = _blank_slide()
    cards = [{"title": str(i), "text": "x"} for i in range(8)]  # 4 rows × 2
    render_card_grid(slide, {"cols": 2, "cards": cards})
    top = _topmost_panel_top_px(slide)
    # A full grid fills the area — no meaningful centering offset.
    assert top is not None
    assert abs(top - SAFE_TOP) <= 8

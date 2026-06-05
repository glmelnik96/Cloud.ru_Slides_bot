"""D10: enforce_canonical text-color fixes for dark slides.

Live run 2026-06-05 produced run2.slide1 with green subtitle text on
a dark background — invisible-ish, definitely off-brand. enforce_canonical
recolored green→graphite, but on dark slides graphite is dark-on-dark,
still invisible. This test pins the dark-context branch: green→white,
not green→graphite.
"""
from __future__ import annotations

import pytest

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from pptx.util import Emu, Pt  # noqa: E402

from enforce_canonical import enforce_canonical_slide  # noqa: E402


@pytest.fixture
def blank_slide():
    prs = Presentation(str(skill_bridge.TEMPLATE_PATH))
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs, prs.slides.add_slide(layout)


def _add_textbox_with_color(slide, text: str, hex_color: str, *, size_pt: int = 14):
    box = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(2000000), Emu(500000))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.color.rgb = RGBColor.from_string(hex_color)
    return box


def _run_color_hex(box) -> str | None:
    """Return the explicit srgbClr value on the first run, or None."""
    txBody = box.text_frame._txBody
    for p in txBody.findall(qn("a:p")):
        for r in p.findall(qn("a:r")):
            rPr = r.find(qn("a:rPr"))
            if rPr is None:
                continue
            sf = rPr.find(qn("a:solidFill"))
            if sf is None:
                continue
            srgb = sf.find(qn("a:srgbClr"))
            if srgb is not None:
                return (srgb.get("val") or "").upper()
    return None


def test_green_text_on_dark_becomes_white(blank_slide) -> None:
    """D10: dark slide + green text → recolor to white (not graphite)."""
    prs, slide = blank_slide
    box = _add_textbox_with_color(slide, "Памятка по сигналам", "26D07C")
    stats = enforce_canonical_slide(slide, dark=True)
    assert stats["green_text"] >= 1
    assert _run_color_hex(box) == "FFFFFF"


def test_green_text_on_light_becomes_graphite(blank_slide) -> None:
    """Light slide preserves prior behaviour: green → graphite."""
    prs, slide = blank_slide
    box = _add_textbox_with_color(slide, "Зелёный текст", "26D07C")
    stats = enforce_canonical_slide(slide, dark=False)
    assert stats["green_text"] >= 1
    assert _run_color_hex(box) == "222222"


def test_white_text_on_dark_kept(blank_slide) -> None:
    """White on dark should NOT be recolored."""
    prs, slide = blank_slide
    box = _add_textbox_with_color(slide, "Белый текст", "FFFFFF")
    stats = enforce_canonical_slide(slide, dark=True)
    assert stats["white_on_light"] == 0
    assert _run_color_hex(box) == "FFFFFF"


def test_white_text_on_light_becomes_graphite(blank_slide) -> None:
    """Existing behaviour preserved: white text on light slide → graphite."""
    prs, slide = blank_slide
    box = _add_textbox_with_color(slide, "Белый текст", "FFFFFF")
    stats = enforce_canonical_slide(slide, dark=False)
    assert stats["white_on_light"] >= 1
    assert _run_color_hex(box) == "222222"

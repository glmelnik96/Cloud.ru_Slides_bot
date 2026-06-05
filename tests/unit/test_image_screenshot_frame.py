"""D5: screenshot images get a brand browser-chrome frame."""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()
from pptx import Presentation  # noqa: E402

from image_renderer import render_image_native  # noqa: E402


def _png(tmp_path):
    from PIL import Image
    p = tmp_path / "shot.png"
    Image.new("RGB", (1200, 700), (240, 240, 240)).save(p)
    return str(p)


def _blank_slide():
    prs = Presentation()
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def _green_strips(slide):
    from pptx.dml.color import RGBColor
    out = []
    for s in slide.shapes:
        if s.shape_type != 1:  # AUTO_SHAPE
            continue
        try:
            if s.fill.type is not None and s.fill.fore_color.rgb == RGBColor(0x26, 0xD0, 0x7C):
                out.append(s)
        except Exception:
            pass
    return out


def test_screenshot_adds_green_titlebar(tmp_path):
    slide = _blank_slide()
    render_image_native(slide, {
        "title": "Консоль управления",
        "image_path": _png(tmp_path),
        "frame": "browser",
    })
    assert len(_green_strips(slide)) >= 1
    # Picture is still present.
    assert any(s.shape_type == 13 for s in slide.shapes)  # PICTURE


def test_plain_image_has_no_chrome(tmp_path):
    slide = _blank_slide()
    render_image_native(slide, {
        "title": "Фото офиса",
        "image_path": _png(tmp_path),
    })
    assert len(_green_strips(slide)) == 0
    assert any(s.shape_type == 13 for s in slide.shapes)

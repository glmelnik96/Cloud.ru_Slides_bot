"""#2: card_grid body text must shrink-to-fit and, as a last resort, truncate
with an ellipsis so a long card body never clips past the card's bottom edge.

The fit logic is extracted into the pure helper ``_fit_card_body`` so the
shrink/truncate behaviour can be unit-tested directly without inspecting XML.
"""
from __future__ import annotations

from worker import skill_bridge

skill_bridge.install()

from flow_renderer import (  # noqa: E402
    _fit_card_body, render_card_grid, GEOFIT_AVAILABLE, _CARD_BODY_MIN_PT,
)


_BASE = 15.0


def test_short_body_unchanged():
    """A short body fits the box at the base size — no shrink, no ellipsis."""
    size, text = _fit_card_body("Короткий текст", 260, 120, _BASE)
    assert size == _BASE
    assert text == "Короткий текст"
    assert "…" not in text


def test_long_body_shrinks_below_base():
    """A long body in a fixed box shrinks below the 15pt default."""
    if not GEOFIT_AVAILABLE:
        import pytest
        pytest.skip("textfit/Pillow unavailable")
    long = ("Очень длинный текст карточки, который не помещается в "
            "стандартном кегле и обязан уменьшиться, " * 2)
    size, text = _fit_card_body(long, 260, 110, _BASE)
    assert size < _BASE
    assert size >= _CARD_BODY_MIN_PT


def test_pathological_body_truncates_with_ellipsis():
    """Text that overflows even at the min font is truncated with an ellipsis
    on a word boundary so nothing clips off the card."""
    if not GEOFIT_AVAILABLE:
        import pytest
        pytest.skip("textfit/Pillow unavailable")
    blob = "слово " * 400  # ~2400 chars, cannot fit a small card at any size
    size, text = _fit_card_body(blob, 240, 90, _BASE)
    assert size == _CARD_BODY_MIN_PT
    assert text.endswith("…")
    assert len(text) < len(blob)
    # Word-boundary truncation: no broken trailing token before the ellipsis.
    assert text.rstrip("…").strip().split()[-1] == "слово"


def test_truncation_not_applied_when_shrink_suffices():
    """If shrinking alone makes it fit, do not add an ellipsis."""
    if not GEOFIT_AVAILABLE:
        import pytest
        pytest.skip("textfit/Pillow unavailable")
    medium = "Текст средней длины для карточки, требует лёгкого уменьшения кегля"
    size, text = _fit_card_body(medium, 260, 120, _BASE)
    assert "…" not in text


def test_d1_shrink_preferred_over_truncation_keeps_trailing_token():
    """D1: a body that fits after a modest shrink (above the min) is NOT
    ellipsis-truncated — the trailing token (e.g. a number) is preserved.

    Sized so the text overflows at the 15pt base but fits comfortably after
    shrinking toward, but not down to, the floor. The old code jumped straight
    to min-size truncation when fit_text's conservative-margin size still
    failed the local capacity check, dropping the last word ("112")."""
    if not GEOFIT_AVAILABLE:
        import pytest
        pytest.skip("textfit/Pillow unavailable")
    body = ("Этот блок карточки содержит несколько фраз и завершается важным "
            "числовым показателем 112")
    # Box forces a shrink below the 15pt base, but the text still fits whole at
    # a size above the floor — so it must NOT reach the truncation branch.
    size, text = _fit_card_body(body, 260, 120, _BASE)
    assert size < _BASE, "this box must engage the shrink path"
    assert size > _CARD_BODY_MIN_PT, "fits above the floor — bounded shrink"
    assert "…" not in text, "modest shrink must avoid the ellipsis branch"
    assert text == body, "no tokens dropped"
    assert text.endswith("112"), "trailing number preserved"


def test_render_card_grid_long_body_does_not_raise():
    """End-to-end: rendering a grid with a long card body succeeds and emits
    a body run whose font size is reduced below the 15pt default."""
    if not GEOFIT_AVAILABLE:
        import pytest
        pytest.skip("textfit/Pillow unavailable")
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    long = "слово " * 120
    render_card_grid(slide, {"cols": 2, "cards": [
        {"title": "CVE", "text": long},
        {"title": "B", "text": "short"},
        {"title": "C", "text": "short"},
        {"title": "D", "text": "short"},
    ]})
    sizes = []
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    sizes.append(r.font.size.pt)
    # Some run must be smaller than the default body size (shrink engaged).
    assert any(s < 15.0 for s in sizes)

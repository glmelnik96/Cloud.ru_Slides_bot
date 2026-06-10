"""Check 1 — can the LLM author on-brand HTML from REAL slide content?

Pulls text from a real test deck (read-only), asks the LLM (reusing the
SLIDE_COMPOSER role) to author a freeform slide body using only brand.css
classes, renders it via the spike renderer, saves a PNG to view.

Isolated: imports llm.client / llm.roles read-only; writes only out/.
"""
from __future__ import annotations

import pathlib
import re
import sys

from pptx import Presentation

from llm.client import LLMCall, call_role
from llm.roles import Role

HERE = pathlib.Path(__file__).resolve().parent
REPO = HERE.parent.parent
OUT = HERE / "out"
OUT.mkdir(exist_ok=True)

DECK = REPO / "tmp" / "live_inputs" / "6851c8d0f0674088.pptx"


def extract_slide_text(path: pathlib.Path, idx: int) -> dict:
    prs = Presentation(str(path))
    slide = prs.slides[idx]
    title = ""
    bullets: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            txt = "".join(r.text for r in para.runs).strip()
            if not txt:
                continue
            if not title:
                title = txt
            else:
                bullets.append(txt)
    return {"title": title, "bullets": bullets[:7]}


SYSTEM = """Ты — арт-директор Cloud.ru 2.0. Верстаешь ОДИН слайд 1280×720 на чистом HTML,
используя ТОЛЬКО классы и переменные из приложенного brand.css. Верни ТОЛЬКО HTML-фрагмент:
один <div class="slide ...">…</div>. Без markdown, без ```html, без <html>/<head>.

Каноны бренда (строго):
- Фон белый (var(--white)). Плоский дизайн: НИКАКИХ градиентов, теней, скруглений, italic, подчёркиваний.
- Шрифт — SB Sans Display (уже в .slide). Заголовок слайда = .slide__header (CAPS, 20pt SemiBold), слева сверху.
- Лого .brand-logo справа сверху. Копирайт .copyright внизу.
- Зелёный #26D07C — ЕДИНСТВЕННЫЙ акцент, применять точечно (тонкая черта-правило над пунктом, квадратный буллет, плашка).
- Текст графитовый #222. Вторичный текст — var(--text-gray).
- Сетка: левое поле 35px. Контент в безопасной зоне (top≈150, bottom≈660).
- Можно использовать готовые компоненты из brand.css (.points/.bullets и их потомки) ИЛИ собрать раскладку inline-стилями НА БАЗЕ переменных бренда.
- Текст бери ДОСЛОВНО из контента, не переписывай и не переводи.
Выбери раскладку, лучше всего подходящую под объём контента (колонка пунктов с зелёными чертами, или список с квадратными буллетами)."""


def author_html(content: dict, brand_css: str) -> str:
    user = (
        f"brand.css:\n```css\n{brand_css}\n```\n\n"
        f"Контент слайда (JSON):\n{content}\n\n"
        "Свёрстай слайд. Верни только HTML-фрагмент."
    )
    res = call_role(LLMCall(
        role=Role.SLIDE_COMPOSER,
        messages=[{"role": "system", "content": SYSTEM}, {"role": "user", "content": user}],
        max_tokens_override=4000,
    ))
    html = res.content.strip()
    html = re.sub(r"^```(?:html)?\s*", "", html)
    html = re.sub(r"\s*```$", "", html)
    print(f"[llm] {res.model} elapsed={res.elapsed_s:.1f}s out_tokens={res.completion_tokens}", file=sys.stderr)
    return html


def main() -> int:
    idx = int(sys.argv[1]) if len(sys.argv) > 1 else 3
    content = extract_slide_text(DECK, idx)
    print(f"[content] title={content['title']!r} bullets={len(content['bullets'])}", file=sys.stderr)
    brand_css = (HERE / "brand.css").read_text(encoding="utf-8")
    html = author_html(content, brand_css)
    tpl = HERE / "templates" / "_llm_check1.html"
    tpl.write_text(html, encoding="utf-8")
    print(f"[saved] {tpl}", file=sys.stderr)
    print("Now run: python probes/html_spike/render.py", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

"""Pack rendered slide PNGs full-bleed into a 16:9 .pptx.

The HTML-render pipeline produces one PNG per slide; this packs them, one image
per slide, edge-to-edge on a blank 13.333"×7.5" (16:9) canvas. Output is a
picture deck by design — non-editable, pixel-faithful to the Chromium render.

The HTML source for each slide is preserved alongside (caller's concern) so the
deck is reproducible; this module only assembles the visual artefact.
"""
from __future__ import annotations

import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# 16:9 at the python-pptx default EMU scale.
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)


def pack_pngs(pngs: list[bytes], out_path: str | Path) -> str:
    """Write a .pptx with each PNG placed full-bleed on its own slide.

    ``pngs`` is a list of PNG byte strings in slide order. Returns the output
    path as a string.
    """
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    blank = prs.slide_layouts[6]  # the truly empty layout (no placeholders)

    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(png), 0, 0, width=_SLIDE_W, height=_SLIDE_H,
        )

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)

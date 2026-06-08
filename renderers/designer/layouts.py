"""Archetype layout skeletons for the designer skill.

Each skeleton owns its slide's layout: it places brand-dressed zones at fixed
coordinates and fills them from a content dict. The composer chooses an
archetype + supplies CONTENT (not coordinates); these functions render it so
output matches the Cloud.ru reference instead of a blind grid raffle.

A skeleton signature is always (slide, content: dict, *, dark: bool) -> None.
content keys are archetype-specific and documented per function.
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt

from renderers.designer import primitives as P

CANVAS_W, CANVAS_H = 1280, 720
M = 60  # cover/section outer margin


# ─── Brand chrome (vector wordmark + footer) ─────────────────────────────────

def brand_logo(slide, *, corner: str = "top_right", dark_bg: bool = False,
               cube_color=None):
    """Minimal vector cloud.ru wordmark: green cube glyph + "cloud.ru" text.

    No raster — a small green square stands in for the cube mark next to the
    wordmark, matching the template's top-corner lockup at slide scale.
    cube_color overrides the green mark (e.g. graphite on a green cover, where
    a green cube would vanish).
    """
    txt_color = P.WHITE if dark_bg else P.GRAPHITE
    cube_rgb = cube_color if cube_color is not None else P.GREEN
    cube = 22
    if corner == "top_left":
        cube_x, cube_y = M, 40
        txt_x = cube_x + cube + 10
    else:  # top_right
        txt_w = 120
        txt_x = CANVAS_W - M - txt_w
        cube_x = txt_x - cube - 10
        cube_y = 40
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, P.px(cube_x), P.px(cube_y),
                                P.px(cube), P.px(cube))
    P._solid(sq, cube_rgb)
    P._no_line(sq)
    tb = slide.shapes.add_textbox(P.px(txt_x), P.px(cube_y - 4), P.px(130), P.px(cube + 8))
    tf = tb.text_frame
    P._zero_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "cloud.ru"
    r.font.name = P.SEMIBOLD_FONT
    r.font.size = Pt(16)
    r.font.color.rgb = txt_color
    return tb


def _content_chrome(slide, *, dark=False):
    """Common chrome for content slides: top-right wordmark."""
    brand_logo(slide, corner="top_right", dark_bg=dark)


# ─── Covers / dividers ───────────────────────────────────────────────────────

def cover_green(slide, content, *, dark=False):
    """Full green fill + portal staircase + dot grid + display title.

    content: {"title": str, "subtitle": str?}
    """
    P.background(slide, "green")
    brand_logo(slide, corner="top_left", dark_bg=False, cube_color=P.GRAPHITE)
    # Dot grid bottom-left, portal staircase bottom-right (template p.1).
    P._dot_pattern(slide, bbox=(M, 470, 520, 700))
    P.portal(slide, (1000, 600, 90), n=4)
    P.display_title(slide, content.get("title") or "", (M, 150, 900, 320),
                    color=P.GRAPHITE)
    sub = (content.get("subtitle") or "").strip()
    if sub:
        P.body_block(slide, [sub], (M, 480, 760, 90), size_pt=20)


def cover_dark(slide, content, *, dark=True):
    """Dark cover: graphite fill, green-outline title box w/ green title text,
    gray descriptor plate, portal + dot grid (template p.6).

    content: {"title": str, "subtitle": str?}
    """
    P.background(slide, "graphite")
    brand_logo(slide, corner="top_left", dark_bg=True)
    # Green outline title box.
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, P.px(M), P.px(150),
                                 P.px(CANVAS_W - 2 * M), P.px(190))
    box.fill.background()
    box.line.color.rgb = P.GREEN
    box.line.width = P.Emu(int(2 * P.EMU_PER_PX))
    P._no_shadow(box)
    P.display_title(slide, content.get("title") or "", (M + 24, 175, CANVAS_W - 2 * M - 48, 150),
                    color=P.GREEN, max_pt=64)
    # Gray descriptor plate.
    sub = (content.get("subtitle") or "").strip()
    plate = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, P.px(M), P.px(370),
                                   P.px(640), P.px(80))
    P._solid(plate, P.RGBColor.from_string("3A3A3A"))
    P._no_line(plate)
    if sub:
        P.body_block(slide, [sub], (M + 20, 388, 600, 50), size_pt=18, dark_bg=True)
    P.portal(slide, (980, 470, 80), n=4, dark_bg=True)


def section_divider(slide, content, *, dark=True):
    """Section break: dark fill + portal + kicker + large section title.

    content: {"title": str, "kicker": str?}
    """
    P.background(slide, "graphite" if dark else "white")
    brand_logo(slide, corner="top_right", dark_bg=dark)
    P.portal(slide, (M, 560, 110), n=3, dark_bg=dark)
    kicker = (content.get("kicker") or "").strip()
    if kicker:
        P.body_block(slide, [kicker.upper()], (M, 230, 700, 50), size_pt=16, dark_bg=dark)
    P.display_title(slide, content.get("title") or "", (M, 290, 1000, 260),
                    dark_bg=dark, max_pt=72)


# ─── Content ─────────────────────────────────────────────────────────────────

def _points_grid(slide, content, cols, rows, *, dark=False):
    """Title band + a cols×rows grid of point_item modules (green divider +
    bold head + body). content: {"title": str, "points": [{"head","text"}...]}.
    Fills row-major; ignores extra points beyond cols*rows.
    """
    P.background(slide, "graphite" if dark else "white")
    _content_chrome(slide, dark=dark)
    P.title_block(slide, content.get("title") or "", (40, 60, 1000, 90),
                  dark_bg=dark, size_pt=34)
    pts = (content.get("points") or [])[: cols * rows]
    area_top, area_bottom = 210, CANVAS_H - 70
    area_h = area_bottom - area_top
    gap_x, gap_y = 30, 36
    cell_w = (CANVAS_W - 2 * 40 - gap_x * (cols - 1)) / cols
    cell_h = (area_h - gap_y * (rows - 1)) / rows
    for i, pt in enumerate(pts):
        r, c = divmod(i, cols)
        left = 40 + c * (cell_w + gap_x)
        top = area_top + r * (cell_h + gap_y)
        P.point_item(slide, pt.get("head") or "", pt.get("text") or "",
                     (left, top, cell_w, cell_h), dark_bg=dark)


def points_3(slide, content, *, dark=False):
    _points_grid(slide, content, cols=3, rows=1, dark=dark)


def points_4(slide, content, *, dark=False):
    _points_grid(slide, content, cols=2, rows=2, dark=dark)


def points_6(slide, content, *, dark=False):
    _points_grid(slide, content, cols=3, rows=2, dark=dark)


def points_8(slide, content, *, dark=False):
    _points_grid(slide, content, cols=4, rows=2, dark=dark)


def bullet_list(slide, content, *, dark=False):
    """Workhorse: title + optional intro + green-tick bullet column.

    content: {"title": str, "intro": str?, "bullets": [str,...]}
    """
    P.background(slide, "graphite" if dark else "white")
    _content_chrome(slide, dark=dark)
    P.title_block(slide, content.get("title") or "", (40, 60, 1000, 90),
                  dark_bg=dark, size_pt=34)
    intro = (content.get("intro") or "").strip()
    top = 210
    if intro:
        P.body_block(slide, [intro], (40, top, 1160, 70), size_pt=18, dark_bg=dark)
        top += 90
    P.body_block(slide, content.get("bullets") or [],
                 (40, top, 1160, CANVAS_H - top - 70), size_pt=18, dark_bg=dark)


# ─── Special (wrap rich primitives) ──────────────────────────────────────────

def table_zebra(slide, content, *, dark=False):
    """Title + full-width branded zebra table.

    content: {"title": str, "headers": [str], "rows": [[str]], "accent_col": int?}
    """
    P.background(slide, "graphite" if dark else "white")
    _content_chrome(slide, dark=dark)
    P.title_block(slide, content.get("title") or "", (40, 60, 1000, 90),
                  dark_bg=dark, size_pt=34)
    P.table_block(slide, content.get("headers") or [], content.get("rows") or [],
                  (40, 190, 1200, CANVAS_H - 190 - 70),
                  accent_col=content.get("accent_col"), dark_bg=dark)


def chart_columns(slide, content, *, dark=False):
    """Title + clustered column chart (5-color brand ramp).

    content: {"title": str, "categories": [str], "series": [{"name","values"}],
              "accent_idx": int?, "data_provenance": str?}
    """
    P.background(slide, "graphite" if dark else "white")
    _content_chrome(slide, dark=dark)
    P.title_block(slide, content.get("title") or "", (40, 60, 1000, 90),
                  dark_bg=dark, size_pt=34)
    P.chart_block(slide, "bar", content.get("categories") or [],
                  content.get("series") or [], (40, 190, 1200, CANVAS_H - 190 - 70),
                  accent_idx=content.get("accent_idx", 0),
                  data_provenance=content.get("data_provenance", "native"), dark_bg=dark)


def roadmap_timeline(slide, content, *, dark=False):
    """Title + horizontal axis + evenly-spaced milestone chips.

    content: {"title": str, "milestones": [{"label","text","accent"?}]}
    """
    P.background(slide, "graphite" if dark else "white")
    _content_chrome(slide, dark=dark)
    P.title_block(slide, content.get("title") or "", (40, 60, 1000, 90),
                  dark_bg=dark, size_pt=34)
    ms = content.get("milestones") or []
    axis_y = 380
    P.timeline_axis(slide, axis_y, 60, 1220, dark_bg=dark)
    if ms:
        n = len(ms)
        span = (1220 - 60) / max(1, n)
        for i, m in enumerate(ms):
            left = 60 + i * span + 10
            P.milestone_tick(slide, m.get("label") or "", m.get("text") or "",
                             (left, axis_y - 150, span - 20, 130),
                             accent=bool(m.get("accent")), dark_bg=dark)

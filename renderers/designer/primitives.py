"""Native python-pptx vector primitives for the Cloud.ru 2.0 designer skill.

Every primitive draws editable native shapes (autoshapes, textboxes, native
charts, freeforms) — NO raster bake. Colors/fonts come from the brand glossary
in llm.prompts._shared so the skill and the prompts share one source of truth.

Brand fidelity (vs the q2 prototype): hero KPI sizing (≤150pt) with a green
accent plate and enlarged %, native green-tick bullets, CAPS section headers,
and geometric overflow protection (Pillow font metrics via the vendored
``textfit``) so text never clips or bleeds off-slide.
"""
from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree

# Brand glossary (single source of truth). Fall back to literals if the import
# path shifts during prototyping.
try:
    from llm.prompts._shared import (
        BRAND_PALETTE,
        EMU_PER_PX,
        PRIMARY_FONT,
        SEMIBOLD_FONT,
    )
except Exception:  # pragma: no cover - prototype safety net
    EMU_PER_PX = 9525
    PRIMARY_FONT = "SB Sans Display"
    SEMIBOLD_FONT = "SB Sans Display Semibold"
    BRAND_PALETTE = {
        "green": "#26D07C", "graphite": "#222222", "gray": "#F2F2F2",
        "white": "#FFFFFF", "stroke": "#C8C8C8",
    }

GREEN = RGBColor.from_string(BRAND_PALETTE["green"].lstrip("#"))
GRAPHITE = RGBColor.from_string(BRAND_PALETTE["graphite"].lstrip("#"))
GRAY = RGBColor.from_string(BRAND_PALETTE["gray"].lstrip("#"))
WHITE = RGBColor.from_string(BRAND_PALETTE["white"].lstrip("#"))
TEXT_GRAY = RGBColor.from_string("5C5C5C")   # secondary text (canon)
ARROW_GRAY = RGBColor.from_string("434343")  # connectors — never green
STROKE = RGBColor.from_string("C8C8C8")

# Non-accent chart ramp. The template's "диаграмма с дополнительными цветами"
# (slides 45/50/61) shows data viz IS allowed multiple colors: a green lead plus
# light BRAND tints (periwinkle, mint, pale-yellow, lilac) — never grays. Green
# stays the lead/accent; these tints fill the remaining series.
NON_ACCENT = [
    RGBColor.from_string("C0E0FC"),  # brand Blue (periwinkle)
    RGBColor.from_string("9FE6C2"),  # mint (light green tint)
    RGBColor.from_string("E8FB7A"),  # pale yellow (Yellow #CFF500 tint)
    RGBColor.from_string("D9C2FF"),  # lilac (Purple #A068FF tint)
]

PT_TO_PX = 96.0 / 72.0


def px(v: float) -> Emu:
    return Emu(int(round(v * EMU_PER_PX)))


def _no_line(shape) -> None:
    shape.line.fill.background()


# Effect tags that carry shadows/glow/reflection — Cloud.ru = ZERO effects.
_EFFECT_TAGS = (
    "a:effectLst", "a:effectDag", "a:outerShdw", "a:innerShdw",
    "a:prstShdw", "a:glow", "a:reflection", "a:softEdge",
)


def _no_shadow(shape) -> None:
    """Force zero effects on a shape (brand: no shadow/glow/reflection ever).

    Mirrors ``effects_util.strip_effects``: clears any inherited effects on the
    shape's spPr and neutralises the theme's effectRef (idx=0). PowerPoint's
    default preset autoshapes ship with an outer shadow via effectRef, so an
    empty ``<a:effectLst/>`` plus ``effectRef idx=0`` is required to kill it.
    """
    el = shape._element
    spPr = el.find(qn("p:spPr"))
    if spPr is None:
        spPr = el.find(qn("a:spPr"))
    if spPr is not None:
        for tag in _EFFECT_TAGS:
            for child in spPr.findall(qn(tag)):
                spPr.remove(child)
        spPr.append(spPr.makeelement(qn("a:effectLst"), {}))
    # Neutralise the style-level effectRef (the real source of preset shadows).
    style = el.find(qn("p:style"))
    if style is not None:
        eref = style.find(qn("a:effectRef"))
        if eref is not None:
            eref.set("idx", "0")
            for child in list(eref):
                eref.remove(child)


def _solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _no_shadow(shape)


def _square_corner(shape) -> None:
    """Clamp a rounded-rect to the brand max corner radius (≤4px)."""
    try:
        w_emu = int(shape.width)
        h_emu = int(shape.height)
        short = min(w_emu, h_emu) / EMU_PER_PX
        shape.adjustments[0] = min(4.0 / max(short, 1.0), 0.08)
    except Exception:  # pragma: no cover - guard against pptx version drift
        pass


# ─── Overflow protection (real font metrics) ─────────────────────────────────

def _fit(text: str, w_px: float, h_px: float | None, base_pt: float,
         min_pt: float, *, semibold: bool, wrap: bool = True,
         balance: bool = False):
    """Largest size ≤ base_pt whose wrapped text fits (w_px, h_px).

    Uses the vendored Pillow-based ``textfit`` when the skill scripts are on
    sys.path (worker / live_run); silently returns ``base_pt`` otherwise so the
    primitive still renders (a hair large) rather than crashing.

    Returns ``(size_pt, anchor_middle, lines)``.
    """
    try:
        import font_resolver
        import textfit
        font_path = font_resolver.resolve(
            SEMIBOLD_FONT if semibold else PRIMARY_FONT, bold=False)
        if not font_path:
            return base_pt, False, 1
        res = textfit.fit_text(
            text,
            box_w_emu=int(w_px * EMU_PER_PX),
            box_h_emu=int(h_px * EMU_PER_PX) if h_px else None,
            font_path=font_path,
            base_pt=base_pt,
            min_pt=min_pt,
            wrap=wrap,
            balance=balance,
        )
        if res is None:
            return base_pt, False, 1
        return res.size_pt, res.anchor_middle, res.lines
    except Exception:
        return base_pt, False, 1


def _zero_margins(tf) -> None:
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0


# ─── Backgrounds ─────────────────────────────────────────────────────────────

def background(slide, kind: str) -> None:
    """Full-bleed background rectangle (Group A brand plashka)."""
    color = {
        "white": WHITE, "graphite": GRAPHITE, "green": GREEN, "dots": WHITE,
    }.get(kind, WHITE)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, px(1280), px(720))
    _solid(rect, color)
    _no_line(rect)
    # Send to back so content draws over it.
    sp = rect._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    if kind == "dots":
        _dot_pattern(slide)


def _dot_pattern(slide, step_px: int = 120, dot_px: int = 4) -> None:
    """Sparse dot texture (brandbook p.26).

    Coarse and HARD-CAPPED so a dots background never exceeds ~60 shapes.
    """
    light = RGBColor.from_string("ECEEF0")
    cap, drawn = 60, 0
    y = 60
    while y < 700 and drawn < cap:
        x = 60
        while x < 1240 and drawn < cap:
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(dot_px), px(dot_px))
            _solid(d, light)
            _no_line(d)
            drawn += 1
            x += step_px
        y += step_px


def _accent_bar(slide, left, top, w, h, color=GREEN):
    """A flat brand accent plate (square corners, no line)."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(left), px(top), px(w), px(h))
    _solid(bar, color)
    _no_line(bar)
    return bar


# ─── Text ────────────────────────────────────────────────────────────────────

def title_block(slide, text: str, rect_px, size_pt: int = 44,
                accent_underline: bool = True, dark_bg: bool = False) -> None:
    """Heading: SemiBold, shrink-to-fit, with a green accent underline plate.

    The green underline is the ONE accent element — it sits directly under the
    fitted text block (its color is the accent, never the letters).
    """
    left, top, w, h = rect_px
    color = WHITE if dark_bg else GRAPHITE
    fit_pt, centred, lines = _fit(text, w, h, base_pt=float(size_pt), min_pt=20.0,
                                  semibold=True, wrap=True, balance=True)

    tb = slide.shapes.add_textbox(px(left), px(top), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if centred else MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = SEMIBOLD_FONT
    run.font.size = Pt(fit_pt)
    run.font.bold = False
    run.font.color.rgb = color

    if accent_underline:
        # Underline width tracks the heading scale (longer for display titles),
        # placed just below the LAST text line so it underlines the whole title.
        line_h = fit_pt * PT_TO_PX * 1.2
        block_h = line_h * max(1, lines)
        bar_w = min(w, max(140.0, fit_pt * 3.6))
        if centred:
            bar_top = top + (h - block_h) / 2 + block_h + 6
        else:
            bar_top = top + block_h + 6
        bar_top = min(bar_top, top + h - 8)
        _accent_bar(slide, left, bar_top, bar_w, 6, color=GREEN)
    return tb


def body_block(slide, bullets, rect_px, size_pt: int = 16, dark_bg: bool = False):
    """Body copy with native green-tick bullets, shrink-to-fit as a block."""
    left, top, w, h = rect_px
    color = WHITE if dark_bg else GRAPHITE
    items = [str(b) for b in bullets if str(b).strip()]
    if not items:
        return None

    # Fit the whole block: measure the longest bullet for width and the joined
    # text for height (line count drives vertical fit).
    joined = "\n".join(items)
    fit_pt, _, _ = _fit(joined, w - 28, h, base_pt=float(size_pt), min_pt=11.0,
                        semibold=False, wrap=True)
    gap_pt = max(4.0, fit_pt * 0.35)

    tb = slide.shapes.add_textbox(px(left), px(top), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = b
        run.font.name = PRIMARY_FONT
        run.font.size = Pt(fit_pt)
        run.font.color.rgb = color
        p.space_after = Pt(gap_pt)
        _native_bullet(p, fit_pt)
    return tb


def _native_bullet(paragraph, size_pt: float) -> None:
    """Give a paragraph a green square bullet + hanging indent (native XML)."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    indent = int(size_pt * PT_TO_PX * 1.4 * EMU_PER_PX)
    pPr.set("marL", str(indent))
    pPr.set("indent", str(-indent))
    # Drop any inherited bullet props, then add ours in schema order.
    for tag in ("a:buClr", "a:buSzPct", "a:buFont", "a:buChar", "a:buNone",
                "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = buClr.makeelement(qn("a:srgbClr"), {"val": "26D07C"})
    buClr.append(srgb)
    buSz = pPr.makeelement(qn("a:buSzPct"), {"val": "70000"})
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u25aa"})  # ▪
    for el in (buClr, buSz, buFont, buChar):
        pPr.append(el)


def kpi_block(slide, num: str, desc: str, rect_px, dark_bg: bool = False):
    """Hero KPI inside the cell: big SemiBold number, enlarged %, green accent
    plate under the number, description below. Ported from the skill's
    kpi_renderer sizing logic so the number fills the cell without overflow."""
    left, top, w, h = rect_px
    color = WHITE if dark_bg else GRAPHITE

    value = str(num).strip()
    has_pct = False
    if value.endswith("%"):
        value = value[:-1].rstrip()
        has_pct = True

    # Number occupies the upper ~62% of the cell; description the rest.
    num_h = h * 0.60
    desc_h = h - num_h - 6
    # Fit the number to the cell width on a single line (no wrap), with room for
    # the % glyph if present.
    avail_w = w * (0.84 if has_pct else 1.0)
    base = min(150.0, num_h / (PT_TO_PX * 1.1))
    fit_pt, _, _ = _fit(value or "0", avail_w, num_h, base_pt=base, min_pt=36.0,
                        semibold=True, wrap=False)

    num_box = slide.shapes.add_textbox(px(left), px(top), px(w), px(num_h))
    tf = num_box.text_frame
    tf.word_wrap = False
    _zero_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = value
    r.font.name = SEMIBOLD_FONT
    r.font.size = Pt(fit_pt)
    r.font.color.rgb = color

    # Estimated width of the number, to place the % and accent bar.
    num_w = min(w, fit_pt * PT_TO_PX * 0.62 * max(1, len(value)) + 8)

    if has_pct:
        pct_pt = max(24.0, fit_pt * 0.5)
        pct_box = slide.shapes.add_textbox(
            px(left + num_w - 4), px(top + 4), px(w - num_w + 8), px(num_h * 0.5))
        ptf = pct_box.text_frame
        _zero_margins(ptf)
        ptf.word_wrap = False
        pp = ptf.paragraphs[0]
        pr = pp.add_run()
        pr.text = "%"
        pr.font.name = SEMIBOLD_FONT
        pr.font.size = Pt(pct_pt)
        pr.font.color.rgb = color

    # Green accent plate under the number (the ONE accent: a plate, not letters).
    _accent_bar(slide, left, top + num_h - 4, min(num_w, w), 6, color=GREEN)

    desc = str(desc).strip()
    if desc and desc_h > 12:
        d_pt, _, _ = _fit(desc, w, desc_h, base_pt=16.0, min_pt=10.0,
                          semibold=False, wrap=True)
        d_box = slide.shapes.add_textbox(px(left), px(top + num_h + 6), px(w), px(desc_h))
        dtf = d_box.text_frame
        dtf.word_wrap = True
        _zero_margins(dtf)
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = desc
        dr.font.name = PRIMARY_FONT
        dr.font.size = Pt(d_pt)
        dr.font.color.rgb = TEXT_GRAY if not dark_bg else WHITE
    return num_box


# ─── Charts (native, Excel-editable) ─────────────────────────────────────────

_CHART_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,    # vertical columns (template s.45/47)
    "hbar": XL_CHART_TYPE.BAR_CLUSTERED,      # horizontal bars (template s.46)
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA_STACKED,       # stacked area (template s.50) — non-
                                              # stacked AREA hides smaller series
                                              # behind larger ones.
    "area_100": XL_CHART_TYPE.AREA_STACKED_100,
}


def chart_block(slide, chart_type, categories, series, rect_px, accent_idx=0,
                data_provenance="native", dark_bg=False):
    left, top, w, h = rect_px
    if data_provenance == "estimated":
        # A report must never silently present read-off numbers as exact.
        h = max(h - 24, 24)
        fn = slide.shapes.add_textbox(px(left), px(top + h), px(w), px(20))
        rp = fn.text_frame.paragraphs[0]
        rr = rp.add_run()
        rr.text = "оценка по графику"
        rr.font.name = PRIMARY_FONT
        rr.font.size = Pt(9)
        rr.font.italic = True
        rr.font.color.rgb = RGBColor.from_string("9AA4AD")
    data = CategoryChartData()
    data.categories = categories
    for s in series:
        data.add_series(s["name"], s["values"])
    gframe = slide.shapes.add_chart(
        _CHART_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
        px(left), px(top), px(w), px(h), data,
    )
    chart = gframe.chart
    chart.has_title = False
    txt_color = WHITE if dark_bg else GRAPHITE
    try:
        chart.font.name = PRIMARY_FONT
        chart.font.size = Pt(11)
        chart.font.color.rgb = txt_color
    except Exception:  # pragma: no cover
        pass
    # Recolor: exactly one green accent series, rest pastel (brandbook rule).
    plot = chart.plots[0]
    if chart_type == "pie":
        pts = plot.series[0].points
        for i, pt in enumerate(pts):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = GREEN if i == accent_idx else NON_ACCENT[i % len(NON_ACCENT)]
    else:
        for i, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = GREEN if i == accent_idx else NON_ACCENT[i % len(NON_ACCENT)]
    if len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    return gframe


# ─── Tables (native, editable zebra) ─────────────────────────────────────────

_TBL_NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # built-in "No Style, No Grid"
_TBL_SEP = RGBColor.from_string("434343")  # vertical separator (canon)


def _tbl_strip_style(table) -> None:
    """Swap the default applied table style for "No Style, No Grid" and drop the
    banded/first-row flags so our manual zebra fills are the only styling."""
    tblPr = table._tbl.find(qn("a:tblPr"))
    if tblPr is None:
        return
    for attr in ("firstRow", "bandRow", "lastRow", "firstCol", "lastCol", "bandCol"):
        if tblPr.get(attr):
            del tblPr.attrib[attr]
    for sid in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(sid)
    sid = etree.SubElement(tblPr, qn("a:tableStyleId"))
    sid.text = _TBL_NO_STYLE


def _tbl_cell_fill(cell, rgb) -> None:
    if rgb is None:
        cell.fill.background()
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb


def _tbl_cell_margins(cell, lr_px: int = 12, tb_px: int = 8) -> None:
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))
    tcPr.set("marL", str(lr_px * EMU_PER_PX))
    tcPr.set("marR", str(lr_px * EMU_PER_PX))
    tcPr.set("marT", str(tb_px * EMU_PER_PX))
    tcPr.set("marB", str(tb_px * EMU_PER_PX))


def _tbl_cell_text(cell, text, size_pt: float, semibold: bool, color: RGBColor) -> None:
    tf = cell.text_frame
    tf.clear()
    cell.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.name = SEMIBOLD_FONT if semibold else PRIMARY_FONT
    run.font.size = Pt(size_pt)
    run.font.bold = False
    run.font.color.rgb = color


def _row_lines(row, col_w, char_w) -> int:
    """Max wrapped-line count across a row's cells (drives that row's height)."""
    max_lines = 1
    for ci, cell in enumerate(row):
        if ci >= len(col_w):
            break
        usable = col_w[ci] - 24  # horizontal margins (12+12)
        cpl = max(1, int(usable / char_w))
        for seg in (str(cell) if cell is not None else "").split("\n"):
            seglen = len(seg) if seg else 1
            lines = max(1, -(-seglen // cpl))  # ceil
            max_lines = max(max_lines, lines)
    return max_lines


def _fit_table_layout(headers, rows, col_w, total_h,
                      start: int = 16, hard_min: int = 8,
                      header_min: int = 34, row_min: int = 26):
    """Pick the largest font (start→hard_min) plus VARIABLE per-row heights so
    the table's total height fits within ``total_h``. Mirrors the skill's
    table_renderer: short rows stay compact, only a long-text row grows tall,
    and the font shrinks until the SUM fits — so the table never bleeds past its
    box the way uniform min-heights do (LibreOffice treats row height as a
    minimum and expands wrapped rows).

    Returns ``(font_pt, header_h, [row_h, ...])``.
    """
    target = total_h * 0.92  # headroom: LO inflates wrapped rows a touch
    chosen = None
    for fpt in range(int(start), int(hard_min) - 1, -1):
        font_px = fpt * 4.0 / 3.0
        char_w = 0.62 * font_px
        line_h = 1.30 * font_px
        header_h = max(header_min, int(_row_lines(headers, col_w, char_w) * line_h + 16))
        row_heights = [max(row_min, int(_row_lines(r, col_w, char_w) * line_h + 16))
                       for r in rows]
        total = header_h + sum(row_heights)
        chosen = (fpt, header_h, row_heights, total)
        if total <= target:
            return fpt, header_h, row_heights
    # Did not fit even at hard_min — compress proportionally (best-effort).
    fpt, header_h, row_heights, total = chosen
    if total > target and total > 0:
        scale = target / total
        header_h = max(int(header_h * scale), 20)
        row_heights = [max(int(h * scale), 16) for h in row_heights]
    return fpt, header_h, row_heights


def table_block(slide, headers, rows, rect_px, accent_col=None,
                first_col_wider: bool = True, dark_bg: bool = False):
    """Native zebra table (template slide-56 style): transparent header row with
    SemiBold graphite text, body rows alternating gray/white, thin vertical
    separators. ``accent_col`` tints one column with a brand periwinkle (not
    green) to highlight it. Square cells, zero effects, autofit font."""
    left, top, w, h = rect_px
    n_cols = len(headers)
    if n_cols == 0 or not rows:
        return None
    # Normalise every row to n_cols (pad short / clip long) so add_table is safe.
    norm: list[list[str]] = []
    for r in rows:
        rr = [str(c) if c is not None else "" for c in r]
        if len(rr) < n_cols:
            rr += [""] * (n_cols - len(rr))
        elif len(rr) > n_cols:
            rr = rr[:n_cols]
        norm.append(rr)
    n_total = len(norm) + 1

    # Column widths (first label column 1.4x wider when asked).
    iw = int(w)
    if first_col_wider and n_cols >= 2:
        rest = iw / (1.4 + n_cols - 1)
        first = int(round(1.4 * rest))
        rest = int(round(rest))
        col_w = [first] + [rest] * (n_cols - 1)
        col_w[-1] = iw - sum(col_w[:-1])
    else:
        col_w = [iw // n_cols] * n_cols
        col_w[-1] = iw - sum(col_w[:-1])

    body_pt, header_h, row_heights = _fit_table_layout(
        [str(x) for x in headers], norm, col_w, h)
    table_h = header_h + sum(row_heights)  # ≤ h (kept inside the grid box)

    gframe = slide.shapes.add_table(n_total, n_cols, px(left), px(top), px(w), px(table_h))
    table = gframe.table
    _tbl_strip_style(table)
    for i, cw in enumerate(col_w):
        table.columns[i].width = px(cw)
    table.rows[0].height = px(header_h)
    for r in range(1, n_total):
        table.rows[r].height = px(row_heights[r - 1])

    head_color = WHITE if dark_bg else GRAPHITE
    base_fill = GRAPHITE if dark_bg else GRAY
    accent_tint = NON_ACCENT[0]  # periwinkle blue (template s.52 highlighted col)

    for c in range(n_cols):
        cell = table.cell(0, c)
        _tbl_cell_fill(cell, None)
        _tbl_cell_margins(cell)
        _tbl_cell_text(cell, headers[c], body_pt, True, head_color)

    for ri, row in enumerate(norm):
        zebra = base_fill if ri % 2 == 0 else None
        for c in range(n_cols):
            cell = table.cell(ri + 1, c)
            tinted = accent_col is not None and c == accent_col
            _tbl_cell_fill(cell, accent_tint if tinted else zebra)
            _tbl_cell_margins(cell)
            txt_color = GRAPHITE if tinted else (WHITE if dark_bg else GRAPHITE)
            _tbl_cell_text(cell, row[c], body_pt, False, txt_color)

    # Vertical separators as connector lines (proven cross-renderer visibility).
    col_x = [left]
    for cw in col_w:
        col_x.append(col_x[-1] + cw)
    sep_color = STROKE if dark_bg else _TBL_SEP
    for i in range(1, len(col_x) - 1):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, px(col_x[i]), px(top), px(col_x[i]), px(top + table_h))
        conn.line.color.rgb = sep_color
        conn.line.width = Emu(int(1 * EMU_PER_PX))
        _no_shadow(conn)
    return gframe


# ─── Decor (outline) ─────────────────────────────────────────────────────────

def outline_corner(slide, anchor: str, dark_bg: bool = False):
    """Thin L-shaped corner bracket (outline-decor obvyazka, ~2px)."""
    color = WHITE if dark_bg else GRAPHITE
    size, off, thick = 60, 40, 2
    pos = {
        "top_left": (off, off, 1, 1),
        "top_right": (1280 - off - size, off, -1, 1),
        "bottom_left": (off, 720 - off - size, 1, -1),
        "bottom_right": (1280 - off - size, 720 - off - size, -1, -1),
    }[anchor]
    x, y, sx, sy = pos
    # horizontal arm
    hor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(size), px(thick))
    _solid(hor, color); _no_line(hor)
    # vertical arm
    ver = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(thick), px(size))
    _solid(ver, color); _no_line(ver)


def sparkle(slide, anchor: str, dark_bg: bool = False):
    """Small 4-point spark (brandbook sparkle motif), drawn as a freeform."""
    color = WHITE if dark_bg else GREEN  # the one green accent on light bg
    off, s = 50, 34
    base = {
        "top_left": (off, off), "top_right": (1280 - off - s, off),
        "bottom_left": (off, 720 - off - s), "bottom_right": (1280 - off - s, 720 - off - s),
    }[anchor]
    bx, by = base
    cx, cy = bx + s / 2, by + s / 2
    fb = slide.shapes.build_freeform(px(cx), px(by), scale=1)
    pts = [
        (cx + s * 0.12, cy - s * 0.12), (bx + s, cy), (cx + s * 0.12, cy + s * 0.12),
        (cx, by + s), (cx - s * 0.12, cy + s * 0.12), (bx, cy),
        (cx - s * 0.12, cy - s * 0.12), (cx, by),
    ]
    fb.add_line_segments([(px(x), px(y)) for x, y in pts], close=True)
    shp = fb.convert_to_shape()
    _solid(shp, color)
    _no_line(shp)
    return shp


# ─── Portal (stepped black graphic, brandbook pp.29-33) ──────────────────────

# Staircase offset per step as a fraction of the square side (PDF: +~24% right,
# -~7.5% up). Squares are identical; later squares z-order on top.
_PORTAL_DX, _PORTAL_DY = 0.24, 0.075


_PORTAL_DARK = RGBColor.from_string("3A3A3A")  # visible staircase on graphite bg


def portal(slide, base_rect_px, n: int = 3, dark_bg: bool = False):
    """N identical squares in an up-right staircase (brand 'портал').

    base_rect_px gives the bottom-left square's (left, top, side) — width is the
    side; the staircase grows up and to the right from there. On a graphite
    background the canonical #222222 would be invisible, so a lighter graphite
    tonal square is used instead.
    """
    left, top, side = base_rect_px
    fill = _PORTAL_DARK if dark_bg else GRAPHITE
    shapes = []
    for i in range(max(1, n)):
        x = left + i * _PORTAL_DX * side
        y = top - i * _PORTAL_DY * side
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(side), px(side))
        _solid(sq, fill)
        _no_line(sq)
        shapes.append(sq)
    return shapes


# ─── Diagram primitives (nodes + arrows) ─────────────────────────────────────

def node_box(slide, text: str, rect_px, accent: bool = False, dark_bg: bool = False):
    """Labelled diagram node: square-corner plate + centered shrink-to-fit text."""
    left, top, w, h = rect_px
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(left), px(top), px(w), px(h))
    _square_corner(box)
    if accent:
        _solid(box, GREEN)
        txt_color = GRAPHITE
    else:
        _solid(box, GRAPHITE if dark_bg else GRAY)
        txt_color = WHITE if dark_bg else GRAPHITE
    _no_line(box)
    fit_pt, _, _ = _fit(str(text), w - 16, h - 12, base_pt=15.0, min_pt=9.0,
                        semibold=True, wrap=True)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = SEMIBOLD_FONT
    r.font.size = Pt(fit_pt)
    r.font.color.rgb = txt_color
    return box


def _set_arrow_head(connector) -> None:
    """Add a triangular tail arrowhead to a connector's line (XML, no API)."""
    ln = connector.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)


def arrow(slide, p0_px, p1_px, rhombus: bool = False, dark_bg: bool = False):
    """Directed connector p0 -> p1, square cap + triangle head.

    The connector line is brand grey (#434343, never green). If rhombus=True,
    seat a 45° green square at the midpoint as a brand backing (its fill MUST
    differ from the line color)."""
    x0, y0 = p0_px
    x1, y1 = p1_px
    line_color = WHITE if dark_bg else ARROW_GRAY
    if rhombus:
        s = 22
        mx, my = (x0 + x1) / 2 - s / 2, (y0 + y1) / 2 - s / 2
        dia = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, px(mx), px(my), px(s), px(s))
        _solid(dia, GREEN)
        _no_line(dia)
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x0), px(y0), px(x1), px(y1))
    conn.line.color.rgb = line_color
    conn.line.width = Emu(int(1 * EMU_PER_PX))  # 1px brand keyline
    _no_shadow(conn)
    _set_arrow_head(conn)
    return conn


# ─── Cards (team / comparison) ───────────────────────────────────────────────

def person_card(slide, heading: str, sub: str, rect_px, plate: bool = True,
                accent: bool = False, dark_bg: bool = False):
    """Card: gray backing plate + green accent keyline + heading + sub line."""
    left, top, w, h = rect_px
    if plate:
        pl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(left), px(top), px(w), px(h))
        _solid(pl, GRAPHITE if dark_bg else GRAY)
        _no_line(pl)
        # Accent is a thin left keyline plate, not a full green fill.
        if accent:
            _accent_bar(slide, left, top, 6, h, color=GREEN)
    pad = 16
    tb = slide.shapes.add_textbox(px(left + pad), px(top + pad), px(w - 2 * pad), px(h - 2 * pad))
    tf = tb.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    h_pt, _, _ = _fit(str(heading), w - 2 * pad, (h - 2 * pad) * 0.55,
                      base_pt=18.0, min_pt=11.0, semibold=True, wrap=True)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = heading
    r.font.name = SEMIBOLD_FONT
    r.font.size = Pt(h_pt)
    r.font.color.rgb = WHITE if dark_bg else GRAPHITE
    if sub:
        s_pt, _, _ = _fit(str(sub), w - 2 * pad, (h - 2 * pad) * 0.4,
                          base_pt=13.0, min_pt=9.0, semibold=False, wrap=True)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = sub
        r2.font.name = PRIMARY_FONT
        r2.font.size = Pt(s_pt)
        r2.font.color.rgb = TEXT_GRAY if not dark_bg else WHITE
    return tb


# ─── Timeline ────────────────────────────────────────────────────────────────

def milestone_tick(slide, label: str, text: str, rect_px, accent: bool = False,
                   dark_bg: bool = False):
    """One milestone: small square tick + label + one-line text under it."""
    left, top, w, h = rect_px
    tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(left), px(top), px(12), px(12))
    _solid(tick, GREEN if accent else (WHITE if dark_bg else GRAPHITE))
    _no_line(tick)
    tb = slide.shapes.add_textbox(px(left), px(top + 18), px(w), px(h - 18))
    tf = tb.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    l_pt, _, _ = _fit(str(label), w, (h - 18) * 0.45, base_pt=20.0, min_pt=12.0,
                      semibold=True, wrap=False)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = label
    r.font.name = SEMIBOLD_FONT
    r.font.size = Pt(l_pt)
    r.font.color.rgb = WHITE if dark_bg else GRAPHITE
    if text:
        t_pt, _, _ = _fit(str(text), w, (h - 18) * 0.5, base_pt=12.0, min_pt=9.0,
                          semibold=False, wrap=True)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = text
        r2.font.name = PRIMARY_FONT
        r2.font.size = Pt(t_pt)
        r2.font.color.rgb = TEXT_GRAY if not dark_bg else WHITE
    return tb


def timeline_axis(slide, y_px: int, x0_px: int = 60, x1_px: int = 1220,
                  dark_bg: bool = False):
    """Thin horizontal axis line for the timeline archetype."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x0_px), px(y_px),
                                      px(x1_px), px(y_px))
    conn.line.color.rgb = WHITE if dark_bg else STROKE
    conn.line.width = Emu(int(1 * EMU_PER_PX))
    _no_shadow(conn)
    return conn

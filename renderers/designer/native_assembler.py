"""Deterministic Composition(DSL) -> native python-pptx slide assembler.

The LLM emits Composition objects on a 12x10 grid; this module turns each
block into native vector shapes via `primitives`. No EMU, no placeholder
indices ever reach the LLM. This is the q2-prototype renderer.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Emu

from renderers.designer import layouts as L
from renderers.designer import primitives as P
from renderers.designer.composition_dsl import (
    Composition,
    GRID_COLS,
    GRID_ROWS,
    Grid,
)

# Archetype name -> skeleton fn. When a Composition carries a known ``layout``
# the skeleton owns the whole slide and the free-grid path is bypassed.
_LAYOUTS = {
    "cover_green": L.cover_green,
    "cover_dark": L.cover_dark,
    "section_divider": L.section_divider,
    "points_3": L.points_3,
    "points_4": L.points_4,
    "points_6": L.points_6,
    "points_8": L.points_8,
    "bullet_list": L.bullet_list,
    "table_zebra": L.table_zebra,
    "chart_columns": L.chart_columns,
    "roadmap_timeline": L.roadmap_timeline,
}

CANVAS_W, CANVAS_H = 1280, 720
MARGIN = 40  # safe-area margin in px (brandbook micro-module multiple)
# Brand safe-area bottom is 660px; reserve a larger bottom margin so a block in
# the last grid row (r=10) never bleeds past it (was ending at 680px).
MARGIN_BOTTOM = 60


def _snap2(v: float) -> float:
    """Snap to the 2px brand micromodule (every offset divisible by 2)."""
    return round(v / 2.0) * 2.0


def _grid_to_px(g: Grid):
    """Convert a grid span to a px rect inside the safe area (2px-snapped)."""
    usable_w = CANVAS_W - 2 * MARGIN
    usable_h = CANVAS_H - MARGIN - MARGIN_BOTTOM
    cell_w = usable_w / GRID_COLS
    cell_h = usable_h / GRID_ROWS
    left = MARGIN + (g.c - 1) * cell_w
    top = MARGIN + (g.r - 1) * cell_h
    w = g.cs * cell_w
    h = g.rs * cell_h
    return (_snap2(left), _snap2(top), _snap2(w), _snap2(h))


def _edge_point(rect, toward):
    """Point on ``rect``'s border along the line from its center to ``toward``.

    Keeps connector endpoints on the node boundary so arrows touch the boxes
    instead of cutting through their centers.
    """
    left, top, w, h = rect
    cx, cy = left + w / 2, top + h / 2
    dx, dy = toward[0] - cx, toward[1] - cy
    if dx == 0 and dy == 0:
        return (cx, cy)
    hw, hh = w / 2, h / 2
    sx = hw / abs(dx) if dx else float("inf")
    sy = hh / abs(dy) if dy else float("inf")
    s = min(sx, sy)
    return (cx + dx * s, cy + dy * s)


# Roles whose blocks occupy area and participate in the de-overlap pass.
_AREA_ROLES = {"title", "body", "kpi", "chart", "table", "node", "card", "milestone"}
_DEOVERLAP_GAP = 12.0   # px vertical gap inserted between separated blocks
_MIN_BLOCK_H = 24.0     # px floor when compressing a block to fit the safe area


def _x_overlap(a, b) -> bool:
    """True when two rects' horizontal spans intersect (share columns)."""
    al, _, aw, _ = a
    bl, _, bw, _ = b
    return al < bl + bw and bl < al + aw


def _deoverlap_rects(items):
    """Vertical-reflow de-overlap for content blocks.

    ``items`` is a list of (key, rect) where rect == (left, top, w, h). Returns
    a dict key -> adjusted rect. Blocks are grouped into horizontal clusters
    (transitively column-overlapping); each cluster is swept top-to-bottom and a
    block colliding with the running bottom of the cluster is pushed down by
    GAP, or compressed if it would breach the safe-area bottom. Blocks that do
    not collide with anything are left untouched (no-op when nothing overlaps).
    """
    rect_for = {key: rect for key, rect in items}
    n = len(items)
    if n < 2:
        return rect_for

    # Union-find over horizontally-overlapping blocks -> independent clusters
    # so side-by-side (non-x-overlapping) columns never affect each other.
    parent = list(range(n))

    def find(i):
        while parent[i] != i:
            parent[i] = parent[parent[i]]
            i = parent[i]
        return i

    def union(i, j):
        parent[find(i)] = find(j)

    for i in range(n):
        for j in range(i + 1, n):
            if _x_overlap(items[i][1], items[j][1]):
                union(i, j)

    clusters: dict[int, list[int]] = {}
    for i in range(n):
        clusters.setdefault(find(i), []).append(i)

    bottom_limit = CANVAS_H - MARGIN_BOTTOM
    for members in clusters.values():
        if len(members) < 2:
            continue
        # Stable top-to-bottom sweep within the cluster.
        members.sort(key=lambda i: (items[i][1][1], items[i][1][0]))
        prev_bottom = None
        for i in members:
            key, (left, top, w, h) = items[i][0], rect_for[items[i][0]]
            if prev_bottom is not None and top < prev_bottom:
                # Collides with the running bottom -> push this block down.
                new_top = _snap2(prev_bottom + _DEOVERLAP_GAP)
                if new_top + h > bottom_limit:
                    # Cannot fit at full height: clamp + compress.
                    new_top = min(new_top, bottom_limit - _MIN_BLOCK_H)
                    new_top = _snap2(max(MARGIN, new_top))
                    h = _snap2(max(_MIN_BLOCK_H, bottom_limit - new_top))
                top = new_top
                rect_for[key] = (left, top, w, h)
            prev_bottom = top + h
    return rect_for


_MIN_CLAMP_W = 60.0   # px floor when shrinking a block to fit the right margin
_MIN_CLAMP_H = 24.0   # px floor when shrinking a block to fit the bottom margin


def _clamp_rect(left, top, w, h):
    """Conform a px rect to the safe area [MARGIN, CANVAS_W-MARGIN] x
    [MARGIN, CANVAS_H-MARGIN_BOTTOM].

    The LLM occasionally emits grid spans that overrun the right edge
    (c+cs-1 > 12) or the bottom edge (r+rs-1 > 10), so the corresponding px
    rect renders off-canvas and is clipped. This pulls every rect back inside
    the safe area: first shrink to fit, then (if the origin itself is past the
    margin) slide the origin back, honouring minimum width/height floors. It is
    a strict no-op for rects already fully inside the safe area.
    """
    right_limit = CANVAS_W - MARGIN          # 1240
    bottom_limit = CANVAS_H - MARGIN_BOTTOM   # 660

    # Left/top into bounds first.
    if left < MARGIN:
        left = MARGIN
    if top < MARGIN:
        top = MARGIN

    # Right edge: shrink width, then pull left back if still overrunning.
    if left + w > right_limit:
        w = right_limit - left
        if w < _MIN_CLAMP_W:
            w = _MIN_CLAMP_W
            left = right_limit - w
            if left < MARGIN:
                left = MARGIN
                w = right_limit - left

    # Bottom edge: shrink height, then pull top back if still overrunning.
    if top + h > bottom_limit:
        h = bottom_limit - top
        if h < _MIN_CLAMP_H:
            h = _MIN_CLAMP_H
            top = bottom_limit - h
            if top < MARGIN:
                top = MARGIN
                h = bottom_limit - top

    return (_snap2(left), _snap2(top), _snap2(w), _snap2(h))


def _portal_base(anchor: str, side: float = 200.0):
    """Bottom-left (left, top, side) for a portal staircase anchored to a corner.

    The staircase grows up-and-right, so leave room above and to the right.
    """
    pad = 40
    grow_w = side * (1 + 2 * 0.24)
    grow_up = side * (1 + 2 * 0.075)
    pos = {
        "top_left": (pad, pad + grow_up - side),
        "top_right": (CANVAS_W - pad - grow_w, pad + grow_up - side),
        "bottom_left": (pad, CANVAS_H - pad - side),
        "bottom_right": (CANVAS_W - pad - grow_w, CANVAS_H - pad - side),
    }[anchor]
    return (pos[0], pos[1], side)


# Deterministic content re-flow ------------------------------------------------
# GLM frequently lays repeated same-role blocks (kpi/card/milestone) out as a
# diagonal "staircase" or overlaps them, and the vertical de-overlap pass can't
# fix a diagonal (the blocks don't share columns, so they cluster separately).
# When a slide carries a group of >=2 such blocks we discard the LLM grid
# placement for that slide and lay the content area out deterministically. The
# threshold is 2 so two-column comparisons (2 cards/kpis) become two equal
# side-by-side cells instead of one card overrunning the right margin.
_GRID_ROLES = ("kpi", "card", "milestone")
_REFLOW_MIN = 2
_REFLOW_GAP = 16.0
# Non-title area roles that flow as full-width rows above/below the grid group.
# 'node' is intentionally excluded: connectors route to node rects, so nodes
# keep their original placement.
_FLOW_ROLES = {"body", "kpi", "chart", "table", "card", "milestone"}


def _cols_for(role: str, n: int) -> int:
    """Column count for a uniform grid of ``n`` same-role blocks."""
    if role == "kpi":
        return n if n <= 4 else 3
    if role == "milestone":
        return n if n <= 5 else 3  # timeline reads better as a single row
    if n <= 1:
        return 1
    return 2 if n <= 6 else 3


def _reflow_slide(blocks, rect_for):
    """Re-flow the content area when a same-role grid group of >=_REFLOW_MIN exists.

    Title band stays; non-title area blocks flow top-to-bottom in document
    order — body/chart/table blocks as full-width rows, the grid group as one
    uniform N×M grid placed where its first member appears. Returns an updated
    rect map; a strict no-op when no trigger group is present.
    """
    counts: dict[str, int] = {}
    for b in blocks:
        if b.role in _GRID_ROLES and id(b) in rect_for:
            counts[b.role] = counts.get(b.role, 0) + 1
    grid_role = next((r for r in _GRID_ROLES if counts.get(r, 0) >= _REFLOW_MIN), None)
    if grid_role is None:
        return rect_for

    # Reserve the title band from the PRIMARY (topmost) title only. The LLM
    # sometimes emits a second 'title' block as a mid-slide section header; that
    # one is flowed as a full-width row below (not counted in the band), so its
    # low position can't inflate content_top and crush the content beneath it.
    titles = [b for b in blocks if b.role == "title" and id(b) in rect_for]
    primary_title = min(titles, key=lambda b: rect_for[id(b)][1]) if titles else None
    title_bottom = MARGIN
    if primary_title is not None:
        _l, t, _w, h = rect_for[id(primary_title)]
        title_bottom = t + h
    content_top = _snap2(title_bottom + _REFLOW_GAP) if title_bottom > MARGIN else MARGIN
    bottom_limit = CANVAS_H - MARGIN_BOTTOM
    left0 = float(MARGIN)
    full_w = float(CANVAS_W - 2 * MARGIN)

    # Ordered content slots; the grid group appears once, at its first member.
    # Secondary titles flow as full-width rows (rendered as titles in place).
    group_members = [b for b in blocks if b.role == grid_role and id(b) in rect_for]
    ordered: list[tuple[str, object]] = []
    seen_group = False
    for b in blocks:
        if b is primary_title or id(b) not in rect_for:
            continue
        is_flow = b.role in _FLOW_ROLES or b.role == "title"
        if not is_flow:
            continue
        if b.role == grid_role:
            if not seen_group:
                ordered.append(("grid", group_members))
                seen_group = True
            continue
        ordered.append(("body", b))

    avail = bottom_limit - content_top
    gaps = _REFLOW_GAP * max(0, len(ordered) - 1)
    body_slots = [s for s in ordered if s[0] == "body"]
    body_h = [rect_for[id(b)][3] for _, b in body_slots]
    body_sum = sum(body_h)
    grid_min = 120.0
    max_body_total = avail - gaps - grid_min
    if body_sum > max_body_total and body_sum > 0:
        scale = max(0.0, max_body_total) / body_sum
        body_h = [max(48.0, h * scale) for h in body_h]
        body_sum = sum(body_h)
    grid_h = max(grid_min, avail - gaps - body_sum)

    new_rect = dict(rect_for)
    y = float(content_top)
    bi = 0
    for kind, payload in ordered:
        if kind == "body":
            h = body_h[bi]
            bi += 1
            new_rect[id(payload)] = (_snap2(left0), _snap2(y),
                                     _snap2(full_w), _snap2(h))
            y += h + _REFLOW_GAP
        else:
            members = payload
            n = len(members)
            cols = _cols_for(grid_role, n)
            rows = (n + cols - 1) // cols
            cell_w = (full_w - _REFLOW_GAP * (cols - 1)) / cols
            cell_h = (grid_h - _REFLOW_GAP * (rows - 1)) / rows
            for i, mb in enumerate(members):
                gl = left0 + (i % cols) * (cell_w + _REFLOW_GAP)
                gt = y + (i // cols) * (cell_h + _REFLOW_GAP)
                new_rect[id(mb)] = (_snap2(gl), _snap2(gt),
                                    _snap2(cell_w), _snap2(cell_h))
            y += grid_h + _REFLOW_GAP
    return new_rect


def assemble_slide(prs: Presentation, comp: Composition) -> None:
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    dark = comp.tone == "dark" or comp.background.kind == "graphite"

    # Skeleton mode: a known archetype owns its full layout. The skeleton paints
    # its own background/chrome, so return before the free-grid path runs.
    if comp.layout and comp.layout in _LAYOUTS:
        _LAYOUTS[comp.layout](slide, comp.content or {}, dark=dark)
        return

    P.background(slide, comp.background.kind)

    # De-overlap pass: compute adjusted px rects for all CONTENT (area) blocks
    # ONCE, keyed by id(blk). The same adjusted rects feed both the node
    # first-pass (so connectors route to the de-overlapped positions) and the
    # render loop, keeping connectors and shapes consistent.
    area_items = [
        (id(blk), _grid_to_px(blk.grid))
        for blk in comp.blocks if blk.role in _AREA_ROLES
    ]
    rect_for = _deoverlap_rects(area_items)

    # Safe-area conformance clamp: applied AFTER de-overlap so the final
    # adjusted rects are guaranteed inside [MARGIN, CANVAS_W-MARGIN] x
    # [MARGIN, CANVAS_H-MARGIN_BOTTOM]. The LLM sometimes emits grid spans that
    # overrun the right/bottom edge; this pulls them back so nothing renders
    # off-canvas. No-op for blocks already in bounds. Nodes (used for connector
    # routing below) read from this same clamped map, so arrows still touch the
    # clamped node boxes.
    rect_for = {key: _clamp_rect(*rect) for key, rect in rect_for.items()}

    # Deterministic re-flow: when the slide carries a same-role group of >=3
    # grid blocks (kpi/card/milestone) the LLM's diagonal/overlapping placement
    # is discarded and the content area is laid out as a clean stack + uniform
    # grid. No-op for slides without such a group (tables, charts, comparisons,
    # simple title-body keep their de-overlapped/clamped rects).
    rect_for = _reflow_slide(comp.blocks, rect_for)

    # First pass: resolve node rects so connectors can route edge-to-edge.
    node_rects: list[tuple[float, float, float, float]] = []
    for blk in comp.blocks:
        if blk.role == "node":
            node_rects.append(rect_for[id(blk)])
    node_centers = [(l + w / 2, t + h / 2) for (l, t, w, h) in node_rects]

    for blk in comp.blocks:
        role = blk.role
        if role == "decor":
            if blk.kind == "sparkle":
                P.sparkle(slide, blk.anchor, dark_bg=dark)
            elif blk.kind == "outline_corner":
                P.outline_corner(slide, blk.anchor, dark_bg=dark)
            elif blk.kind == "portal":
                P.portal(slide, _portal_base(blk.anchor), n=blk.portal_squares,
                         dark_bg=dark)
            continue
        if role == "connector":
            if 0 <= blk.src < len(node_centers) and 0 <= blk.dst < len(node_centers):
                p0 = _edge_point(node_rects[blk.src], node_centers[blk.dst])
                p1 = _edge_point(node_rects[blk.dst], node_centers[blk.src])
                P.arrow(slide, p0, p1, rhombus=blk.rhombus, dark_bg=dark)
            continue
        rect = rect_for[id(blk)]
        if role == "title":
            P.title_block(slide, blk.text, rect, size_pt=blk.size_pt,
                          accent_underline=blk.accent_underline, dark_bg=dark)
        elif role == "body":
            P.body_block(slide, blk.bullets, rect, size_pt=blk.size_pt, dark_bg=dark)
        elif role == "kpi":
            P.kpi_block(slide, blk.num, blk.desc, rect, dark_bg=dark)
        elif role == "chart":
            P.chart_block(
                slide, blk.chart_type, blk.categories,
                [s.model_dump() for s in blk.series], rect, accent_idx=blk.accent_idx,
                data_provenance=blk.data_provenance, dark_bg=dark,
            )
        elif role == "table":
            P.table_block(
                slide, blk.headers, blk.rows, rect,
                accent_col=blk.accent_col, first_col_wider=blk.first_col_wider,
                dark_bg=dark,
            )
        elif role == "node":
            P.node_box(slide, blk.text, rect, accent=blk.accent, dark_bg=dark)
        elif role == "card":
            P.person_card(slide, blk.heading, blk.sub, rect, plate=blk.plate,
                          accent=blk.accent, dark_bg=dark)
        elif role == "milestone":
            P.milestone_tick(slide, blk.label, blk.text, rect, accent=blk.accent,
                             dark_bg=dark)


def build_deck(comps: list[Composition], out_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Emu(CANVAS_W * P.EMU_PER_PX)
    prs.slide_height = Emu(CANVAS_H * P.EMU_PER_PX)
    for comp in comps:
        assemble_slide(prs, comp)
    prs.save(out_path)
    return out_path

#!/usr/bin/env python3
"""make_chart_test_deck.py — build a synthetic .pptx with REAL native charts.

None of the 5 benchmark decks carry actual PPTX chart objects, so the native-
chart extraction path (parse_pptx._extract_chart → _inject_parsed_charts →
chart_native_pptx.render_chart_pptx_slide) was never exercised end-to-end.

This deterministically produces a small deck whose slides hold genuine
``add_chart`` graphicFrame objects (column-clustered, pie, line) plus a plain
text slide, so the whole extraction→inject→build chain can be validated.

Usage:
    python scripts/make_chart_test_deck.py [out.pptx]

Default output: tests/fixtures/synthetic_chart_deck.pptx
"""
from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu


DEFAULT_OUT = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "tests", "fixtures", "synthetic_chart_deck.pptx",
)

# A wide canvas so the charts have room; deterministic geometry.
_LEFT, _TOP, _W, _H = Emu(457200), Emu(914400), Emu(8229600), Emu(4572000)


def _add_blank(prs):
    # layout 6 is the blank layout in the default python-pptx template.
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_title_only(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    return slide


def build(out_path: str = DEFAULT_OUT) -> str:
    prs = Presentation()

    # ── Slide 1: COLUMN_CLUSTERED (bar), 2 series, 4 categories, titled ──
    s1 = _add_blank(prs)
    cd = CategoryChartData()
    cd.categories = ["2022", "2023", "2024", "2025"]
    cd.add_series("Облако", (12.0, 18.5, 24.0, 31.2))
    cd.add_series("On-prem", (40.0, 35.0, 28.0, 20.0))
    gf = s1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, _LEFT, _TOP, _W, _H, cd)
    gf.chart.has_title = True
    gf.chart.chart_title.text_frame.text = "Выручка по сегментам"

    # ── Slide 2: PIE, 1 series, 4 categories ──
    s2 = _add_blank(prs)
    cd = CategoryChartData()
    cd.categories = ["IaaS", "PaaS", "SaaS", "Прочее"]
    cd.add_series("Доля рынка", (45.0, 25.0, 20.0, 10.0))
    gf = s2.shapes.add_chart(XL_CHART_TYPE.PIE, _LEFT, _TOP, _W, _H, cd)
    gf.chart.has_title = True
    gf.chart.chart_title.text_frame.text = "Структура выручки"

    # ── Slide 3: LINE, 2 series, 5 categories ──
    s3 = _add_blank(prs)
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3", "Q4", "Q5"]
    cd.add_series("ARR", (10.0, 14.0, 19.0, 26.0, 35.0))
    cd.add_series("Churn", (5.0, 4.5, 4.0, 3.2, 2.8))
    gf = s3.shapes.add_chart(XL_CHART_TYPE.LINE, _LEFT, _TOP, _W, _H, cd)
    gf.chart.has_title = True
    gf.chart.chart_title.text_frame.text = "Динамика ARR и оттока"

    # ── Slide 4: plain text (so the deck isn't chart-only) ──
    _add_title_only(prs, "Итоги квартала")

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path


def main() -> None:
    out = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_OUT
    path = build(out)
    print(f"Saved {path}")


if __name__ == "__main__":
    main()

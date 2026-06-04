"""Create a tiny 4-slide test draft .pptx for the live pipeline run."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def build(out: Path) -> Path:
    prs = Presentation()
    # Slide 1 — title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Cloud.ru Quarterly Review"
    s.placeholders[1].text = "Q1 2026 — Sales & product highlights"

    # Slide 2 — bullet text
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Key Achievements"
    tf = s.placeholders[1].text_frame
    tf.text = "Revenue grew 32% year-over-year"
    p = tf.add_paragraph(); p.text = "New enterprise clients: 14"
    p = tf.add_paragraph(); p.text = "Platform uptime: 99.97%"
    p = tf.add_paragraph(); p.text = "Launched 3 new product lines"

    # Slide 3 — KPI-style content (short labels + numbers)
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Financial Metrics"
    tf = s.placeholders[1].text_frame
    tf.text = "MRR: 12.4M RUB (+18% QoQ)"
    p = tf.add_paragraph(); p.text = "Gross margin: 67%"
    p = tf.add_paragraph(); p.text = "Customer acquisition cost: 84k RUB"
    p = tf.add_paragraph(); p.text = "LTV/CAC ratio: 4.2x"

    # Slide 4 — comparison / closing
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Next Quarter Priorities"
    tf = s.placeholders[1].text_frame
    tf.text = "Scale infrastructure to 2x capacity"
    p = tf.add_paragraph(); p.text = "Hire 8 senior engineers"
    p = tf.add_paragraph(); p.text = "Launch managed AI tier"
    p = tf.add_paragraph(); p.text = "Open Almaty data center"

    prs.save(str(out))
    return out


if __name__ == "__main__":
    out = Path(sys.argv[1] if len(sys.argv) > 1 else "test_draft.pptx")
    p = build(out)
    print(f"wrote {p} ({p.stat().st_size} bytes, {Presentation(str(p)).slides.__len__()} slides)")
